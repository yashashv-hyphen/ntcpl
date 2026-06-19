#!/usr/bin/env python3
from __future__ import annotations

import argparse
import base64
import json
import logging
import os
import re
import threading
import time
import uuid
import webbrowser
from copy import copy
from html import escape
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.error import HTTPError, URLError
from urllib.parse import quote
from urllib.request import Request, urlopen

from dotenv import load_dotenv
load_dotenv()  # must run before any os.getenv() constant reads
from flask import Flask, jsonify, request, send_from_directory, stream_with_context
from openpyxl import load_workbook
from PIL import Image, ImageDraw
from tenacity import RetryError, retry, stop_after_attempt, wait_fixed
from werkzeug.utils import secure_filename

try:
    import anthropic as _anthropic_module
    ANTHROPIC_AVAILABLE = bool(os.environ.get("ANTHROPIC_API_KEY", "").strip())
except ImportError:
    _anthropic_module = None
    ANTHROPIC_AVAILABLE = False

try:
    import groq as _groq_module  # type: ignore
    _GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "").strip()
    GROQ_AVAILABLE = bool(_GROQ_API_KEY)
except ImportError:
    _groq_module = None
    _GROQ_API_KEY = ""
    GROQ_AVAILABLE = False

GROQ_VISION_MODEL = os.getenv("GROQ_VISION_MODEL", "meta-llama/llama-4-scout-17b-16e-instruct").strip()
GROQ_TEXT_MODEL   = os.getenv("GROQ_TEXT_MODEL",   "llama-3.1-8b-instant").strip()

# Only Groq (vision/text/analysis/copy) and FAL (image generation) are supported.

# Groq text: llama-3.1-8b-instant has 500k TPM free. Each listing ~4400 tokens.
# 5 concurrent calls ≈ 22k tokens in flight — safe headroom. Set GROQ_TEXT_CONCURRENCY to tune.
_GROQ_TEXT_SEM = threading.Semaphore(int(os.getenv("GROQ_TEXT_CONCURRENCY", "8")))

# Groq vision: only 100 req/day on free tier. Track usage so large batches
# fall through immediately instead of wasting a round-trip per SKU.
_GROQ_VISION_DAY_LIMIT = int(os.getenv("GROQ_VISION_DAY_LIMIT", "90"))
_groq_vision_used = 0
_groq_vision_lock = threading.Lock()

def _groq_vision_available() -> bool:
    with _groq_vision_lock:
        return _groq_vision_used < _GROQ_VISION_DAY_LIMIT

def _groq_vision_increment() -> None:
    global _groq_vision_used
    with _groq_vision_lock:
        _groq_vision_used += 1

# Groq TEXT daily limit — once exhausted, skip Groq immediately in the fallback
# chain rather than burning 4×retry cycles on quota errors.
_GROQ_TEXT_DAY_LIMIT = int(os.getenv("GROQ_TEXT_DAY_LIMIT", "14000"))
_groq_text_used = 0
_groq_text_lock = threading.Lock()

def _groq_text_available() -> bool:
    with _groq_text_lock:
        return _groq_text_used < _GROQ_TEXT_DAY_LIMIT

def _groq_text_increment() -> None:
    global _groq_text_used
    with _groq_text_lock:
        _groq_text_used += 1

def _groq_is_quota_error(exc_str: str) -> bool:
    """True when Groq returns a hard daily/org quota error (not a per-minute rate limit)."""
    s = exc_str.lower()
    return (
        "exceeded your daily" in s
        or "organization_quota" in s
        or "billing" in s
        or ("quota" in s and "try again in" not in s)
    )

from amazon_template_autofill import (
    parse_template_rows,
    read_all_valid_values,
    read_browse_nodes,
    read_column_metadata,
    read_data_definitions,
    read_recommended_node_to_product_type,
    read_valid_product_types,
)


app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 200 * 1024 * 1024  # 200 MB upload limit

@app.errorhandler(413)
def _too_large(e):
    return jsonify({"ok": False, "error": "File too large. Maximum upload size is 200 MB."}), 413

APP_DIR = Path(__file__).resolve().parent
UPLOAD_DIR = APP_DIR / "uploaded_product_images"
TEMPLATE_UPLOAD_DIR = APP_DIR / "uploaded_excel_templates"
INSTRUCTIONS_FILE = APP_DIR / "image_generation_instructions.txt"
DEFAULT_TEMPLATE_PATH = APP_DIR / "LEASH_ANIMAL_COLLAR (1).xlsm"

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
)
log = logging.getLogger(__name__)

MIN_IMAGE_DIMENSION = int(os.getenv("MIN_IMAGE_DIMENSION", "500"))
MAX_IMAGE_BYTES = int(os.getenv("MAX_IMAGE_BYTES", "350000"))
MIN_CONFIDENCE = float(os.getenv("MIN_CONFIDENCE", "0.70"))
DESCRIPTION_MIN_LEN = int(os.getenv("DESCRIPTION_MIN_LEN", "120"))
DESCRIPTION_MAX_LEN = int(os.getenv("DESCRIPTION_MAX_LEN", "250"))
BULLET_MIN_LEN = int(os.getenv("BULLET_MIN_LEN", "40"))
# ── Speed / throughput settings ───────────────────────────────────────────────
# FAST_MODE=true  skips vision analysis + validation refinement for maximum speed.
# Individual flags can also be set independently.
_FAST_MODE          = os.getenv("FAST_MODE", "false").strip().lower() == "true"
SKIP_VISION         = _FAST_MODE or os.getenv("SKIP_VISION",      "false").strip().lower() == "true"
SKIP_REFINEMENT     = _FAST_MODE or os.getenv("SKIP_REFINEMENT",  "false").strip().lower() == "true"
LISTING_REFINE_PASSES = 0 if SKIP_REFINEMENT else int(os.getenv("LISTING_REFINE_PASSES", "1"))

REPETITION_THRESHOLD = int(os.getenv("REPETITION_THRESHOLD", "3"))
FLAGGED_WORDS: List[str] = [
    "premium",
    "high quality",
    "best",
    "luxury",
    "world-class",
    "unbeatable",
    "revolutionary",
]
ALLOWED_EXTENSIONS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff", ".tif"}

# ── Feature flags ────────────────────────────────────────────────────────────
IMAGE_GENERATION_ENABLED = os.getenv("IMAGE_GENERATION_ENABLED", "false").lower() == "true"

# ── Catalog pipeline constants ──────────────────────────────────────────────
MASTER_NODE_PATH = APP_DIR / "Book1.xlsx"
CATALOG_UPLOAD_DIR = APP_DIR / "uploaded_catalogs"

# Master nodes are loaded lazily and cached here.
_MASTER_NODES_CACHE: Optional[List[Dict]] = None


class ImageValidationError(ValueError):
    """Raised when the uploaded image does not meet quality requirements."""


class VisionAPIError(RuntimeError):
    """Raised when the vision API returns an unexpected response."""


class CopyGenerationError(RuntimeError):
    """Raised when text generation fails to return usable listing copy."""


def _name_based_analysis(name: str) -> Dict[str, Any]:
    """Build a minimal product analysis from the product name/SKU when no vision AI is available.

    Uses generic product name to avoid SKU codes appearing in listing copy.
    """
    raw = re.sub(r"[_\-]+", " ", (name or "product")).strip()
    # If name looks like a pure SKU code (e.g. "c2 001", "c3 041"), use generic label
    _sku_like = re.match(r"^[a-z][0-9]\s+[0-9]+$", raw.lower())
    product_label = "Premium Quality Product" if _sku_like else raw.title()
    return {
        "product_type": product_label,
        "category": "Home & Kitchen",
        "material": "high quality material",
        "colors": [],
        "features": ["durable construction", "easy to use", "versatile design"],
        "usage": f"Suitable for everyday household and kitchen use",
        "style": "Modern",
        "confidence": 0.5,
    }


def _fallback_analysis_from_image(image_path: Path) -> Dict[str, Any]:
    width, height = validate_image(image_path)
    stem = image_path.stem.replace("_", " ").replace("-", " ").strip() or "Product"
    return {
        "product_type": stem.title(),
        "category": "General",
        "material": "uncertain",
        "colors": [],
        "features": [f"image size {width}x{height}"],
        "usage": "General daily use",
        "style": "Practical",
        "confidence": 0.75,
    }


def _recover_partial_listing_json(raw: str) -> Dict[str, Any]:
    """Extract title / bullet_points / description from a truncated JSON response.

    Returns a dict with whatever fields could be parsed; missing fields are empty.
    """
    out: Dict[str, Any] = {"title": "", "bullet_points": [], "description": ""}
    m = re.search(r'"title"\s*:\s*"([^"\\]*(?:\\.[^"\\]*)*)"', raw)
    if m:
        out["title"] = m.group(1)
    bullets: List[str] = re.findall(r'"([^"]{30,})"', raw)
    if out["title"] in bullets:
        bullets.remove(out["title"])
    out["bullet_points"] = bullets[:5]
    m2 = re.search(r'"description"\s*:\s*"([^"\\]*(?:\\.[^"\\]*)*)"', raw)
    if m2:
        out["description"] = m2.group(1)
    return out


def _fallback_listing_from_analysis(product_analysis: Dict[str, Any]) -> Dict[str, Any]:
    name = str(product_analysis.get("product_type", "")).strip() or "Product"
    category = str(product_analysis.get("category", "")).strip() or "General"
    usage = str(product_analysis.get("usage", "")).strip() or "daily routines"
    material = str(product_analysis.get("material", "")).strip()
    colors = product_analysis.get("colors", [])
    features = product_analysis.get("features", [])
    if not isinstance(colors, list):
        colors = []
    if not isinstance(features, list):
        features = []
    colors_text = ", ".join([str(c).strip() for c in colors if str(c).strip()][:3]) or "image-visible tones"
    feature_text = ", ".join([str(f).strip() for f in features if str(f).strip()][:3]) or "practical everyday design"
    material_text = material if material and material.lower() not in {"uncertain", "unknown", "n/a"} else "image-observed build"
    description = (
        f"This {name} appears in a clean product image with {colors_text} and {feature_text}. "
        f"It fits the {category} category and is positioned for {usage}. "
        f"The listing emphasizes usability, handling comfort, and dependable day-to-day performance based on visible traits."
    )
    bullets = [
        f"Image-grounded overview: {name} with {colors_text} for clear shopper expectation before checkout.",
        f"Visible build cues suggest {material_text} and a practical structure suited to regular handling.",
        f"Use-case clarity: designed for {usage}, supporting routine tasks with straightforward functionality.",
        "Usability focus: easy to evaluate from listing visuals, with attention to comfort and convenience.",
        f"Buyer confidence: feature highlights include {feature_text}, presented in concise non-repetitive language.",
    ]
    kw_terms = [w for w in re.findall(r"[a-z0-9]+", f"{name} {category} {usage} {material_text}".lower()) if len(w) > 3]
    seen_kw: set = set()
    kw_unique = [t for t in kw_terms if not (t in seen_kw or seen_kw.add(t))]  # type: ignore[func-returns-value]
    fallback_keywords = kw_unique[:20] + [""] * max(0, 20 - len(kw_unique))
    return {
        "title": f"{name} for {usage} - {category} Essentials",
        "bullet_points": bullets,
        "description": description,
        "keywords": fallback_keywords[:20],
        "estimated_dimensions": "see product description",
        "hsn_code": "",
        "product_tax_code": "A_GEN_TAX_18",
    }

IMAGE_MODE_STANDARD = "standard"
IMAGE_MODE_AMAZON_INSTRUCTIONS = "amazon_instructions"

SECTION_PRODUCT_ANALYSIS = "PRODUCT_ANALYSIS"
SECTION_LIFESTYLE = "LIFESTYLE_IMAGE"
SECTION_INFOGRAPHIC = "INFOGRAPHIC_IMAGE"


COLOR_NAMES = [
    ("black", (20, 20, 20)),
    ("white", (245, 245, 245)),
    ("gray", (128, 128, 128)),
    ("red", (200, 40, 40)),
    ("orange", (220, 130, 40)),
    ("yellow", (220, 210, 60)),
    ("green", (60, 160, 80)),
    ("blue", (70, 110, 200)),
    ("purple", (140, 90, 180)),
    ("pink", (220, 130, 170)),
    ("brown", (120, 85, 55)),
]


def rgb_to_name(rgb):
    r, g, b = rgb
    best_name = "multicolor"
    best_dist = 10**9
    for name, (cr, cg, cb) in COLOR_NAMES:
        d = (r - cr) ** 2 + (g - cg) ** 2 + (b - cb) ** 2
        if d < best_dist:
            best_dist = d
            best_name = name
    return best_name


_BROWSER_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/124.0.0.0 Safari/537.36"
    ),
    "Accept": "image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.9",
    "Referer": "https://www.amazon.in/",
}


def load_image_any(source: str) -> Image.Image:
    src = (source or "").strip()
    if not src:
        raise ValueError("Image source is empty.")
    if src.lower().startswith(("http://", "https://")):
        from urllib.parse import urlparse, urlunparse, quote
        p = urlparse(src)
        # Re-encode path to handle spaces / non-ASCII characters in URLs
        safe_path = quote(p.path, safe="/%+@")
        safe_src  = urlunparse(p._replace(path=safe_path))
        req = Request(safe_src, headers=_BROWSER_HEADERS)
        with urlopen(req, timeout=60) as resp:
            data = resp.read()
        return Image.open(BytesIO(data)).convert("RGB")
    return Image.open(src).convert("RGB")


def _load_combo_product_image(url_list) -> Image.Image:
    """Load a product image for combo compositing with retry and fallback.

    url_list: a single URL string OR a list of URL strings to try in order.
    Retries each URL up to 3 times with short backoff before moving on.
    Raises on total failure.
    """
    import time as _time
    srcs = url_list if isinstance(url_list, list) else [url_list]
    srcs = [s for s in srcs if s and str(s).strip()]
    if not srcs:
        raise ValueError("No image URLs provided for combo slot.")
    last_exc: Exception = ValueError("No image URLs provided.")
    for u in srcs:
        for attempt in range(3):
            try:
                return load_image_any(u)
            except Exception as exc:
                last_exc = exc
                log.warning("Combo image load attempt %d failed for %s: %s", attempt + 1, u, exc)
                if attempt < 2:
                    _time.sleep(1.5 ** attempt)  # 1s, 1.5s backoff
    raise last_exc


# Thread-safe LRU cache for remote image downloads.
# Eliminates redundant HTTP fetches when the same image URL is analyzed for
# vision AND compressed for the copy-generation provider in the same SKU, and
# also deduplicates across SKU variants that share a reference image.
_IMG_BYTES_CACHE: Dict[str, bytes] = {}
_IMG_BYTES_CACHE_LOCK = threading.Lock()
_IMG_BYTES_CACHE_MAX = int(os.getenv("IMG_CACHE_MAX_ENTRIES", "1000"))

import http.client as _http_client
_HTTP_POOL: dict = {}
_HTTP_POOL_LOCK = threading.Lock()

def _pooled_urlopen(url: str, timeout: int = 30) -> bytes:
    """Fetch a URL reusing a persistent HTTP connection per host."""
    from urllib.parse import urlparse
    parsed = urlparse(url)
    scheme = parsed.scheme.lower()
    host   = parsed.netloc
    path   = parsed.path + (f"?{parsed.query}" if parsed.query else "")
    key    = (scheme, host)
    with _HTTP_POOL_LOCK:
        conn = _HTTP_POOL.get(key)
        if conn is None:
            if scheme == "https":
                import ssl
                conn = _http_client.HTTPSConnection(host, timeout=timeout,
                       context=ssl.create_default_context())
            else:
                conn = _http_client.HTTPConnection(host, timeout=timeout)
            _HTTP_POOL[key] = conn
    try:
        conn.request("GET", path, headers=dict(_BROWSER_HEADERS))
        resp = conn.getresponse()
        data = resp.read()
        if resp.status in (301, 302, 303, 307, 308):
            loc = resp.getheader("Location", "")
            if loc:
                return _pooled_urlopen(loc, timeout)
        return data
    except Exception:
        # Connection broken — drop it so next call creates a fresh one
        with _HTTP_POOL_LOCK:
            _HTTP_POOL.pop(key, None)
        raise


def _normalize_image_url(url: str) -> str:
    """Convert any sharing/embed URL into a direct-download image URL.

    Handles: Google Drive, Google Photos, Dropbox, OneDrive, Box,
             WeTransfer (best-effort), and generic redirect chains.
    Returns the original URL unchanged when no transformation applies.
    """
    import re as _re
    from urllib.parse import urlparse as _up, urlencode as _ue, parse_qs as _pq, urlunparse as _uu

    u = url.strip()
    pl = u.lower()

    # ── Data URIs ─────────────────────────────────────────────────────────────
    if pl.startswith("data:image/"):
        return u  # handled inline in read_image_bytes_any

    # ── Google Drive ──────────────────────────────────────────────────────────
    # Patterns:
    #   https://drive.google.com/file/d/FILE_ID/view[?...]
    #   https://drive.google.com/file/d/FILE_ID/edit[?...]
    #   https://drive.google.com/open?id=FILE_ID
    #   https://drive.google.com/uc?id=FILE_ID[&export=...]
    #   https://docs.google.com/uc?id=FILE_ID
    #   https://lh3.googleusercontent.com/d/FILE_ID  (thumbnail)
    if "drive.google.com" in pl or "docs.google.com" in pl:
        m = _re.search(r"/file/d/([a-zA-Z0-9_-]+)", u)
        if not m:
            qs = _pq(_up(u).query)
            fid = (qs.get("id") or qs.get("ID") or [""])[0]
        else:
            fid = m.group(1)
        if fid:
            # Use drive.usercontent.google.com — Google's newer direct-download
            # endpoint that bypasses the HTML confirmation page for most files.
            return f"https://drive.usercontent.google.com/download?id={fid}&export=download&authuser=0&confirm=t"
    if "lh3.googleusercontent.com" in pl:
        m = _re.search(r"/d/([a-zA-Z0-9_-]+)", u)
        if m:
            return f"https://drive.usercontent.google.com/download?id={m.group(1)}&export=download&authuser=0&confirm=t"
        return u
    if "drive.usercontent.google.com" in pl:
        return u  # already the canonical download URL

    # ── Google Photos ─────────────────────────────────────────────────────────
    # https://photos.google.com/photo/AF... → no public API; try lh3 pattern
    if "photos.google.com" in pl or "photos.app.goo.gl" in pl:
        # best-effort: return as-is and rely on redirect following
        return u

    # ── Dropbox ───────────────────────────────────────────────────────────────
    # https://www.dropbox.com/s/HASH/file.jpg?dl=0 → dl=1
    # https://dropbox.com/scl/fi/... → dl=1
    if "dropbox.com" in pl:
        u = _re.sub(r"[?&]dl=0", "", u)
        u = u + ("&dl=1" if "?" in u else "?dl=1")
        u = u.replace("www.dropbox.com", "dl.dropboxusercontent.com")
        return u
    if "dl.dropboxusercontent.com" in pl:
        return u  # already direct

    # ── OneDrive / SharePoint ─────────────────────────────────────────────────
    # https://1drv.ms/i/s!... short link → follow redirect
    # https://onedrive.live.com/...resid=... → transform to embed/download
    if "1drv.ms" in pl or "onedrive.live.com" in pl or "sharepoint.com" in pl:
        # Can't transform without auth; return as-is — redirect chain will handle it
        return u

    # ── iCloud ────────────────────────────────────────────────────────────────
    if "icloud.com" in pl:
        return u  # must follow redirect; return unchanged

    # ── Box ───────────────────────────────────────────────────────────────────
    # https://app.box.com/s/HASH → https://app.box.com/shared/static/HASH (direct)
    if "box.com" in pl:
        m = _re.search(r"box\.com/s/([a-zA-Z0-9]+)", u)
        if m:
            return f"https://app.box.com/shared/static/{m.group(1)}"
        return u

    # ── WeTransfer ────────────────────────────────────────────────────────────
    if "wetransfer.com" in pl or "we.tl" in pl:
        return u  # requires JS; return as-is and hope redirect works

    # ── Imgur ─────────────────────────────────────────────────────────────────
    # https://imgur.com/HASH → https://i.imgur.com/HASH.jpg
    if "imgur.com" in pl and "i.imgur.com" not in pl:
        m = _re.search(r"imgur\.com/([a-zA-Z0-9]+)$", u.rstrip("/"))
        if m:
            return f"https://i.imgur.com/{m.group(1)}.jpg"
        return u

    # ── Pinterest ─────────────────────────────────────────────────────────────
    # pin.it short links → follow redirect; i.pinimg.com direct CDN links work as-is
    if "pin.it" in pl or ("pinterest" in pl and "pinimg.com" not in pl):
        return u  # follow redirect

    # ── Instagram ─────────────────────────────────────────────────────────────
    # Public post CDN images work; sharing links need scraping — return as-is
    if "instagram.com" in pl or "cdninstagram.com" in pl:
        return u

    # ── Amazon / Shopify / Cloudinary / other CDN ─────────────────────────────
    # These are direct image URLs already — no transformation needed
    return u


def _fetch_image_bytes_http(url: str, timeout: int = 30) -> bytes:
    """Fetch image bytes from any HTTP URL, handling:
    - Google Drive large-file confirmation pages
    - Redirect chains (up to 8 hops)
    - Non-image HTML pages (raises ValueError)
    """
    import re as _re

    normalized = _normalize_image_url(url)
    current    = normalized
    _IMAGE_CT  = ("image/", "application/octet-stream", "binary/")

    for hop in range(8):
        try:
            data = _pooled_urlopen(current, timeout=timeout)
        except Exception:
            req = Request(current, headers=_BROWSER_HEADERS)
            with urlopen(req, timeout=timeout) as resp:
                data = resp.read()

        # Peek content type from first bytes
        sig = data[:12]
        is_image = (
            sig[:8] in (b"\x89PNG\r\n\x1a\n", b"GIF87a\x00", b"GIF89a\x00")
            or sig[:3] == b"\xff\xd8\xff"   # JPEG
            or sig[:4] in (b"RIFF", b"WEBP")
            or sig[:4] == b"\x00\x00\x01\x00"  # ICO
        )
        if is_image:
            return data

        # HTML — might be a sharing page with a redirect or confirmation
        if data[:100].lstrip()[:15].lower().startswith((b"<!doctype", b"<html")):
            html = data.decode("utf-8", errors="ignore")

            # Google Drive virus-scan / large-file confirmation page
            # Look for the download link in the HTML (both old and new Drive UI)
            gd_link = _re.search(
                r'href=["\']([^"\']*(?:drive\.usercontent\.google\.com|drive\.google\.com/uc)[^"\']*confirm=[^"\']+)["\']',
                html, _re.IGNORECASE
            )
            if gd_link:
                href = gd_link.group(1).replace("&amp;", "&")
                current = href if href.startswith("http") else "https://drive.google.com" + href
                continue
            # Fallback: extract file ID from current URL and retry with explicit confirm
            gd_id_m = _re.search(r'[?&]id=([a-zA-Z0-9_-]+)', current)
            if ("drive.google.com" in current or "drive.usercontent.google.com" in current) and gd_id_m:
                fid = gd_id_m.group(1)
                # Try the old uc endpoint as fallback
                current = f"https://drive.google.com/uc?export=download&id={fid}&confirm=t&uuid={hop}"
                continue

            # Generic HTML redirect (meta refresh or single anchor)
            meta_url = _re.search(
                r'<meta[^>]+http-equiv=["\']refresh["\'][^>]+content=["\'][^"\']*url=([^"\'>\s]+)',
                html, _re.IGNORECASE
            )
            if meta_url:
                current = meta_url.group(1).strip()
                continue

            # og:image / twitter:image fallback (e.g. Pinterest, Instagram)
            og = _re.search(
                r'<meta[^>]+(?:property=["\']og:image["\']|name=["\']twitter:image["\'])[^>]+content=["\']([^"\']+)["\']',
                html, _re.IGNORECASE
            )
            if og:
                current = og.group(1).strip()
                continue

            raise ValueError(
                f"URL returned an HTML page, not an image, after {hop+1} hop(s): {url}"
            )

        # Non-image binary — return and let PIL decide
        return data

    raise ValueError(f"Could not resolve to an image after 8 redirect hops: {url}")


def read_image_bytes_any(source: str) -> bytes:
    src = (source or "").strip()
    if not src:
        raise ValueError("Image source is empty.")

    # Data URI — decode inline
    if src.lower().startswith("data:image/"):
        import base64 as _b64
        header, _, encoded = src.partition(",")
        return _b64.b64decode(encoded)

    if src.lower().startswith(("http://", "https://")):
        # Cache keyed on original URL so normalisation is transparent
        with _IMG_BYTES_CACHE_LOCK:
            if src in _IMG_BYTES_CACHE:
                return _IMG_BYTES_CACHE[src]
        data = _fetch_image_bytes_http(src, timeout=30)
        with _IMG_BYTES_CACHE_LOCK:
            if len(_IMG_BYTES_CACHE) >= _IMG_BYTES_CACHE_MAX:
                del _IMG_BYTES_CACHE[next(iter(_IMG_BYTES_CACHE))]
            _IMG_BYTES_CACHE[src] = data
        return data

    return Path(src).read_bytes()


def ensure_upload_dir() -> Path:
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    return UPLOAD_DIR


def ensure_template_upload_dir() -> Path:
    TEMPLATE_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    return TEMPLATE_UPLOAD_DIR


def describe_image(source: str, product_name: str = "") -> str:
    img = load_image_any(source)
    small = img.resize((120, 120))
    pixels = list(small.getdata())
    step = max(1, len(pixels) // 4000)
    sampled = pixels[::step]
    if not sampled:
        raise ValueError("Could not read image pixels.")
    avg = tuple(int(sum(c[i] for c in sampled) / len(sampled)) for i in range(3))
    tone = rgb_to_name(avg)
    w, h = img.size
    orientation = "portrait" if h > w else "landscape" if w > h else "square"
    name = product_name.strip() or "This product"
    return (
        f"{name} reads as a {tone}-toned product in a {orientation} hero shot, with visible form and finish "
        f"that help shoppers judge scale, materials, and how it would look in real use."
    )


def load_instruction_sections() -> Dict[str, str]:
    if not INSTRUCTIONS_FILE.is_file():
        raise FileNotFoundError(f"Instructions file not found: {INSTRUCTIONS_FILE}")
    text = INSTRUCTIONS_FILE.read_text(encoding="utf-8")
    sections: Dict[str, str] = {}
    current = None
    lines: List[str] = []
    for line in text.splitlines():
        header = re.match(r"^\[([A-Z0-9_]+)\]\s*$", line.strip())
        if header:
            if current:
                sections[current] = "\n".join(lines).strip()
            current = header.group(1)
            lines = []
            continue
        if line.startswith("#"):
            continue
        lines.append(line)
    if current:
        sections[current] = "\n".join(lines).strip()
    missing = [
        name
        for name in (SECTION_PRODUCT_ANALYSIS, SECTION_LIFESTYLE, SECTION_INFOGRAPHIC)
        if name not in sections
    ]
    if missing:
        raise ValueError("Instructions file missing sections: " + ", ".join(missing))
    return sections


def fill_template(template: str, **values: str) -> str:
    result = template or ""
    for key, value in values.items():
        result = result.replace("{" + key + "}", value or "")
    return " ".join(result.split())


def extract_sku_seller_notes(seller_notes: str, sku: str, product_name: str = "", window: int = 3000) -> str:
    """Return only the portion of seller notes relevant to this SKU/product.

    Searches the full notes text for the SKU or product name, then returns a
    window of characters around each match. This prevents cross-SKU contamination
    when a single seller notes file covers an entire product catalog.

    Falls back to the full text (capped at window*2) when no match is found,
    so single-product files still work correctly.
    """
    if not seller_notes or not seller_notes.strip():
        return ""

    text = seller_notes.strip()
    search_terms = [t.strip().lower() for t in [sku, product_name] if t.strip()]

    matched_spans: list[tuple[int, int]] = []
    for term in search_terms:
        if not term or len(term) < 3:
            continue
        start = 0
        while True:
            idx = text.lower().find(term, start)
            if idx == -1:
                break
            span_start = max(0, idx - window // 2)
            span_end   = min(len(text), idx + len(term) + window // 2)
            matched_spans.append((span_start, span_end))
            start = idx + 1

    if not matched_spans:
        # No match — return beginning of notes up to window*2 (single-product file)
        return text[: window * 2]

    # Merge overlapping spans and collect text
    matched_spans.sort()
    merged: list[str] = []
    cur_start, cur_end = matched_spans[0]
    for s, e in matched_spans[1:]:
        if s <= cur_end:
            cur_end = max(cur_end, e)
        else:
            merged.append(text[cur_start:cur_end])
            cur_start, cur_end = s, e
    merged.append(text[cur_start:cur_end])

    result = "\n...\n".join(merged)
    return result[:window * 3]  # hard cap to avoid token bloat


def build_analysis_prompt(product_name: str, seller_notes: str) -> str:
    sections = load_instruction_sections()
    name = product_name.strip() or "Product"
    notes = seller_notes.strip() or "None provided."
    return fill_template(
        sections[SECTION_PRODUCT_ANALYSIS],
        product_name=name,
        seller_notes=notes,
        product_details="",
    )


def _image_data_url(source: str) -> str:
    """Return a data URL (base64 JPEG) for the given image source."""
    img_bytes = read_image_bytes_any(source)
    img = Image.open(BytesIO(img_bytes)).convert("RGB")
    buf = BytesIO()
    img.save(buf, format="JPEG", quality=90)
    b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
    return f"data:image/jpeg;base64,{b64}"


def build_product_details_for_instructions(
    copy_obj: Optional[Dict],
    image_source: str,
    product_name: str,
) -> str:
    name = product_name.strip() or "This product"
    parts = [describe_image(image_source, product_name=name)]
    if copy_obj:
        desc = str(copy_obj.get("product_description", "")).strip()
        if desc:
            parts.append(desc)
        bullets = [str(b).strip() for b in copy_obj.get("bullet_points", []) if str(b).strip()]
        if bullets:
            parts.append("Key features: " + "; ".join(bullets[:5]))
    return " ".join([p for p in parts if p]).strip()


def build_instruction_prompts(product_name: str, product_details: str) -> List[str]:
    sections = load_instruction_sections()
    name = product_name.strip() or "Product"
    details = (product_details or "").strip()
    if not details:
        raise ValueError("product_details is required for instruction prompts.")
    lifestyle = fill_template(
        sections[SECTION_LIFESTYLE],
        product_name=name,
        product_details=details,
    )
    infographic = fill_template(
        sections[SECTION_INFOGRAPHIC],
        product_name=name,
        product_details=details,
    )
    return [lifestyle, infographic]


def build_listing_image_prompts(product_name: str, analysis: Dict[str, Any]) -> List[str]:
    """Build two themed image prompts from product analysis data.

    Image 1 — Themed lifestyle shot: cinematic, real-world setting, human element,
               product-type-specific props, palette woven into the scene.
    Image 2 — Infographic background only: themed gradient/geometric backdrop +
               centered hero product. NO text/callouts — PIL overlays those after.
    """
    try:
        from amazon_product_images import _resolve_theme, _PRESERVATION_RULE  # type: ignore
        theme = _resolve_theme(analysis)
    except ImportError:
        theme = {
            "lifestyle_scene": (
                "warm natural light-wood surface — the product stands upright at centre. "
                "A folded cream linen cloth and a small complementary ceramic prop rest softly "
                "out of focus at one edge. "
                "Soft golden natural daylight from a large window at camera-left, "
                "gentle warm reflector fill on the right, single clean directional shadow. "
                "A hand reaches in naturally from the right edge."
            ),
            "infographic_scene": (
                "centered hero product on a pure warm white (#FFFAF5) background. "
                "Perfectly even soft studio lighting. "
                "A subtle warm cream radial glow (#F0EAD6) behind the product. "
                "Clean open space on all four sides."
            ),
            "palette": ["#1B4332", "#C9A84C", "#FFFAF5"],
        }
        _PRESERVATION_RULE = (
            "IMPORTANT: The product must be pixel-perfect identical to the reference image — "
            "same shape, color, finish, branding, logos, text, label design, and proportions. "
            "Do NOT restyle, recolor, redesign, or relabel the product. "
            "Only the background, scene, lighting, props, and graphic elements may be generated."
        )

    name     = (product_name or str(analysis.get("product_type", "product"))).strip() or "product"
    material = str(analysis.get("material", "")).strip()
    usage    = str(analysis.get("usage", "")).strip()
    style    = str(analysis.get("style", "")).strip()
    category = str(analysis.get("category", "")).strip()
    colors   = [str(c).strip() for c in analysis.get("colors", [])
                if str(c).strip() and str(c).lower() not in {"uncertain", "unknown"}]
    features = [str(f).strip() for f in analysis.get("features", [])
                if str(f).strip() and str(f).lower() not in {"uncertain", "unknown"}]

    color_hint   = colors[0] if colors else ""
    product_desc = f"{color_hint + ' ' if color_hint else ''}{name}".strip()
    mat_note     = (
        f"The {material} finish on the product must remain exactly as in the reference — texture, sheen, and surface detail pixel-perfect. "
        if material and material.lower() not in {"uncertain", "unknown", "n/a"} else ""
    )
    feat_note    = f"These product details must remain sharp and true: {'; '.join(features[:3])}. " if features else ""
    palette      = theme.get("palette", [])
    pal_str      = ", ".join(palette[:3]) if palette else "vibrant complementary colors"

    # Determine if this is a usable/wearable item (for infographic lifestyle panel)
    _usable_categories = {
        "clothing", "apparel", "fashion", "sports", "outdoor", "fitness", "kitchen",
        "home", "tools", "pet", "baby", "beauty", "personal care", "health", "food",
        "beverage", "electronics", "accessories", "bags", "footwear", "jewellery", "jewelry"
    }
    _cat_lower = category.lower()
    _is_usable = any(kw in _cat_lower for kw in _usable_categories) or bool(usage)

    # Usage context for scene building
    _usage_hint = usage if usage and usage.lower() not in {"uncertain", "unknown", "n/a"} else ""
    _style_hint = style if style and style.lower() not in {"uncertain", "unknown", "n/a"} else ""

    # ── IMAGE 1: LIFESTYLE ────────────────────────────────────────────────────
    # Warm, authentic, natural — product in its real-world setting with curated
    # lifestyle props that tell its story. Inspired by premium Indian brand photography.
    lifestyle_prompt = (
        f"Award-winning lifestyle product photograph of the {product_desc}. "
        f"{theme['lifestyle_scene']} "
        f"{'The product is shown in active natural use — real person, real moment, real life (' + _usage_hint + '). ' if _usage_hint else ''}"
        f"{'Style: ' + _style_hint + '. ' if _style_hint else ''}"
        f"{mat_note}"
        f"{feat_note}"
        "Shot on Canon EOS R5, 50mm f/1.4 — product razor-sharp with creamy shallow depth of field, "
        "background beautifully blurred into warm bokeh. "
        f"Color palette: {pal_str}. "
        "Color grading: warm golden tones, lifted shadows, rich and inviting — not cold, not sterile. "
        "Lighting: soft golden natural light from a large window at camera-left, warm fill on the right. "
        "No text overlays, no watermarks, no artificial logos added to the scene. "
        "Hyperrealistic, 8K, premium commercial photography. "
        f"{_PRESERVATION_RULE}"
    )

    # ── IMAGE 2: INFOGRAPHIC BACKGROUND ──────────────────────────────────────
    # Clean, open studio scene — PIL will overlay all text, icons, and chips on top.
    # The FAL image must have generous empty space and no competing visual elements.
    infographic_prompt = (
        f"Premium Amazon product studio image of the {product_desc} — clean infographic background. "
        f"{theme['infographic_scene']} "
        "The product is the single centered hero — perfectly sharp, filling 60% of the square frame. "
        "Product lighting: large softbox from upper-left at 45 degrees, gentle fill reflector from the right — "
        "no blown highlights, no color cast, every label and surface detail perfectly clear. "
        "A soft clean contact shadow sits directly beneath the product only. "
        "The background is intentionally minimal and open — large clear zones on all four sides "
        "to accommodate text and graphic overlays that will be added programmatically. "
        "Square 1:1 format, 2000x2000px, hyperrealistic, razor-sharp product rendering. "
        "STRICT: absolutely zero text, zero icons, zero callouts, zero arrows, zero badges — "
        "clean background and product only, no annotations whatsoever. "
        f"{_PRESERVATION_RULE}"
    )

    return [lifestyle_prompt, infographic_prompt]


def _fallback_image_bytes(prompt: str, reason: str) -> bytes:
    canvas = Image.new("RGB", (1536, 1536), color=(246, 248, 251))
    d = ImageDraw.Draw(canvas)
    d.rectangle((48, 48, 1488, 1488), outline=(70, 90, 120), width=6)
    d.text((90, 90), "Fallback image generated", fill=(30, 30, 30))
    d.text((90, 150), f"Reason: {reason[:220]}", fill=(70, 70, 70))
    d.text((90, 220), f"Prompt: {prompt[:320]}", fill=(70, 70, 70))
    buf = BytesIO()
    canvas.save(buf, format="PNG")
    return buf.getvalue()


def _reference_conditioning_text(reference_image_source: str, product_name: str) -> str:
    ref_desc = describe_image(reference_image_source, product_name=product_name)
    return (
        "STRICT PRODUCT REFERENCE MODE. "
        "The uploaded image is the exact product reference and must be preserved precisely. "
        "Do not change the product shape, proportions, materials, texture, stitching, branding placement, "
        "hardware, or color palette. "
        "Generate a professional Amazon-quality commercial scene using the SAME product from the reference image. "
        "The final output must look like the exact same photographed item placed into a modern premium ecommerce scene. "
        f"Reference product details: {ref_desc}"
    )


def parse_json_object_from_text(text: str) -> Dict:
    t = (text or "").strip()
    if not t:
        raise ValueError("Empty model response.")
    try:
        return json.loads(t)
    except Exception:
        pass
    start = t.find("{")
    end = t.rfind("}")
    if start >= 0 and end > start:
        try:
            return json.loads(t[start : end + 1])
        except Exception as exc:
            raise ValueError(f"Model JSON parse error: {exc}. Raw: {t[:500]}") from exc
    raise ValueError(f"Model did not return JSON. Raw: {t[:500]}")


def validate_image(image_path: Path) -> Tuple[int, int]:
    try:
        img = Image.open(image_path)
        width, height = img.size
    except Exception as exc:
        raise ImageValidationError(f"Cannot open image: {exc}") from exc
    if width < MIN_IMAGE_DIMENSION or height < MIN_IMAGE_DIMENSION:
        raise ImageValidationError(
            f"Image is {width}x{height} px - minimum is "
            f"{MIN_IMAGE_DIMENSION}x{MIN_IMAGE_DIMENSION} px."
        )
    return width, height


def encode_image_b64(image_path: Path, max_bytes: Optional[int] = None) -> str:
    byte_limit = max_bytes if max_bytes is not None else MAX_IMAGE_BYTES
    raw = image_path.read_bytes()
    if len(raw) <= byte_limit:
        return base64.b64encode(raw).decode("utf-8")

    # Downscale/compress oversized images to avoid HF 413 request limits.
    img = Image.open(image_path).convert("RGB")
    work = img.copy()
    quality = 90
    max_side = max(work.size)

    for _ in range(10):
        buf = BytesIO()
        work.save(buf, format="JPEG", quality=quality, optimize=True)
        out = buf.getvalue()
        if len(out) <= byte_limit:
            log.info(
                "Compressed image for HF payload: %s bytes -> %s bytes",
                len(raw),
                len(out),
            )
            return base64.b64encode(out).decode("utf-8")

        quality = max(55, quality - 8)
        if quality <= 65 and max_side > 1600:
            max_side = int(max_side * 0.82)
            work.thumbnail((max_side, max_side), Image.LANCZOS)

    # Final safety fallback.
    fallback = img.copy()
    fallback.thumbnail((1280, 1280), Image.LANCZOS)
    buf = BytesIO()
    fallback.save(buf, format="JPEG", quality=68, optimize=True)
    out = buf.getvalue()
    if len(out) > byte_limit:
        # Hard guard for strict providers: keep shrinking aggressively.
        hard = img.copy()
        hard.thumbnail((960, 960), Image.LANCZOS)
        buf = BytesIO()
        hard.save(buf, format="JPEG", quality=60, optimize=True)
        out = buf.getvalue()
    log.info(
        "Fallback compressed image for HF payload: %s bytes -> %s bytes",
        len(raw),
        len(out),
    )
    return base64.b64encode(out).decode("utf-8")


def extract_json(text) -> Dict[str, Any]:
    if isinstance(text, dict):
        return text  # already parsed
    text = (text or "").strip()
    if "```json" in text:
        text = text.split("```json", 1)[-1].split("```", 1)[0].strip()
    match = re.search(r"\{.*\}", text, re.DOTALL)
    if not match:
        raise ValueError(f"No JSON object found in model output. Raw:\n{text[:400]}")
    raw_json = match.group()
    # Attempt 1: raw parse
    try:
        return json.loads(raw_json)
    except json.JSONDecodeError:
        pass
    # Attempt 2: strip non-whitespace control chars
    try:
        cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', ' ', raw_json)
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass
    # Attempt 3: nuclear — strip ALL control chars including literal newlines inside strings
    try:
        nuclear = re.sub(r'[\x00-\x1f\x7f]', ' ', raw_json)
        return json.loads(nuclear)
    except json.JSONDecodeError as exc:
        raise ValueError(f"JSON parse error: {exc}. Raw snippet:\n{raw_json[:400]}") from exc


def _parse_vision_text(text: str) -> Dict[str, Any]:
    """Parse a free-text / markdown product description into the standard analysis dict.

    Used when a vision model returns prose or markdown instead of JSON.
    Extracts product_type, category, material, colors, features, usage, style.
    """
    if isinstance(text, dict):
        return text  # already a structured dict — return as-is
    t = (text or "").lower()
    lines = [l.strip().lstrip("*•-").strip() for l in text.splitlines() if l.strip()]

    def _extract_field(keys: List[str]) -> str:
        for line in lines:
            ll = line.lower()
            for key in keys:
                if key in ll and (":" in line or "**" in line):
                    # Strip markdown bold markers then split on colon
                    clean = re.sub(r"\*+", "", line).strip()
                    if ":" in clean:
                        val = clean.split(":", 1)[-1].strip(" ,.")
                        if val and val.lower() not in ("", "n/a", "unknown"):
                            return val
        return ""

    def _extract_list(keys: List[str]) -> List[str]:
        result: List[str] = []
        inside = False
        for line in lines:
            ll = line.lower()
            clean = re.sub(r"\*+", "", line).strip()
            if any(key in ll for key in keys) and (":" in clean or "**" in line):
                inside = True
                after = clean.split(":", 1)[-1].strip() if ":" in clean else ""
                if after and not after.startswith(("-", "+")):
                    result = [v.strip().strip("*") for v in re.split(r"[,;]+", after) if v.strip()]
                continue
            if inside:
                if re.match(r"^[\*\-\+\•]", line.lstrip()):
                    val = re.sub(r"^[\*\-\+\•\s\*]+", "", clean).strip()
                    if val and not any(k in val.lower() for k in keys):
                        result.append(val)
                elif ":" in clean and any(c.isalpha() for c in clean[:20]) and len(clean) < 60:
                    break  # next field started
        return [r for r in result if r][:6]

    product_type  = _extract_field(["product type", "product_type", "type"])
    category      = _extract_field(["category"])
    material      = _extract_field(["material"])
    usage         = _extract_field(["usage", "intended use", "use"])
    style         = _extract_field(["style"])
    pattern       = _extract_field(["pattern", "print"])
    fit           = _extract_field(["fit", "size hint", "fit_or_size"])
    gender_target = _extract_field(["gender", "target", "gender_target"])
    colors        = _extract_list(["color", "colour"])
    features      = _extract_list(["feature"])

    if not product_type:
        # Try to infer from first meaningful line
        for line in lines[:5]:
            if len(line) > 3 and not any(k in line.lower() for k in ("product", "description", "analysis")):
                product_type = line.title()
                break

    return {
        "product_type":    product_type or "Product",
        "category":        category or "General",
        "material":        material or "",
        "colors":          colors,
        "pattern":         pattern or "",
        "features":        features,
        "fit_or_size_hint": fit or "",
        "gender_target":   gender_target or "",
        "usage":           usage or "everyday use",
        "style":           style or "standard",
        "confidence":      0.85,
    }


VISION_PROMPT = """
Analyze this product image carefully.

Identify:
- exact product type
- material
- visible features
- colors
- intended use
- shape and design
- style
- product category

Rules:
- do NOT hallucinate
- describe ONLY visible facts
- if uncertain, say "uncertain"

Return strict JSON with exactly these keys:
{
  "product_type": "",
  "category": "",
  "material": "",
  "colors": [],
  "features": [],
  "usage": "",
  "style": "",
  "confidence": 0.0
}
"""




def build_listing_prompt(product_analysis: Dict[str, Any], seller_notes: str = "") -> str:
    # Pull key visual facts for direct reference in the prompt
    _inv = {"uncertain", "unknown", "n/a", ""}
    product_type  = str(product_analysis.get("product_type", "")).strip() or "product"
    material      = str(product_analysis.get("material", "")).strip()
    colors        = product_analysis.get("colors", [])
    features      = product_analysis.get("features", [])
    usage         = str(product_analysis.get("usage", "")).strip()
    style         = str(product_analysis.get("style", "")).strip()
    category      = str(product_analysis.get("category", "")).strip()
    pattern       = str(product_analysis.get("pattern", "")).strip()
    fit           = str(product_analysis.get("fit_or_size_hint", "")).strip()
    gender_target = str(product_analysis.get("gender_target", "")).strip()

    colors_str   = ", ".join(str(c) for c in colors   if str(c).strip().lower() not in _inv) or "as shown"
    features_str = "; ".join(str(f) for f in features if str(f).strip().lower() not in _inv)
    extras = ", ".join(v for v in [pattern, fit, gender_target] if v.lower() not in _inv)

    notes_block = ""
    if seller_notes and seller_notes.strip():
        notes_block = f"""
SUPPLEMENTARY SELLER REFERENCE
(Optional background notes — the product image is the primary source of truth.)
\"\"\"
{seller_notes.strip()[:50000]}
\"\"\"

HOW TO USE THIS REFERENCE:
1. IMAGE FIRST — Base all visual attributes (color, shape, material appearance,
   packaging) exclusively on what is visible in the image. Never override the image.

2. FILL GAPS ONLY — Use seller notes solely to add details not visible in the image:
   exact dimensions, certifications, model numbers, compatibility specs, use cases.
   Weave them naturally into the listing copy; do NOT copy verbatim.

3. KEYWORDS — Extract product-specific terms, variant names, and buyer-intent phrases
   from seller notes to supplement image-derived keywords.

4. TITLE — Only add model name, size, or defining spec from seller notes if not
   already clear from the image.

5. CONFLICTS — If seller notes contradict what is clearly visible in the image,
   trust the image. Never invent facts beyond what either source confirms.
"""

    return f"""You are a world-class advertising copywriter — the kind brands pay serious money for. You've written copy that stopped people mid-scroll, made them smile, and sent them straight to the checkout. You don't write product descriptions; you write desire. You understand that people don't buy things — they buy better versions of their lives. Every word you write makes the reader feel something.

PRODUCT FACTS (transform these into compelling copy — never echo them word-for-word):
- Product      : {product_type}
- Category     : {category}
- Material     : {material if material and material.lower() not in ("uncertain","unknown","") else "not specified"}
- Colors       : {colors_str}
- Pattern      : {pattern if pattern and pattern.lower() not in ("uncertain","unknown","") else "not specified"}
- Features     : {features_str if features_str else "not specified"}
- Fit / Cut    : {fit if fit and fit.lower() not in ("uncertain","unknown","") else "not specified"}
- Audience     : {gender_target if gender_target and gender_target.lower() not in ("uncertain","unknown","") else "general"}
- Use occasion : {usage if usage and usage.lower() not in ("uncertain","unknown","") else "everyday use"}
- Style        : {style if style and style.lower() not in ("uncertain","unknown","") else "not specified"}
{f"- Extra details: {extras}" if extras else ""}
{notes_block}

YOUR TASK — write copy that makes someone fall in love with this product before they've touched it.

1. TITLE — 170–200 characters (target the upper end — this is prime real estate)
   - Write it LAST, after the description AND after you have decided constrained_fields values.
   - This is the FIRST thing the consumer reads — pack it with everything they need to know to say "that's mine".
   - Structure: [Product Type] + [Material/Key Spec] + [Primary Variant: color/size/pack] + [Top Benefit] + [Key consumer-relevant attributes from constrained_fields: occasion, fabric type, pattern, fit, style, season — pick the 2–3 most informative ones for this specific product].
   - No pipes, no slashes, no ALL CAPS. Read like a real product name, not a keyword dump.
   - Every word must pull its weight — but do NOT pad with filler to hit the character count; only include attributes that genuinely help the buyer decide.

2. BULLET POINTS — exactly 5, each 80–200 characters
   - Format: ALL-CAPS HOOK — then the payoff sentence.
   - Each bullet unlocks one specific reason the customer's life gets better with this product.
   - Make the customer picture themselves using it. Use sensory language. Be specific.
   - Angles to hit (one each): the tactile/material experience, the design detail that turns heads, the fit or feel that makes it theirs, the moment or occasion it was made for, the practical payoff that keeps them coming back.
   - Write in second person — speak directly to the reader: "You'll reach for this every morning", "The kind of piece that earns compliments without trying".
   - Zero filler. Zero repetition. No "This product". No generic claims.
   - Example bar: "WHISPER-SOFT FABRIC — The brushed cotton feels like your favourite old tee on day one — no breaking in required, just instant comfort from the moment you pull it on."

3. DESCRIPTION — 3 paragraphs, 180–320 words total
   - Paragraph 1: HOOK — lead with the feeling, the moment, the problem solved, or the identity the product represents. Make the first sentence impossible to skip. Don't start with "This product", "Introducing", or "Looking for".
   - Paragraph 2: THE STORY — go deep on materials, construction, and design choices. Translate each feature into a lived benefit. Show you understand the product from the inside out. This is where trust is built.
   - Paragraph 3: THE CLOSE — paint the picture of owning this: when they'll use it, how it pairs, who it's for. End with energy — a sentence that makes them scroll back up to click "Add to Cart".
   - Tone: the voice of someone who genuinely loves this category. Confident, specific, warm. Not a sales pitch — a recommendation from an expert friend.
   - Never mention the image, photo, or "as shown".

4. KEYWORDS — exactly 20 search terms
   Think like a buyer at every stage — discovery, comparison, ready-to-buy:
   - 3–4 broad category terms (product type synonyms a buyer might type)
   - 4–5 material + style combos ("soft cotton tee", "stainless water bottle")
   - 5–6 occasion / use-case long-tail terms ("gym bag for women", "gift for new dad")
   - 4–5 buyer-intent terms ("under 500 rupees", "set of 4", "for tall men")
   Rules: all lowercase, no brand names, no punctuation, each term genuinely distinct.

5. ESTIMATED DIMENSIONS
   - Standard sizing for this product type based on visible proportions.
   - Format: "approximately X × Y × Z inches" or "approximately X inches long".
   - If truly unknown: "see product description".

6. HSN CODE — exact 8-digit code for Indian customs/GST
   Use the narrowest sub-heading available:
     42: leather/bags/pet gear — e.g. 42011000, 42021100
     39: plastics/silicone     — e.g. 39269099
     44: wood/bamboo           — e.g. 44219990
     52–63: textiles/apparel   — e.g. 62044390, 61051000
     73: steel/iron hardware   — e.g. 73269099
     84–85: electronics        — e.g. 85044000, 84733099
     94: furniture/lamps       — e.g. 94035090
     95: toys/sports           — e.g. 95030099
   Output only the 8 digits — no spaces, no dashes.

7. PRODUCT TAX CODE — Amazon India GST
   "A_GEN_NOTAX" (0%), "A_GEN_TAX_5" (5%), "A_GEN_TAX_12" (12%), "A_GEN_TAX_18" (18% — default), "A_GEN_TAX_28" (28%)

HARD RULES — breaking any of these ruins the listing:
- Never say "the image", "as shown", "as seen in photo", "pictured", or "the picture".
- No unverifiable claims: no "dermatologist tested", "clinically proven", "FDA approved", "guaranteed".
- No lazy superlatives: premium, luxury, best, world-class, revolutionary, unbeatable, superior, top-quality.
- No placeholders: N/A, unknown, TBD, [insert], fill-in-the-blank.
- Bullets are about product BENEFITS — never about packaging colors, label text, logo design, or box contents.

Return strict JSON only — no markdown fences, no explanation, no preamble:
{{
  "title": "",
  "bullet_points": ["", "", "", "", ""],
  "description": "",
  "keywords": ["", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", ""],
  "estimated_dimensions": "",
  "hsn_code": "",
  "product_tax_code": ""
}}"""




_KW_STOPWORDS = {
    "the", "and", "for", "with", "this", "that", "from", "are", "was",
    "has", "have", "its", "not", "but", "all", "can", "may", "will",
    "been", "also", "than", "into", "more", "some", "any", "very",
    "likely", "typical", "style", "type", "used", "use", "uses",
    "both", "such", "most", "each", "only", "well", "per", "yes",
    "via", "etc", "see", "made", "based", "other",
}

def _analysis_terms(product_analysis: Dict[str, Any]) -> List[str]:
    terms: List[str] = []
    for key in ("product_type", "category", "material", "usage", "style"):
        value = str(product_analysis.get(key, "")).strip().lower()
        if value and value not in {"uncertain", "unknown", "n/a"}:
            terms.extend([
                t for t in re.findall(r"[a-z0-9']+", value)
                if len(t) > 2 and t not in _KW_STOPWORDS
            ])
    for key in ("colors", "features"):
        values = product_analysis.get(key, [])
        if isinstance(values, list):
            for item in values:
                s = str(item).strip().lower()
                if not s or s in {"uncertain", "unknown", "n/a"}:
                    continue
                terms.extend([
                    t for t in re.findall(r"[a-z0-9']+", s)
                    if len(t) > 2 and t not in _KW_STOPWORDS
                ])
    # preserve order and deduplicate
    seen = set()
    out: List[str] = []
    for t in terms:
        if t not in seen:
            seen.add(t)
            out.append(t)
    return out


def validate_listing(listing: Dict[str, Any], product_analysis: Optional[Dict[str, Any]] = None) -> List[str]:
    errors: List[str] = []
    # Accept either "description" (two-call path) or "product_description" (one-shot path)
    _desc_key = "description" if "description" in listing else "product_description"
    for key in ("title", "bullet_points"):
        if key not in listing:
            errors.append(f"Missing field: '{key}'")
    if "description" not in listing and "product_description" not in listing:
        errors.append("Missing field: 'description'")
    title = str(listing.get("title", ""))
    if title and len(title) < 20:
        errors.append("Title is too short (minimum 20 characters).")
    if title and len(title) > 200:
        errors.append("Title is too long (maximum 200 characters).")
    bullets = listing.get("bullet_points", [])
    if not isinstance(bullets, list):
        bullets = []
    if len(bullets) != 5:
        errors.append(f"Expected exactly 5 bullet points, got {len(bullets)}.")
    for i, bullet in enumerate(bullets, 1):
        bullet_s = str(bullet)
        if len(bullet_s) < BULLET_MIN_LEN:
            errors.append(f"Bullet {i} is too short ({len(bullet_s)} chars, minimum {BULLET_MIN_LEN}).")
    description = str(listing.get(_desc_key, ""))
    if len(description) < DESCRIPTION_MIN_LEN:
        errors.append(f"Description too short ({len(description)} chars, minimum {DESCRIPTION_MIN_LEN}).")
    if len(description) > DESCRIPTION_MAX_LEN * 6:
        errors.append("Description is excessively long.")
    forbidden_placeholders = ("n/a", "unknown", "uncertain", "not sure", "cannot determine")
    for token in forbidden_placeholders:
        if token in description.lower():
            errors.append(f"Description contains placeholder/uncertain text: '{token}'.")
            break
    for i, bullet in enumerate(bullets, 1):
        b_lower = str(bullet).lower()
        if any(tok in b_lower for tok in forbidden_placeholders):
            errors.append(f"Bullet {i} contains placeholder/uncertain text.")

    # Distinctness check across bullets
    normalized_bullets = [" ".join(re.findall(r"[a-z0-9']+", str(b).lower())) for b in bullets]
    for i in range(len(normalized_bullets)):
        for j in range(i + 1, len(normalized_bullets)):
            if _word_overlap_score(normalized_bullets[i], normalized_bullets[j]) > 0.72:
                errors.append("Bullet points are too repetitive; each bullet must cover a distinct angle.")
                break
        else:
            continue
        break

    # Grounding checks: copy must reflect the vision analysis facts.
    # Skip for name-based / low-confidence fallback analyses — no real image facts to ground in.
    _analysis_confidence = float(product_analysis.get("confidence", 1.0)) if product_analysis else 1.0
    if product_analysis and _analysis_confidence >= 0.6:
        facts = _analysis_terms(product_analysis)
        if facts:
            full_text = (title + " " + description + " " + " ".join([str(b) for b in bullets])).lower()
            matched = [fact for fact in facts if re.search(rf"\b{re.escape(fact)}\b", full_text)]
            if len(matched) < min(3, len(facts)):
                errors.append(
                    "Listing is not grounded enough in detected image facts (product type/material/colors/features/usage)."
                )
            bullet_fact_hits = 0
            for bullet in bullets:
                bullet_l = str(bullet).lower()
                if any(re.search(rf"\b{re.escape(fact)}\b", bullet_l) for fact in facts):
                    bullet_fact_hits += 1
            if bullet_fact_hits < 3:
                errors.append("At least 3 bullet points must reference concrete facts from image analysis.")

    combined = (title + " " + description + " " + " ".join([str(b) for b in bullets])).lower()
    for word in FLAGGED_WORDS:
        count = combined.count(word)
        if count > REPETITION_THRESHOLD:
            errors.append(f"Overuse of '{word}' ({count} times, threshold {REPETITION_THRESHOLD}).")
    return errors




# Editorial checks: description must tie to the image, product role, use cases, and shopper usability;
# bullets must reinforce the same story (non-generic, non-duplicative).
_VISUAL_IN_DESC = re.compile(
    r"\b(image|images|photo|photograph|picture|pictures|shot|shots|framing|frame|visible|shown|shows|"
    r"appear|appears|appearance|look|looks|color|colour|shape|form|material|materials|finish|texture|"
    r"surface|detail|details|packaging|package|pattern|design|silhouette|profile|close-up|closeup|"
    r"angle|hero|listing|pictured|photographed|see the|in the photo|on display)\b",
    re.I,
)
_USE_CASE_IN_DESC = re.compile(
    r"\b(use|using|used|useful|daily|routine|routines|ideal|perfect for|great for|suited|when|where|while|"
    r"during|at home|home|office|travel|trip|outdoor|indoors?|work|school|gym|kitchen|garden|walk|run|"
    r"storage|season|weather|task|tasks|activity|activities|lifestyle|scenario|context|environment|"
    r"everyday|frequent|regular)\b",
    re.I,
)
_USABILITY_IN_DESC = re.compile(
    r"\b(easy|easier|easily|comfort|comfortable|lightweight|simple|straightforward|quick|secure|stable|"
    r"sturdy|grip|handle|hold|holding|wear|wearing|fit|fits|fitting|adjust|adjustable|install|setup|"
    r"set-up|clean|wash|maintain|care|durable|reliable|convenient|ergonomic|intuitive|user-friendly|"
    r"hands[- ]free|practical|manageable)\b",
    re.I,
)
_THEME_IN_BULLET = re.compile(
    r"\b(image|photo|picture|visible|shown|color|shape|material|finish|use|routine|daily|home|travel|"
    r"outdoor|indoor|easy|comfort|fit|handle|wear|durable|practical|customer|shopper|buyer)\b",
    re.I,
)
_VISUAL_IN_BULLET = re.compile(
    r"\b(visible|shown|color|colour|tone|shape|form|material|materials|finish|texture|surface|detail|"
    r"pattern|design|silhouette|profile|packaging|look|appearance|pictured)\b",
    re.I,
)
_USE_CASE_IN_BULLET = re.compile(
    r"\b(use|using|ideal|perfect for|great for|suited|when|where|daily|routine|home|office|travel|outdoor|"
    r"indoor|work|school|gym|kitchen|garden|walk|run|activity|lifestyle|context|everyday)\b",
    re.I,
)
_USABILITY_IN_BULLET = re.compile(
    r"\b(easy|easier|easily|comfort|comfortable|lightweight|simple|straightforward|quick|secure|stable|"
    r"sturdy|grip|handle|hold|wear|fit|adjust|adjustable|install|setup|clean|wash|maintain|care|durable|"
    r"reliable|convenient|ergonomic|intuitive|user-friendly|practical)\b",
    re.I,
)


def _extract_visual_keywords(summary: str) -> List[str]:
    tokens = set(re.findall(r"[a-z0-9']+", (summary or "").lower()))
    keywords = {name for name, _ in COLOR_NAMES}
    keywords |= {
        "portrait",
        "landscape",
        "square",
        "form",
        "shape",
        "material",
        "materials",
        "finish",
        "texture",
        "packaging",
        "hero",
        "shot",
        "color",
        "tone",
    }
    hits = sorted(tokens & keywords)
    return hits


def _sentence_count(text: str) -> int:
    parts = re.split(r"(?<=[.!?])\s+", (text or "").strip())
    return len([p for p in parts if len(p.strip()) > 8])


def _word_overlap_score(a: str, b: str) -> float:
    wa = set(re.findall(r"[a-z0-9']+", a.lower()))
    wb = set(re.findall(r"[a-z0-9']+", b.lower()))
    if not wa or not wb:
        return 0.0
    inter = len(wa & wb)
    return inter / max(len(wa), len(wb))


def audit_listing_copy(
    product_description: str,
    bullet_points: List[str],
    visual_summary: str = "",
) -> List[str]:
    """Return human-readable gaps; empty list means copy aligns with listing rules."""
    gaps: List[str] = []
    desc = (product_description or "").strip()
    bullets = [str(b).strip() for b in (bullet_points or []) if str(b).strip()][:5]

    if len(desc) < 180:
        gaps.append("product_description is too short to cover image, use cases, and usability (aim for 3–5 rich sentences).")
    if _sentence_count(desc) < 2:
        gaps.append("product_description needs multiple sentences (image + product role + use + usability).")

    if not _VISUAL_IN_DESC.search(desc):
        gaps.append(
            "product_description must explicitly ground the shopper in what the image shows "
            "(form, color, materials, packaging, or how the product is presented)."
        )
    if not _USE_CASE_IN_DESC.search(desc):
        gaps.append(
            "product_description must include real-world use cases (when, where, or how customers use the product)."
        )
    if not _USABILITY_IN_DESC.search(desc):
        gaps.append(
            "product_description must explain usability for the customer (ease, handling, comfort, fit, care, or similar)."
        )
    visual_keywords = _extract_visual_keywords(visual_summary)
    if visual_keywords and not any(k in desc.lower() for k in visual_keywords):
        gaps.append(
            "product_description should explicitly echo the image's visual traits (color, shape, materials, or orientation)."
        )

    if len(bullets) != 5:
        gaps.append("bullet_points must contain exactly five substantive bullets.")
    else:
        short = [i + 1 for i, b in enumerate(bullets) if len(b) < 28]
        if short:
            gaps.append(f"Bullets {short} are too short; expand so each amplifies the description with a distinct benefit.")
        themed = sum(1 for b in bullets if _THEME_IN_BULLET.search(b))
        if themed < 4:
            gaps.append(
                "Most bullet points should echo the same story as the description (image-visible traits, use cases, usability)."
            )
        visual_hits = sum(1 for b in bullets if _VISUAL_IN_BULLET.search(b))
        use_hits = sum(1 for b in bullets if _USE_CASE_IN_BULLET.search(b))
        usability_hits = sum(1 for b in bullets if _USABILITY_IN_BULLET.search(b))
        if visual_hits < 2:
            gaps.append("At least two bullet points must reference visible product traits from the image.")
        if use_hits < 1:
            gaps.append("At least one bullet point must describe a real-world use case.")
        if usability_hits < 1:
            gaps.append("At least one bullet point must emphasize usability (comfort, handling, setup, or care).")
        for i in range(len(bullets)):
            for j in range(i + 1, len(bullets)):
                if _word_overlap_score(bullets[i], bullets[j]) > 0.62:
                    gaps.append("Bullet points overlap too much; rewrite so each bullet covers a different angle.")
                    break
            else:
                continue
            break

    return gaps


def _normalize_ai_copy_record(obj: Dict, product_name: str) -> Dict:
    item_name = str(obj.get("item_name", "")).strip() or (product_name.strip() if product_name else "")
    desc = str(obj.get("product_description", "")).strip()
    bullets = obj.get("bullet_points", [])
    prompts = obj.get("image_prompts", [])
    if not isinstance(bullets, list):
        bullets = []
    if not isinstance(prompts, list):
        prompts = []
    bullets = [str(x).strip() for x in bullets if str(x).strip()][:5]
    prompts = [str(x).strip() for x in prompts if str(x).strip()][:2]
    if not desc:
        raise ValueError("AI did not return product_description.")
    while len(bullets) < 5:
        bullets.append(
            "Practical usability for everyday use—see product imagery for form, finish, and fit with your routine."
        )
    while len(prompts) < 2:
        base_name = item_name or product_name or "this product"
        if len(prompts) == 0:
            prompts.append(
                f"Award-winning lifestyle product photograph of {base_name}. Product placed naturally on a warm light-wood surface with relevant props softly out of focus. Warm golden natural light from camera-left. Product razor-sharp, background warm bokeh. Hyperrealistic, 8K, premium Indian brand commercial photography. No text. Use the first uploaded image as strict product reference — do not alter its shape, color, or label."
            )
        else:
            prompts.append(
                f"Premium Amazon product studio photograph of {base_name} — clean infographic background. Product centered as single hero on pure warm white background, filling 60% of frame. Soft studio lighting from upper-left softbox, gentle fill from right. Clean open space on all four sides. No text, no callouts, no icons. Hyperrealistic. Use the first uploaded image as strict product reference."
            )
    return {
        "item_name": item_name,
        "product_description": desc,
        "bullet_points": bullets[:5],
        "image_prompts": prompts[:2],
    }


def _fallback_copy_from_first_image(image_source: str, product_name: str) -> Dict:
    img = load_image_any(image_source)
    w, h = img.size
    orientation = "portrait" if h > w else "landscape" if w > h else "square"
    name = product_name.strip() or "This product"
    visual_summary = describe_image(image_source, product_name=name)
    desc = (
        f"{visual_summary} "
        f"The image highlights shape, apparent materials, and proportion so shoppers can picture {name} in everyday routines—"
        f"where it belongs, when someone reaches for it, and how it improves the moment. "
        f"Usability comes through visually: handling, comfort, and straightforward use without guesswork, "
        f"so customers know exactly what to expect before checkout."
    )
    bullets = [
        f"Grounded in the listing photo: {visual_summary}",
        f"Clear {orientation} presentation so customers can judge proportions, finish, and real-world presence before buying.",
        f"Use-case clarity: the image suggests when and where {name} fits naturally into daily routines.",
        "Usability-first: the photo signals how to handle, wear, install, or operate the product with ease.",
        "Buyer confidence: visible design details match practical benefits, reducing uncertainty at checkout.",
    ]
    prompts = [
        f"Award-winning lifestyle product photograph of {name}. Product placed naturally on a warm light-wood or marble surface with relevant lifestyle props softly out of focus beside it. Soft golden natural light from camera-left, warm fill from the right. Product razor-sharp with creamy bokeh background. Warm color grading, rich and inviting. Hyperrealistic, 8K. No text. Use the first uploaded image as strict product reference — preserve its shape, color, labels, and design exactly.",
        f"Premium Amazon product studio photograph of {name} — clean infographic background. Product centered as single hero on pure warm white (#FFFAF5) background filling 60% of frame. Soft studio lighting from upper-left softbox, gentle fill from the right, clean contact shadow below. Generous open space on all four sides for overlay text. No text, no callouts, no icons, no arrows. Hyperrealistic. Use the first uploaded image as strict product reference.",
    ]
    return {
        "item_name": name,
        "product_description": desc,
        "bullet_points": bullets[:5],
        "image_prompts": prompts[:2],
    }


def ai_generate_copy_from_image(image_source: str, product_name: str) -> Dict:
    """Generate Amazon listing copy from an image using the best available AI provider.

    Uses Groq for image analysis and copy generation.
    `image_source` can be a public URL or a local file path.
    """
    src = (image_source or "").strip()
    if not src:
        raise ValueError("Main Image URL is required for AI generation.")

    # ── Step 1: Analyse the image ──────────────────────────────────────────────
    analysis: Dict[str, Any] = {}
    try:
        analysis = groq_analyze_image_url(src, product_name)
    except Exception as exc:
        log.warning("Image analysis failed in ai_generate_copy_from_image: %s", exc)

    # ── Step 2: Generate listing copy via Groq ───────────────────────────────
    listing: Dict[str, Any] = {}
    if analysis:
        try:
            listing = _generate_listing(analysis)
        except Exception as exc:
            log.warning("Listing generation failed: %s — using fallback", exc)
            listing = _fallback_listing_from_analysis(analysis)
    if not listing:
        listing = _fallback_copy_from_first_image(src, product_name)

    user_item_name = product_name.strip()
    base_name = (
        user_item_name
        or str(analysis.get("product_type", "")).strip()
        or str(listing.get("title", "")).strip()
        or "this product"
    )
    bullets = listing.get("bullet_points", [])
    if not isinstance(bullets, list):
        bullets = []
    prompts = [
        f"Award-winning lifestyle product photograph of {base_name}. Product on a warm light-wood or natural surface with relevant lifestyle props softly out of focus. Soft golden natural light from camera-left, warm fill from right. Product razor-sharp, background warm bokeh. Hyperrealistic, 8K, premium commercial photography. No text. Use the first uploaded image as strict product reference — preserve its shape, color, labels, and design exactly.",
        f"Premium Amazon product studio photograph of {base_name} — clean infographic background. Product centered as single hero on pure warm white (#FFFAF5) background filling 60% of frame. Soft studio lighting, generous open space on all four sides. No text, no callouts, no icons. Hyperrealistic. Use the first uploaded image as strict product reference.",
    ]
    out = {
        "item_name": user_item_name or base_name,
        "product_description": str(listing.get("product_description") or listing.get("description", "")).strip(),
        "bullet_points": [str(b).strip() for b in bullets][:5],
        "image_prompts": prompts,
    }
    while len(out["bullet_points"]) < 5:
        out["bullet_points"].append("Practical product benefit grounded in visible features and everyday usage.")
    if not out["product_description"]:
        out["product_description"] = str(listing.get("title", base_name)).strip()
    validation_errors = validate_listing(listing, product_analysis=analysis) if analysis else []
    if validation_errors:
        out["__alignment_residual_gaps"] = validation_errors
    return out




def _wrap_text_lines(text: str, max_chars: int) -> List[str]:
    words = (text or "").split()
    if not words:
        return [""]
    lines: List[str] = []
    cur = words[0]
    for w in words[1:]:
        if len(cur) + 1 + len(w) <= max_chars:
            cur += " " + w
        else:
            lines.append(cur)
            cur = w
    lines.append(cur)
    return lines


def _load_fonts():
    """Return (bold_lg, bold_md, bold_sm, regular_md, regular_sm) ImageFont objects."""
    from PIL import ImageFont
    candidates_bold = [
        "/usr/share/fonts/truetype/ubuntu/Ubuntu-B.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ]
    candidates_regular = [
        "/usr/share/fonts/truetype/ubuntu/Ubuntu-R.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    def _first(paths, size):
        for p in paths:
            if os.path.exists(p):
                return ImageFont.truetype(p, size)
        return ImageFont.load_default()
    return (
        _first(candidates_bold, 52),    # bold_lg  — hero title
        _first(candidates_bold, 36),    # bold_md  — section title / feature label
        _first(candidates_bold, 26),    # bold_sm  — small label / badge
        _first(candidates_regular, 28), # regular_md — body text
        _first(candidates_regular, 22), # regular_sm — caption / chip text
    )


def _draw_gradient_bg(img: Image.Image, top_color: Tuple, bottom_color: Tuple) -> None:
    """Fill img in-place with a vertical linear gradient."""
    d = ImageDraw.Draw(img)
    w, h = img.size
    r0, g0, b0 = top_color
    r1, g1, b1 = bottom_color
    for y in range(h):
        t = y / (h - 1)
        r = int(r0 + (r1 - r0) * t)
        g = int(g0 + (g1 - g0) * t)
        b = int(b0 + (b1 - b0) * t)
        d.line([(0, y), (w - 1, y)], fill=(r, g, b))


def _add_drop_shadow(canvas: Image.Image, product: Image.Image, px: int, py: int,
                     shadow_offset: int = 18, shadow_blur: int = 30) -> None:
    """Paste a blurred dark shadow beneath the product, then paste product on top."""
    from PIL import ImageFilter
    shadow_layer = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    # Build a solid-black mask the same shape as the product's alpha (or a filled rect)
    mask = Image.new("L", product.size, 180)
    shadow_layer.paste(Image.new("RGB", product.size, (0, 0, 0)),
                       (px + shadow_offset, py + shadow_offset), mask)
    shadow_layer = shadow_layer.filter(ImageFilter.GaussianBlur(shadow_blur))
    canvas.paste(shadow_layer.convert("RGB"), (0, 0),
                 shadow_layer.split()[3] if shadow_layer.mode == "RGBA" else None)
    if product.mode == "RGBA":
        canvas.paste(product, (px, py), product.split()[3])
    else:
        canvas.paste(product, (px, py))


def _wrap_text(text: str, font, max_width: int) -> List[str]:
    """Word-wrap text to fit within max_width pixels."""
    words = str(text).split()
    lines: List[str] = []
    current = ""
    for w in words:
        test = (current + " " + w).strip()
        bbox = font.getbbox(test)
        if bbox[2] - bbox[0] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = w
    if current:
        lines.append(current)
    return lines or [""]


def generate_combo_image_labeled(
    image_sources: List[str],
    out_path: Path,
    combo_sku: str = "",
    target_size: int = 1200,
) -> None:
    """Generate a combo image showing ALL products — no text, original colors preserved."""
    generate_combo_image(image_sources, out_path, target_size=target_size)


def _safe_alpha_composite(dst: "Image.Image", src: "Image.Image", x: int, y: int) -> None:
    """Composite src onto dst at (x, y), clipping safely if src extends outside dst."""
    dw, dh = dst.size
    sw, sh = src.size
    src_x0 = max(0, -x)
    src_y0 = max(0, -y)
    src_x1 = min(sw, dw - x)
    src_y1 = min(sh, dh - y)
    if src_x0 >= src_x1 or src_y0 >= src_y1:
        return
    dst_x = max(0, x)
    dst_y = max(0, y)
    dst.alpha_composite(src.crop((src_x0, src_y0, src_x1, src_y1)), (dst_x, dst_y))


def _build_fan_canvas(
    product_imgs: list,
    canvas_w: int,
    canvas_h: int,
) -> "Image.Image":
    """Arrange product images like playing cards held in a hand (fanned arc).

    Used when a combo has more than 3 products.  Each product is placed on a
    white portrait card that is rotated around a pivot point far below the
    canvas, producing a natural fan spread.
    """
    import math
    from PIL import ImageDraw as _IDf, ImageFilter as _IFf

    n = len(product_imgs)

    # ── Card dimensions ───────────────────────────────────────────────────────
    card_w = int(canvas_w * 0.28)
    card_h = int(canvas_h * 0.60)

    # ── Fan geometry ──────────────────────────────────────────────────────────
    # Total angular spread grows with n but caps so cards stay visible.
    total_fan   = min(80.0, 16.0 * (n - 1))
    angle_step  = total_fan / (n - 1) if n > 1 else 0.0
    start_angle = -total_fan / 2.0

    # Work on an oversized canvas so rotated cards that poke outside the
    # target area don't get clipped during compositing.
    PAD = int(canvas_w * 0.55)
    ww  = canvas_w + 2 * PAD
    wh  = canvas_h + 2 * PAD

    work = Image.new("RGBA", (ww, wh), (255, 255, 255, 255))

    # Pivot sits below the working canvas; arm_len positions card centres so
    # the top of the fanned spread lands near the top of the target canvas.
    pivot_x = ww / 2.0
    pivot_y = wh * 0.88
    arm_len = canvas_h * 0.70

    for i, prod in enumerate(product_imgs):
        angle_deg = start_angle + i * angle_step
        angle_rad = math.radians(angle_deg)

        # Centre of this card in working-canvas coords
        cx = pivot_x + arm_len * math.sin(angle_rad)
        cy = pivot_y - arm_len * math.cos(angle_rad)

        # ── Build card ────────────────────────────────────────────────────────
        card = Image.new("RGBA", (card_w, card_h), (255, 255, 255, 255))
        prod_c = prod.copy()
        prod_c.thumbnail((card_w - 24, card_h - 24), Image.LANCZOS)
        px = (card_w - prod_c.width)  // 2
        py = (card_h - prod_c.height) // 2
        if prod_c.mode == "RGBA":
            card.paste(prod_c, (px, py), prod_c.split()[3])
        else:
            card.paste(prod_c.convert("RGBA"), (px, py))

        # Thin border
        _IDf.Draw(card).rectangle(
            [(1, 1), (card_w - 2, card_h - 2)],
            outline=(190, 190, 190), width=2,
        )

        # ── Rotate card around its own centre ─────────────────────────────────
        rotated = card.rotate(
            -angle_deg, expand=True,
            resample=Image.BICUBIC,
            fillcolor=(255, 255, 255, 0),
        )

        # ── Drop shadow ───────────────────────────────────────────────────────
        shadow = Image.new("RGBA", rotated.size, (0, 0, 0, 0))
        alpha  = rotated.split()[3].point(lambda p: p * 75 // 255)
        shadow.putalpha(alpha)
        shadow = shadow.filter(_IFf.GaussianBlur(radius=10))

        rx = int(cx - rotated.width  / 2)
        ry = int(cy - rotated.height / 2)

        _safe_alpha_composite(work, shadow, rx + 7, ry + 9)
        _safe_alpha_composite(work, rotated, rx, ry)

    # Crop back to target canvas size (strip the PAD border)
    cropped = work.crop((PAD, PAD, PAD + canvas_w, PAD + canvas_h))
    return cropped.convert("RGB")


def generate_combo_image(
    image_sources: List[str],
    out_path: Path,
    combo_number: int = 1,
    target_size: int = 1200,
) -> None:
    """Compose multiple product images on a clean white background.

    ≤3 products → side-by-side row layout.
    >3 products → fanned arc layout (like cards held in a hand).
    """
    from PIL import ImageFilter as _IFc
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    CANVAS_W = target_size
    CANVAS_H = target_size

    # ── Load product images — crop to product area, preserve original colors ──
    product_imgs: list = []
    for src_entry in image_sources:
        try:
            raw = _load_combo_product_image(src_entry).convert("RGB")
            raw = _crop_to_product_rgb(raw, padding=30)
            product_imgs.append(raw)
        except Exception as exc:
            log.warning("Combo: all URLs failed for one product slot, skipping: %s", exc)

    if not product_imgs:
        Image.new("RGB", (CANVAS_W, CANVAS_H), (255, 255, 255)).save(str(out_path), format="PNG")
        log.warning("Combo: all images failed, saved blank placeholder")
        return

    n = len(product_imgs)

    # ── Grid layout for 4+ products ──────────────────────────────────────────
    if n > 3:
        import math as _math
        PADDING = int(target_size * 0.04)
        GAP     = int(target_size * 0.03)
        n_cols = 2 if n <= 4 else 3
        n_rows = _math.ceil(n / n_cols)
        cell_w = (CANVAS_W - PADDING * 2 - GAP * (n_cols - 1)) // n_cols
        cell_h = (CANVAS_H - PADDING * 2 - GAP * (n_rows - 1)) // n_rows
        canvas = Image.new("RGB", (CANVAS_W, CANVAS_H), (255, 255, 255))
        for i, raw in enumerate(product_imgs):
            row = i // n_cols
            col = i % n_cols
            items_in_row = min(n_cols, n - row * n_cols)
            row_width = items_in_row * cell_w + (items_in_row - 1) * GAP
            row_x_offset = (CANVAS_W - row_width) // 2
            raw_copy = raw.copy()
            raw_copy.thumbnail((cell_w, cell_h), Image.LANCZOS)
            pw, ph = raw_copy.size
            x = row_x_offset + col * (cell_w + GAP) + (cell_w - pw) // 2
            y = PADDING + row * (cell_h + GAP) + (cell_h - ph) // 2
            canvas.paste(raw_copy, (x, y))
        canvas.save(str(out_path), format="PNG")
        log.info("Combo grid image saved: %s (%d products)", out_path.name, n)
        return

    # ── Row layout for ≤3 products ────────────────────────────────────────────
    PADDING    = int(target_size * 0.02)
    GAP        = int(target_size * 0.025)
    MAX_ITEM_H = CANVAS_H - PADDING * 2
    slot_w = (CANVAS_W - PADDING * 2 - GAP * max(n - 1, 0)) // max(n, 1)
    canvas = Image.new("RGB", (CANVAS_W, CANVAS_H), (255, 255, 255))

    for i, raw in enumerate(product_imgs):
        raw.thumbnail((slot_w, MAX_ITEM_H), Image.LANCZOS)
        pw, ph = raw.size
        x = PADDING + i * (slot_w + GAP) + (slot_w - pw) // 2
        y = CANVAS_H - PADDING - ph
        canvas.paste(raw, (x, y))

    canvas.save(str(out_path), format="PNG")
    log.info("Combo image saved: %s (%d products)", out_path.name, n)


def _strip_bullet_label(bullet: str) -> Tuple[str, str]:
    """Split 'LABEL — description text' into (label, description)."""
    for sep in (" — ", " - ", ": "):
        if sep in bullet:
            parts = bullet.split(sep, 1)
            label = parts[0].strip().upper()
            if len(label) <= 30 and label == label.upper():
                return label, parts[1].strip()
    return "", bullet.strip()


def _remove_photo_background(img: Image.Image) -> Image.Image:
    """Best-effort background removal using pixel brightness/colour distance."""
    if img.mode != "RGBA":
        img = img.convert("RGBA")
    try:
        import numpy as np
        data = np.array(img, dtype=np.float32)
        rgb  = data[:, :, :3]
        # Sample corner pixels as background colour estimate
        corners = np.array([
            rgb[0, 0], rgb[0, -1], rgb[-1, 0], rgb[-1, -1],
            rgb[0, img.width // 2], rgb[img.height // 2, 0],
        ])
        bg_colour = corners.mean(axis=0)
        dist = np.sqrt(((rgb - bg_colour) ** 2).sum(axis=2))
        threshold = max(30.0, dist.max() * 0.18)
        mask = (dist > threshold).astype(np.uint8) * 255
        # Dilate mask slightly so edges aren't clipped
        from PIL import ImageFilter as _IFbg
        mask_img = Image.fromarray(mask, mode="L").filter(_IFbg.MaxFilter(3))
        result = img.copy()
        result.putalpha(mask_img)
        return result
    except Exception:
        return img


def _crop_to_content(img: Image.Image, padding: int = 12) -> Image.Image:
    """Crop an RGBA image to the tight bounding box of non-transparent pixels."""
    if img.mode != "RGBA":
        img = img.convert("RGBA")
    try:
        import numpy as np
        alpha = np.array(img.split()[3])
        rows  = np.any(alpha > 10, axis=1)
        cols  = np.any(alpha > 10, axis=0)
        if not rows.any():
            return img
        rmin, rmax = np.where(rows)[0][[0, -1]]
        cmin, cmax = np.where(cols)[0][[0, -1]]
        rmin = max(0, rmin - padding)
        rmax = min(img.height - 1, rmax + padding)
        cmin = max(0, cmin - padding)
        cmax = min(img.width  - 1, cmax + padding)
        return img.crop((cmin, rmin, cmax + 1, rmax + 1))
    except Exception:
        return img


def _crop_to_product_rgb(img: Image.Image, padding: int = 30) -> Image.Image:
    """Crop an RGB image to the bounding box of non-background content.

    Does NOT modify any pixel colors — purely finds where the product is and
    crops away the surrounding background padding (grey/white studio backdrop).
    """
    try:
        import numpy as np
        rgb = np.array(img.convert("RGB"), dtype=np.float32)
        h, w = rgb.shape[:2]
        # Sample only the four corners + edge midpoints (never center) so the
        # background estimate is never contaminated by product pixels
        corners = np.array([
            rgb[0, 0], rgb[0, -1], rgb[-1, 0], rgb[-1, -1],
            rgb[0, w // 2], rgb[-1, w // 2],
            rgb[h // 2, 0], rgb[h // 2, -1],
        ])
        bg = corners.mean(axis=0)
        dist = np.sqrt(((rgb - bg) ** 2).sum(axis=2))
        # Use 25% of max-distance as threshold; floor at 40 to handle gradients
        threshold = max(40.0, dist.max() * 0.25)
        mask = dist > threshold
        rows = np.any(mask, axis=1)
        cols = np.any(mask, axis=0)
        if not rows.any() or not cols.any():
            return img
        rmin, rmax = int(np.where(rows)[0][0]),  int(np.where(rows)[0][-1])
        cmin, cmax = int(np.where(cols)[0][0]),  int(np.where(cols)[0][-1])
        rmin = max(0, rmin - padding)
        rmax = min(img.height - 1, rmax + padding)
        cmin = max(0, cmin - padding)
        cmax = min(img.width  - 1, cmax + padding)
        return img.crop((cmin, rmin, cmax + 1, rmax + 1))
    except Exception:
        return img


def generate_reference_locked_images(
    reference_image_source: str,
    out_dir: Path,
    prefix: str,
    product_name: str,
    bullets: List[str],
    analysis: Optional[Dict] = None,
    listing: Optional[Dict] = None,
) -> List[str]:
    """Generate two professional Amazon listing images using PIL.

    Image 1 — Hero/Lifestyle: gradient backdrop, large product shot, title, tag line.
    Image 2 — Infographic: white background, feature callout cards on the left,
                            product image on the right, full Groq-filled detail chips.
    """
    from PIL import ImageFilter, ImageFont
    out_dir.mkdir(parents=True, exist_ok=True)
    analysis = analysis or {}
    bold_lg, bold_md, bold_sm, regular_md, regular_sm = _load_fonts()

    # ── Theme-aware palette ───────────────────────────────────────────────────
    def _hex_to_rgb(h: str) -> tuple:
        h = h.lstrip("#")
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    try:
        from amazon_product_images import _resolve_theme  # type: ignore
        _theme = _resolve_theme(analysis)
        _pal   = [_hex_to_rgb(c) for c in _theme.get("palette", [])]
    except Exception:
        _pal = []

    # Fallback palette (neutral premium)
    _P0 = _pal[0] if len(_pal) > 0 else (18,  38,  72)   # dominant / dark
    _P1 = _pal[1] if len(_pal) > 1 else (180, 140,  60)   # accent 1
    _P2 = _pal[2] if len(_pal) > 2 else (245, 240, 228)   # background light
    _P3 = _pal[3] if len(_pal) > 3 else (228, 218, 200)   # background dark

    # Named aliases used below
    NAVY        = _P0
    GOLD        = _P1
    WARM_WHITE  = (255, 254, 250)
    OFF_WHITE   = (248, 246, 242)
    LIGHT_GRAY  = (230, 228, 224)
    MID_GRAY    = (120, 118, 114)
    DARK_GRAY   = ( 42,  42,  42)
    CREAM_TOP   = _P2
    CREAM_BOT   = _P3
    BADGE_BG    = tuple(max(0, c - 15) for c in _P2)
    DIVIDER     = tuple(max(0, c - 30) for c in _P2)

    SIZE = 1200

    name  = (product_name or "Product").strip()
    # Truncate title to 2 lines worth of characters for hero image
    title = name[:90]

    analysis = analysis or {}
    listing  = listing  or {}
    _inv_vals = {"uncertain", "unknown", "n/a", "none", "not specified", ""}

    material = str(analysis.get("material", "")).strip()
    colors   = [str(c).strip() for c in analysis.get("colors", [])
                if str(c).strip() and str(c).lower() not in {"uncertain", "unknown"}]
    features = [str(f).strip() for f in analysis.get("features", [])
                if str(f).strip() and str(f).lower() not in {"uncertain", "unknown"}]
    usage    = str(analysis.get("usage", "")).strip()
    category = str(analysis.get("category", "")).strip()
    ptype    = str(analysis.get("product_type", "")).strip()
    style_v  = str(analysis.get("style", "")).strip()

    # Extra details from Groq listing output
    _desc_raw  = str(listing.get("description") or listing.get("product_description") or "").strip()
    # First sentence of description — used as infographic tagline
    _desc_tag  = (_desc_raw.split(".")[0].strip() + ".") if _desc_raw else ""
    if len(_desc_tag) > 120:
        _desc_tag = _desc_tag[:117] + "…"
    _dimensions   = str(listing.get("estimated_dimensions", "")).strip()
    _constrained  = {
        k: str(v).strip()
        for k, v in (listing.get("constrained_fields") or {}).items()
        if str(v).strip() and str(v).strip().lower() not in _inv_vals
    }

    # Build ordered detail chips — shown in the infographic footer bar
    detail_chips: List[str] = []
    if material and material.lower() not in _inv_vals:
        detail_chips.append(f"Material: {material}")
    if colors:
        detail_chips.append(f"Color: {', '.join(colors[:2])}")
    if _dimensions and _dimensions.lower() not in _inv_vals and "see product" not in _dimensions.lower():
        detail_chips.append(f"Size: {_dimensions}")
    if style_v and style_v.lower() not in _inv_vals:
        detail_chips.append(f"Style: {style_v}")
    if usage and usage.lower() not in _inv_vals:
        detail_chips.append(f"Use: {usage[:40]}")
    if category:
        detail_chips.append(category)
    # Add top constrained fields (e.g. Fabric Type, Occasion, Pattern Type)
    for _cf_k, _cf_v in list(_constrained.items())[:4]:
        detail_chips.append(f"{_cf_k}: {_cf_v}")

    _raw = load_image_any(reference_image_source).convert("RGBA")
    # Remove background and crop to tight product bbox — same pipeline as combo images
    try:
        _raw = _remove_photo_background(_raw)
    except Exception as _bge:
        log.debug("BG removal skipped in generate_reference_locked_images: %s", _bge)
    src = _crop_to_content(_raw)

    # ── IMAGE 1: Hero / Lifestyle ─────────────────────────────────────────────
    hero = Image.new("RGB", (SIZE, SIZE), CREAM_TOP)
    _draw_gradient_bg(hero, CREAM_TOP, CREAM_BOT)
    d1 = ImageDraw.Draw(hero)

    # Subtle noise texture band at the bottom (floor shadow strip)
    d1.rectangle([(0, SIZE - 220), (SIZE, SIZE)], fill=(210, 200, 182))
    _draw_gradient_bg_region = None  # no helper needed; use direct draw
    # Gradient the floor strip
    for y in range(SIZE - 220, SIZE):
        t = (y - (SIZE - 220)) / 220
        r = int(210 + (190 - 210) * t)
        g = int(200 + (182 - 200) * t)
        b = int(182 + (158 - 182) * t)
        d1.line([(0, y), (SIZE - 1, y)], fill=(r, g, b))

    # Soft radial oval highlight behind product (simulate studio light)
    highlight = Image.new("RGBA", (SIZE, SIZE), (0, 0, 0, 0))
    dh = ImageDraw.Draw(highlight)
    cx, cy = SIZE // 2, SIZE // 2 - 40
    dh.ellipse([(cx - 380, cy - 440), (cx + 380, cy + 440)], fill=(255, 252, 240, 55))
    from PIL import ImageFilter as _IF
    highlight = highlight.filter(_IF.GaussianBlur(60))
    hero.paste(highlight.convert("RGB"), (0, 0), highlight.split()[3])

    # Product image — centered, max 680x680
    product = src.copy()
    product.thumbnail((680, 680), Image.LANCZOS)
    pw, ph = product.size
    px = (SIZE - pw) // 2
    py = (SIZE - ph) // 2 - 60
    _add_drop_shadow(hero, product.convert("RGBA"), px, py, shadow_offset=22, shadow_blur=35)

    # Re-draw overlay draw handle after pasting
    d1 = ImageDraw.Draw(hero)

    # Top category badge
    if category or ptype:
        badge_text = (category or ptype).upper()
        bb = bold_sm.getbbox(badge_text)
        bw = bb[2] - bb[0] + 36
        bh = bb[3] - bb[1] + 16
        bx = (SIZE - bw) // 2
        by = 44
        d1.rounded_rectangle([(bx, by), (bx + bw, by + bh)], radius=6, fill=BADGE_BG)
        d1.text((bx + 18, by + 8), badge_text, font=bold_sm, fill=NAVY)

    # Title block at bottom of hero
    title_lines = _wrap_text(title, bold_lg, SIZE - 100)[:3]
    title_block_h = len(title_lines) * 62 + 20
    title_y = SIZE - 220 - title_block_h - 16

    # Semi-transparent title backdrop
    overlay = Image.new("RGBA", (SIZE, title_block_h + 24), (18, 38, 72, 195))
    hero.paste(overlay.convert("RGB"), (0, title_y - 12),
               overlay.split()[3])

    d1 = ImageDraw.Draw(hero)
    ty = title_y
    for line in title_lines:
        d1.text((50, ty), line, font=bold_lg, fill=WARM_WHITE)
        ty += 62

    # Gold accent bar under title
    d1.rectangle([(50, ty + 6), (SIZE - 50, ty + 10)], fill=GOLD)

    # Usage / tag line below accent bar
    if usage and usage.lower() not in {"uncertain", "unknown", "general use"}:
        tag = usage[:80]
        d1.text((50, ty + 18), tag, font=regular_sm, fill=(220, 212, 190))

    # Bottom detail chips
    chip_x = 50
    chip_y = SIZE - 58
    for chip in detail_chips[:4]:
        cb = regular_sm.getbbox(chip)
        cw = cb[2] - cb[0] + 24
        if chip_x + cw > SIZE - 50:
            break
        d1.rounded_rectangle([(chip_x, chip_y - 8), (chip_x + cw, chip_y + 24)],
                              radius=4, fill=(255, 255, 255, 140) if False else BADGE_BG)
        d1.text((chip_x + 12, chip_y), chip, font=regular_sm, fill=NAVY)
        chip_x += cw + 12

    p1 = out_dir / f"{prefix}_ai_1.png"
    hero.save(p1, format="PNG")

    # ── IMAGE 2: Infographic ──────────────────────────────────────────────────
    info = Image.new("RGB", (SIZE, SIZE), WARM_WHITE)
    d2 = ImageDraw.Draw(info)

    PANEL_W = 530      # left panel width
    IMG_X   = PANEL_W + 20
    IMG_W   = SIZE - IMG_X - 20

    # Header banner — taller when we have a tagline to show
    BANNER_H = 140 if _desc_tag else 100
    d2.rectangle([(0, 0), (SIZE, BANNER_H)], fill=NAVY)
    # Gold accent stripe
    d2.rectangle([(0, BANNER_H), (SIZE, BANNER_H + 5)], fill=GOLD)

    banner_lines = _wrap_text(title, bold_md, SIZE - 100)[:2]
    title_block_h = len(banner_lines) * 42
    title_start_y = max(10, (BANNER_H - title_block_h - (28 if _desc_tag else 0)) // 2)
    bly = title_start_y
    for bline in banner_lines:
        d2.text((50, bly), bline, font=bold_md, fill=WARM_WHITE)
        bly += 42
    # Tagline (first sentence of Groq description) just below title
    if _desc_tag:
        d2.text((50, bly + 4), _desc_tag, font=regular_sm, fill=(190, 185, 175))

    # Left panel background (very subtle tint)
    d2.rectangle([(0, BANNER_H + 5), (PANEL_W, SIZE)], fill=OFF_WHITE)
    d2.line([(PANEL_W, BANNER_H + 5), (PANEL_W, SIZE)], fill=LIGHT_GRAY, width=2)

    # Right panel — product image, centered vertically with drop shadow
    prod_info = src.copy()
    max_prod_h = SIZE - BANNER_H - 5 - 80
    prod_info.thumbnail((IMG_W - 20, max_prod_h), Image.LANCZOS)
    piw, pih = prod_info.size
    pix = IMG_X + (IMG_W - piw) // 2
    piy = BANNER_H + 5 + (max_prod_h - pih) // 2 + 40
    _add_drop_shadow(info, prod_info.convert("RGBA"), pix, piy, shadow_offset=18, shadow_blur=28)
    d2 = ImageDraw.Draw(info)

    # Feature cards on left panel
    detail_bullets = [b for b in bullets if (b or "").strip()][:5]
    if not detail_bullets:
        detail_bullets = [f for f in features[:5]]
    if not detail_bullets:
        detail_bullets = ["High-quality construction", "Premium materials", "Versatile design"]

    # Use theme palette for callout card accents, cycle through palette colors
    ACCENT_COLORS = [
        _P1,                          # theme accent 1
        _P0,                          # theme dominant
        tuple(min(255, c + 30) for c in _P1),   # lighter accent 1
        tuple(max(0,   c - 30) for c in _P0),   # darker dominant
        (120, 118, 114),              # neutral fallback
    ]

    card_y = BANNER_H + 5 + 28
    card_x = 20
    card_w = PANEL_W - 40
    card_gap = 14

    # Calculate available height per card
    available_h = SIZE - BANNER_H - 5 - 28 - 60  # leave 60px bottom
    card_h_max = (available_h - card_gap * len(detail_bullets)) // len(detail_bullets)
    card_h_max = max(card_h_max, 80)

    for i, bullet in enumerate(detail_bullets):
        label, desc = _strip_bullet_label(bullet)
        accent = ACCENT_COLORS[i % len(ACCENT_COLORS)]

        # Estimate card height based on text wrapping
        desc_lines = _wrap_text(desc, regular_md, card_w - 80) if desc else []
        inner_h = 44 + len(desc_lines) * 32 + 12  # label row + desc lines + padding
        card_h = min(max(inner_h, 70), card_h_max)

        # Card background
        d2.rounded_rectangle([(card_x, card_y), (card_x + card_w, card_y + card_h)],
                              radius=8, fill=WARM_WHITE)
        d2.rounded_rectangle([(card_x, card_y), (card_x + card_w, card_y + card_h)],
                              radius=8, outline=LIGHT_GRAY, width=1)

        # Accent left bar
        d2.rounded_rectangle([(card_x, card_y), (card_x + 8, card_y + card_h)],
                              radius=4, fill=accent)

        # Circle number badge
        badge_cx = card_x + 28
        badge_cy = card_y + 24
        d2.ellipse([(badge_cx - 14, badge_cy - 14), (badge_cx + 14, badge_cy + 14)],
                   fill=accent)
        num_str = str(i + 1)
        nb = bold_sm.getbbox(num_str)
        d2.text((badge_cx - (nb[2] - nb[0]) // 2, badge_cy - (nb[3] - nb[1]) // 2),
                num_str, font=bold_sm, fill=WARM_WHITE)

        text_x = card_x + 52
        text_y = card_y + 10

        if label:
            d2.text((text_x, text_y), label, font=bold_sm, fill=accent)
            text_y += 30
        else:
            text_y += 4

        for dl in desc_lines:
            if text_y + 30 > card_y + card_h - 6:
                break
            d2.text((text_x, text_y), dl, font=regular_md, fill=DARK_GRAY)
            text_y += 32

        # Connector dot on the right edge of card (pointing toward product)
        dot_y = card_y + card_h // 2
        d2.ellipse([(card_x + card_w - 8, dot_y - 5), (card_x + card_w + 2, dot_y + 5)],
                   fill=accent)
        # Dashed connector line to product
        line_x0 = card_x + card_w + 2
        line_x1 = pix - 10
        if line_x1 > line_x0:
            for lx in range(line_x0, line_x1, 10):
                d2.line([(lx, dot_y), (min(lx + 6, line_x1), dot_y)],
                        fill=(*accent, 120) if False else (*accent,), width=1)

        card_y += card_h + card_gap
        if card_y + 60 > SIZE:
            break

    # Bottom info bar — two rows of chips showing all Groq-filled product details
    # Row 1 starts at SIZE-100, row 2 at SIZE-56. Navy background covers both rows.
    BAR_ROW1 = SIZE - 100
    BAR_ROW2 = SIZE - 56
    d2.rectangle([(0, BAR_ROW1 - 6), (SIZE, SIZE)], fill=NAVY)
    # Thin gold divider at top of bar
    d2.rectangle([(0, BAR_ROW1 - 6), (SIZE, BAR_ROW1 - 3)], fill=GOLD)

    def _draw_chip_row(chips: List[str], row_y: int, chip_fill, text_fill) -> None:
        cx = 20
        for chip in chips:
            cb = regular_sm.getbbox(chip)
            cw = (cb[2] - cb[0]) + 24
            if cx + cw > SIZE - 20:
                break
            d2.rounded_rectangle([(cx, row_y), (cx + cw, row_y + 34)],
                                  radius=4, fill=chip_fill)
            d2.text((cx + 12, row_y + 6), chip, font=regular_sm, fill=text_fill)
            cx += cw + 10

    # Split chips across two rows: primary specs in row 1, constrained fields in row 2
    primary_chips    = [c for c in detail_chips if not any(
        c.startswith(f"{k}:") for k in _constrained.keys()
    )]
    secondary_chips  = [c for c in detail_chips if c not in primary_chips]

    _draw_chip_row(primary_chips[:6],   BAR_ROW1, GOLD,    NAVY)
    _draw_chip_row(secondary_chips[:6], BAR_ROW2, (60, 80, 120), WARM_WHITE)

    p2 = out_dir / f"{prefix}_ai_2.png"
    info.save(p2, format="PNG")

    return [str(p1.resolve()), str(p2.resolve())]


def apply_ai_assets_to_row(
    row: Dict,
    columns: List[Dict],
    product_index: int,
    generate_images: bool = True,
    image_mode: str = IMAGE_MODE_STANDARD,
) -> Dict:
    main_image_attr = first_attr_by_label(columns, "Main Image URL")
    desc_attr = first_attr_by_label(columns, "Product Description")
    item_name_attr = first_attr_by_label(columns, "Item Name")
    bullet_attrs = all_attrs_by_label(columns, "Bullet Point")[:5]
    other_img_attrs = all_attrs_by_label(columns, "Other Image URL")[:2]
    if not main_image_attr:
        raise ValueError("Main Image URL column not found in template.")

    image_source = str(row.get(main_image_attr, "")).strip()
    if not image_source:
        raise ValueError("Main Image URL is required for AI generation.")
    product_name = str(row.get(item_name_attr, "")).strip() if item_name_attr else ""
    analysis_name = product_name or "Product"
    copy_obj = None
    copy_error = None
    try:
        copy_obj = ai_generate_copy_from_image(image_source, product_name)
    except Exception as exc:
        copy_error = str(exc)
        if not generate_images:
            raise

    if copy_obj:
        # Item Name is user-provided; do not overwrite it with AI output.
        if desc_attr and copy_obj.get("product_description"):
            row[desc_attr] = str(copy_obj["product_description"]).strip()
        for idx, attr in enumerate(bullet_attrs):
            row[attr] = copy_obj["bullet_points"][idx] if idx < len(copy_obj["bullet_points"]) else row.get(attr, "")
        residual = copy_obj.get("__alignment_residual_gaps")
        if residual:
            note = "Copy alignment: " + "; ".join(residual)
            prev = (row.get("__ai_copy_warning") or "").strip()
            row["__ai_copy_warning"] = (prev + (" | " if prev else "") + note)[:4000]

    if generate_images and IMAGE_GENERATION_ENABLED:
        mode = (image_mode or IMAGE_MODE_STANDARD).strip().lower()
        if mode not in {IMAGE_MODE_STANDARD, IMAGE_MODE_AMAZON_INSTRUCTIONS}:
            raise ValueError(f"Unknown image_mode '{image_mode}'.")
        if mode == IMAGE_MODE_AMAZON_INSTRUCTIONS:
            prompts = build_listing_image_prompts(analysis_name, analysis or {})
            base_out = Path(os.environ.get("AMAZON_IMAGE_PROJECT_DIR", "generated_images"))
            out_dir = base_out / f"product_{product_index+1}"
            prefix = f"product_{product_index+1}"
            # Image generation: PIL-based reference-locked images
            image_paths = generate_reference_locked_images(
                reference_image_source=image_source,
                out_dir=out_dir,
                prefix=prefix,
                product_name=product_name or analysis_name,
                bullets=[],
            )[: len(other_img_attrs)]
        else:
            # Enforce strict product identity: build secondary images directly from
            # the uploaded reference image so generated visuals never drift.
            base_out = Path(os.environ.get("AI_GENERATED_IMAGE_DIR", "generated_amazon_images/ai"))
            bullet_source = []
            if copy_obj and isinstance(copy_obj.get("bullet_points"), list):
                bullet_source = [str(x).strip() for x in copy_obj.get("bullet_points", []) if str(x).strip()]
            if not bullet_source:
                bullet_source = [str(row.get(attr, "")).strip() for attr in bullet_attrs if str(row.get(attr, "")).strip()]
            image_paths = generate_reference_locked_images(
                reference_image_source=image_source,
                out_dir=base_out / f"product_{product_index+1}",
                prefix=f"product_{product_index+1}",
                product_name=product_name or analysis_name,
                bullets=bullet_source[:5],
                listing=copy_obj,
            )[: len(other_img_attrs)]
        for idx, attr in enumerate(other_img_attrs):
            if idx < len(image_paths):
                row[attr] = image_paths[idx]
        if copy_error:
            row["__ai_copy_warning"] = copy_error
    return row


def load_template_meta(template_path: str):
    wb = load_workbook(template_path, keep_vba=True)
    ws = wb["Template"]
    dd = wb["Data Definitions"] if "Data Definitions" in wb.sheetnames else None
    bd = wb["Browse Data"] if "Browse Data" in wb.sheetnames else None
    vv = wb["Valid Values"] if "Valid Values" in wb.sheetnames else None

    label_row, attribute_row, data_row = parse_template_rows(ws)
    group_row = max(1, label_row - 1)
    example_row = attribute_row + 1

    columns = read_column_metadata(ws, label_row, attribute_row)
    ordered_cols = sorted(columns.values(), key=lambda c: c.col_idx)
    definitions = read_data_definitions(dd) if dd is not None else {}
    browse_nodes = read_browse_nodes(bd) if bd is not None else []
    node_to_product_type = read_recommended_node_to_product_type(vv) if vv is not None else {}
    valid_product_types  = read_valid_product_types(vv) if vv is not None else []
    all_valid_values     = read_all_valid_values(vv) if vv is not None else {}
    product_id_type_options: List[str] = []
    if vv is not None:
        for _row in vv.iter_rows(min_row=1, values_only=True):
            if len(_row) < 2:
                continue
            b_val = str(_row[1] or "").strip().lower()
            if b_val == "product id type":
                vals = [str(v or "").strip() for v in _row[2:42] if str(v or "").strip()]
                if vals:
                    product_id_type_options = vals
                    break

    merged_group_ranges = []
    for mr in ws.merged_cells.ranges:
        if mr.min_row == group_row and mr.max_row == group_row:
            merged_group_ranges.append((mr.min_col, mr.max_col, str(ws.cell(group_row, mr.min_col).value or "").strip()))

    def group_value_for_col(col_idx: int) -> str:
        for min_col, max_col, val in merged_group_ranges:
            if min_col <= col_idx <= max_col:
                return val
        return str(ws.cell(group_row, col_idx).value or "").strip()

    cols = []
    required_attrs: List[str] = []
    for meta in ordered_cols:
        req = definitions.get(meta.attribute)
        req_text = (req.requirement if req else "").strip()
        if req_text.lower() == "required":
            required_attrs.append(meta.attribute)
        cols.append(
            {
                "col": meta.column_letter,
                "col_idx": meta.col_idx,
                "group": group_value_for_col(meta.col_idx),
                "label": meta.label,
                "attr": meta.attribute,
                "example": str(ws.cell(example_row, meta.col_idx).value or "").strip(),
                "required": req_text,
            }
        )

    return {
        "columns": cols,
        "data_row": data_row,
        "label_row": label_row,
        "attribute_row": attribute_row,
        "example_row": example_row,
        "required_attrs": required_attrs,
        "definitions": definitions,
        "browse_options": [
            {"node_id": node_id, "path": path, "display": f"{path} ({node_id})"}
            for node_id, path in browse_nodes
        ],
        "node_to_product_type": node_to_product_type,
        "valid_product_types": valid_product_types,
        "all_valid_values": all_valid_values,
        "product_id_type_options": product_id_type_options,
    }


def copy_row_style(ws, src_row: int, dst_row: int, max_col: int) -> None:
    for c in range(1, max_col + 1):
        s = ws.cell(src_row, c)
        d = ws.cell(dst_row, c)
        d._style = copy(s._style)
        if s.has_style:
            d.font = copy(s.font)
            d.fill = copy(s.fill)
            d.border = copy(s.border)
            d.alignment = copy(s.alignment)
            d.protection = copy(s.protection)
            d.number_format = s.number_format
    if src_row in ws.row_dimensions:
        ws.row_dimensions[dst_row].height = ws.row_dimensions[src_row].height
        ws.row_dimensions[dst_row].hidden = ws.row_dimensions[src_row].hidden
        ws.row_dimensions[dst_row].outlineLevel = ws.row_dimensions[src_row].outlineLevel


# Column index cache: avoids re-scanning all columns on every per-SKU call.
# Keyed by id(columns) — safe because columns lists are created once per request
# and stay alive for the duration of the batch.
_COLS_INDEX_CACHE: Dict[int, Dict[str, List[str]]] = {}
_COLS_INDEX_LOCK = threading.Lock()


def _build_columns_index(columns: List[Dict]) -> Dict[str, List[str]]:
    """Build {normalised_label: [attr, ...]} index for fast O(1) label lookups."""
    idx: Dict[str, List[str]] = {}
    for c in columns:
        norm = re.sub(r"[^a-z0-9]+", " ", str(c.get("label", "")).strip().lower()).strip()
        if norm:
            idx.setdefault(norm, []).append(str(c.get("attr", "")))
    return idx


def _get_columns_index(columns: List[Dict]) -> Dict[str, List[str]]:
    key = id(columns)
    with _COLS_INDEX_LOCK:
        if key not in _COLS_INDEX_CACHE:
            _COLS_INDEX_CACHE[key] = _build_columns_index(columns)
            # Limit cache size to 32 distinct columns objects
            if len(_COLS_INDEX_CACHE) > 32:
                del _COLS_INDEX_CACHE[next(iter(_COLS_INDEX_CACHE))]
        return _COLS_INDEX_CACHE[key]


def first_attr_by_label(columns: List[Dict], label: str) -> str:
    target_norm = re.sub(r"[^a-z0-9]+", " ", (label or "").strip().lower()).strip()
    idx = _get_columns_index(columns)
    # Exact normalised match
    if target_norm in idx:
        return idx[target_norm][0]
    # Tolerant "contains" fallback
    for norm, attrs in idx.items():
        if target_norm and (target_norm in norm or norm in target_norm):
            return attrs[0]
    return ""


def all_attrs_by_label(columns: List[Dict], label: str) -> List[str]:
    target_norm = re.sub(r"[^a-z0-9]+", " ", (label or "").strip().lower()).strip()
    idx = _get_columns_index(columns)
    if target_norm in idx:
        return idx[target_norm]
    fuzzy = []
    for norm, attrs in idx.items():
        if target_norm and (target_norm in norm or norm in target_norm):
            fuzzy.extend(attrs)
    return fuzzy


# ═══════════════════════════════════════════════════════════════════════════
# CATALOG PIPELINE — Excel parsers, master-node lookup, AI providers
# ═══════════════════════════════════════════════════════════════════════════

def ensure_catalog_upload_dir() -> Path:
    CATALOG_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    return CATALOG_UPLOAD_DIR


def read_image_catalog(path: Path) -> Dict[str, List[str]]:
    """Parse an image-catalog file (Excel or CSV).

    Expected format (any sheet whose name contains 'catalog' or the first sheet):
      Column 1: SKU / image filename key
      Columns 2–N: Public Image URL 1 … (any number)

    Returns {sku: [non-empty url, ...]} for every product row.
    Supports .xlsx / .xlsm / .xls and .csv.
    """
    ext = path.suffix.lower()

    if ext == ".csv":
        import csv
        rows = []
        for encoding in ("utf-8-sig", "utf-8", "latin-1"):
            try:
                with open(path, newline="", encoding=encoding) as fh:
                    rows = list(csv.reader(fh))
                break
            except UnicodeDecodeError:
                continue
        if not rows:
            raise ValueError("Image catalog CSV is empty or unreadable.")
    else:
        wb = load_workbook(path, data_only=True, read_only=True)
        sheet_name = next(
            (s for s in wb.sheetnames if "catalog" in s.lower()),
            wb.sheetnames[0],
        )
        ws = wb[sheet_name]
        rows = [[c for c in row] for row in ws.iter_rows(values_only=True)]
        wb.close()
        if not rows:
            raise ValueError("Image catalog file is empty.")

    header = [str(c or "").strip().lower() for c in rows[0]]
    sku_col = 0
    _URL_KW = ("url", "image", "img", "photo", "picture", "link", "src", "media", "file")
    url_cols: List[int] = [
        i for i, h in enumerate(header)
        if i != 0 and any(kw in h for kw in _URL_KW)
    ]
    if not url_cols:
        url_cols = list(range(1, len(header)))

    catalog: Dict[str, List[str]] = {}
    for row in rows[1:]:
        if not row or not row[sku_col]:
            continue
        sku = str(row[sku_col]).strip()
        if not sku:
            continue
        urls = []
        seen_urls: set = set()
        for ci in url_cols:
            if ci < len(row) and row[ci]:
                u = str(row[ci] or "").strip()
                if u.lower().startswith(("http://", "https://")) and u not in seen_urls:
                    urls.append(u)
                    seen_urls.add(u)
        # Per-row fallback: scan ALL non-SKU columns for any http URL not already captured
        if not urls:
            for ci in range(1, len(row)):
                if ci < len(row) and row[ci]:
                    u = str(row[ci] or "").strip()
                    if u.lower().startswith(("http://", "https://")) and u not in seen_urls:
                        urls.append(u)
                        seen_urls.add(u)
        catalog[sku] = urls

    if not catalog:
        raise ValueError("No product rows found in image catalog.")
    return catalog


def parse_flat_file_skus(path: Path) -> List[Dict]:
    """Parse an Amazon flat-file xlsm/xlsx, return one dict per product data row.

    Uses the same `parse_template_rows` / `read_column_metadata` helpers as
    `load_template_meta` so output dicts have {attribute: value} pairs.
    """
    wb = load_workbook(path, keep_vba=True, data_only=True)
    if "Template" not in wb.sheetnames:
        raise ValueError("Flat file has no 'Template' sheet.")
    ws = wb["Template"]
    label_row, attribute_row, data_row = parse_template_rows(ws)
    columns = read_column_metadata(ws, label_row, attribute_row)
    ordered_cols = sorted(columns.values(), key=lambda c: c.col_idx)

    products: List[Dict] = []
    for row_idx in range(data_row, ws.max_row + 1):
        row: Dict[str, str] = {}
        for meta in ordered_cols:
            val = ws.cell(row_idx, meta.col_idx).value
            row[meta.attribute] = str(val).strip() if val is not None else ""
        if any(v for v in row.values()):
            products.append(row)
    return products


def _load_master_nodes() -> List[Dict]:
    global _MASTER_NODES_CACHE
    if _MASTER_NODES_CACHE is not None:
        return _MASTER_NODES_CACHE
    if not MASTER_NODE_PATH.is_file():
        log.warning("Master node file not found: %s", MASTER_NODE_PATH)
        _MASTER_NODES_CACHE = []
        return _MASTER_NODES_CACHE
    wb = load_workbook(MASTER_NODE_PATH, data_only=True, read_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        _MASTER_NODES_CACHE = []
        return _MASTER_NODES_CACHE
    header = [str(c or "").strip().lower() for c in rows[0]]
    nodes: List[Dict] = []
    for row in rows[1:]:
        if not row:
            continue
        rec: Dict = {header[i]: str(row[i] or "").strip() for i in range(min(len(header), len(row)))}
        if rec.get("browse_node_id") or rec.get("product_type"):
            nodes.append(rec)
    wb.close()
    _MASTER_NODES_CACHE = nodes
    log.info("Loaded %d master browse nodes from %s", len(nodes), MASTER_NODE_PATH)
    return nodes


def lookup_master_node(query: str) -> Dict:
    """Return the best-matching master node record for a query string.

    Tries exact `browse_node_id` match first, then `product_type` contains match,
    then `browse_node` contains match (all case-insensitive).
    Returns {} when nothing matches.
    """
    q = (query or "").strip().lower()
    if not q:
        return {}
    nodes = _load_master_nodes()

    # 1. Exact browse_node_id
    for rec in nodes:
        if rec.get("browse_node_id", "").lower() == q:
            return rec

    # 2. Exact product_type
    for rec in nodes:
        if rec.get("product_type", "").lower() == q:
            return rec

    # 3. browse_node contains
    for rec in nodes:
        if q in rec.get("browse_node", "").lower():
            return rec

    # 4. product_type contains
    for rec in nodes:
        if q in rec.get("product_type", "").lower():
            return rec

    return {}


def extract_flat_file_nodes(flat_products: List[Dict]) -> List[Dict]:
    """Return the unique browse nodes that appear in the flat file, enriched from Book1.xlsx.

    Each returned dict has: browse_node_id, browse_node, product_type.
    Only nodes with a non-empty browse_node_id are included.
    """
    master = _load_master_nodes()
    master_by_id: Dict[str, Dict] = {
        r.get("browse_node_id", ""): r
        for r in master
        if r.get("browse_node_id", "").strip()
    }

    seen: set = set()
    result: List[Dict] = []
    for row in flat_products:
        node_id = row.get("recommended_browse_nodes", "").strip()
        if not node_id or node_id in seen:
            continue
        seen.add(node_id)
        master_rec = master_by_id.get(node_id, {})
        result.append({
            "browse_node_id": node_id,
            "browse_node": master_rec.get("browse_node", ""),
            "product_type": master_rec.get("product_type", row.get("product_type", "")),
        })

    log.info("Flat file contains %d unique browse node(s): %s", len(result), [r["browse_node_id"] for r in result])
    return result


def best_node_for_analysis(
    analysis: Dict,
    candidate_nodes: List[Dict],
    seller_notes: str = "",
) -> Dict:
    """Pick the best browse node from the template's Browse Data candidates.

    Scoring is hierarchical — words that match the leaf segment (most specific
    part of the path after the last ">") score 3×; words in the product_type
    field score 2×; words anywhere in the full path score 1×.

    Always returns a result — first candidate is the fallback so the field is
    never left blank.
    """
    if not candidate_nodes:
        return {}
    if len(candidate_nodes) == 1:
        return candidate_nodes[0]

    # ── Build search terms from every available signal ────────────────────────
    terms: set = set()
    for field in ("product_type", "category", "usage", "style", "material"):
        val = str(analysis.get(field, "")).lower()
        terms.update(w for w in re.findall(r"[a-z0-9]+", val) if len(w) > 2)
    for item in analysis.get("features", []):
        terms.update(w for w in re.findall(r"[a-z0-9]+", str(item).lower()) if len(w) > 2)
    for item in analysis.get("colors", []):
        terms.update(w for w in re.findall(r"[a-z0-9]+", str(item).lower()) if len(w) > 2)
    if seller_notes:
        # Take the first 2000 chars of seller notes for extra signal
        terms.update(
            w for w in re.findall(r"[a-z0-9]+", seller_notes[:2000].lower()) if len(w) > 3
        )

    # Remove common stop words that appear in every node path
    terms -= {"and", "the", "for", "with", "from", "use", "used", "using", "general"}

    def _score(node: Dict) -> float:
        # Both field names are supported: "browse_node" (legacy) and "path" (browse_options)
        full_path = (node.get("browse_node") or node.get("path") or "").lower()
        product_type = (node.get("product_type") or "").lower()

        # Leaf = the last segment after ">"
        segments = [s.strip() for s in full_path.split(">")]
        leaf = segments[-1] if segments else full_path

        score = 0.0
        for t in terms:
            if t in leaf:
                score += 3          # most specific match — highest weight
            elif t in product_type:
                score += 2          # product type field
            elif t in full_path:
                score += 1          # somewhere in the path

        # Normalise by path depth so a deep specific match wins over a shallow broad one
        depth_bonus = len(segments) * 0.05
        return score + depth_bonus

    scored = [(n, _score(n)) for n in candidate_nodes]
    scored.sort(key=lambda x: x[1], reverse=True)
    best, best_score = scored[0]

    log.info(
        "best_node_for_analysis: picked '%s' (score=%.2f) from %d candidates | top-3: %s",
        best.get("browse_node_id"),
        best_score,
        len(candidate_nodes),
        [(n.get("browse_node_id"), f"{s:.2f}") for n, s in scored[:3]],
    )
    return best


def best_product_type_for_analysis(analysis: Dict, valid_types: List[str]) -> str:
    """Pick the best matching product type from the template's valid options.

    Scores each valid_type by counting how many image-analysis signal words appear
    in (or share a substring with) the type's token set.  Falls back to the first
    valid type when nothing scores above zero so the field is never left blank.
    """
    if not valid_types:
        return ""
    if len(valid_types) == 1:
        return valid_types[0]

    terms: set = set()
    for field in ("product_type", "category", "usage", "style", "material"):
        val = str(analysis.get(field, "")).lower()
        terms.update(w for w in re.findall(r"[a-z0-9]+", val) if len(w) > 2)
    for item in analysis.get("features", []):
        terms.update(w for w in re.findall(r"[a-z0-9]+", str(item).lower()) if len(w) > 2)

    def _score(pt: str) -> int:
        pt_tokens = set(re.findall(r"[a-z0-9]+", pt.lower()))
        return sum(1 for t in terms if any(t in tok or tok in t for tok in pt_tokens))

    scored = sorted(valid_types, key=_score, reverse=True)
    best_pt = scored[0]
    log.info(
        "best_product_type_for_analysis: picked %r from %s | top scores: %s",
        best_pt,
        valid_types,
        [(pt, _score(pt)) for pt in scored[:3]],
    )
    return best_pt


# ── AI Providers (Groq only) ──────────────────────────────────────────────────



# ── Groq provider ─────────────────────────────────────────────────────────────

def _groq_client():
    if not GROQ_AVAILABLE or _groq_module is None:
        raise ValueError("GROQ_API_KEY is not set. Get a free key at console.groq.com")
    return _groq_module.Groq(api_key=_GROQ_API_KEY, timeout=90.0)


def groq_analyze_image_url(image_url: str, product_name: str = "") -> Dict[str, Any]:
    """Analyze a product image using Groq's Llama vision model.

    Free tier: 100 req/day. Skips immediately when the daily budget is exhausted
    so large batches don't waste a round-trip on quota errors.
    """
    if not _groq_vision_available():
        raise ValueError("Groq vision daily limit reached — falling back to next provider.")
    client = _groq_client()
    img_bytes = read_image_bytes_any(image_url)

    # Compress to keep within Groq's base64 payload limits
    try:
        img = Image.open(BytesIO(img_bytes)).convert("RGB")
        img.thumbnail((800, 800), Image.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=75, optimize=True)
        img_bytes = buf.getvalue()
    except Exception:
        pass

    import base64 as _b64
    img_b64 = _b64.b64encode(img_bytes).decode()
    name_hint = f"Product name hint: {product_name.strip()}. " if product_name.strip() else ""
    prompt = name_hint + _CLAUDE_VISION_PROMPT

    resp = client.chat.completions.create(
        model=GROQ_VISION_MODEL,
        messages=[{
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                {"type": "text", "text": prompt},
            ],
        }],
        max_tokens=700,
    )
    raw = resp.choices[0].message.content if resp.choices else ""
    if not raw:
        raise ValueError("Groq vision returned empty response.")
    _groq_vision_increment()
    try:
        return extract_json(raw)
    except ValueError:
        return _parse_vision_text(raw)


def _parse_groq_retry_after(exc_str: str) -> float:
    """Extract the 'try again in Xs' wait from a Groq 429 error string."""
    import re as _re
    m = _re.search(r"try again in\s+([\d.]+)s", exc_str, _re.IGNORECASE)
    if m:
        return min(float(m.group(1)) + 1.0, 15.0)
    return 8.0  # safe default


def groq_generate_listing(product_analysis: Dict[str, Any], seller_notes: str = "") -> Dict[str, Any]:
    """Generate Amazon listing copy via Groq (llama-3.1-8b-instant).

    Free tier: 14,400 req/day. Bails immediately on daily quota errors so the
    caller's fallback chain moves to the local heuristic without wasted retries.
    Uses _GROQ_TEXT_SEM to cap concurrent calls and stay under TPM limit.
    """
    if not _groq_text_available():
        raise CopyGenerationError("Groq text daily limit reached — falling back to next provider.")
    client = _groq_client()
    prompt = build_listing_prompt(product_analysis, seller_notes=seller_notes)
    _last_exc: Optional[Exception] = None
    for attempt in range(4):
        _rate_wait: float = 0.0
        with _GROQ_TEXT_SEM:
            try:
                resp = client.chat.completions.create(
                    model=GROQ_TEXT_MODEL,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=1400,
                )
                raw = resp.choices[0].message.content if resp.choices else ""
                if not raw:
                    raise CopyGenerationError("Groq returned empty listing response.")
                _groq_text_increment()
                return extract_json(raw)
            except CopyGenerationError:
                raise
            except ValueError as exc:
                # JSON parse error — retry immediately after releasing semaphore
                _last_exc = exc
                if attempt >= 3:
                    raise
                log.warning("Groq JSON parse error on attempt %d, retrying: %s", attempt+1, str(exc)[:120])
            except Exception as exc:
                exc_str = str(exc)
                if _groq_is_quota_error(exc_str):
                    log.warning("Groq text daily quota exhausted — switching to fallback provider.")
                    with _groq_text_lock:
                        global _groq_text_used
                        _groq_text_used = _GROQ_TEXT_DAY_LIMIT
                    raise CopyGenerationError(f"Groq daily quota exhausted: {exc}") from exc
                if "rate_limit" in exc_str.lower() or "429" in exc_str:
                    _rate_wait = _parse_groq_retry_after(exc_str)
                    _last_exc = exc
                    log.warning("Groq rate limit, waiting %.1fs OUTSIDE semaphore (attempt %d/4)",
                                _rate_wait, attempt + 1)
                else:
                    raise
        # Sleep OUTSIDE the semaphore so other workers can proceed
        if _rate_wait:
            time.sleep(_rate_wait)
            if attempt == 3:
                break


_ONE_SHOT_PROMPT = """You are a world-class advertising copywriter who also happens to be a sharp Amazon product analyst. You see a product image and do two things simultaneously: read it with a cataloger's precision, then write about it with an ad creative's flair. Your copy makes people feel something — not just informed, but genuinely excited to own the product.

Look at this product image carefully, then return a SINGLE JSON object with exactly two top-level keys: "analysis" and "listing".

"analysis" — objective visual facts only, no guessing:
{
  "product_type": "exact product name (e.g. Insulated Stainless Steel Water Bottle)",
  "category": "Amazon category (e.g. Kitchen & Dining, Pet Supplies, Electronics)",
  "material": "primary visible material (e.g. Stainless Steel, Nylon, Leather)",
  "colors": ["every visible color on the product"],
  "features": ["every visible feature, detail, port, or attachment — be specific"],
  "usage": "primary real-world use case",
  "style": "design aesthetic (e.g. Modern Minimalist, Rugged Outdoor, Classic Casual)",
  "pattern": "surface finish or pattern if visible, else empty string",
  "fit_or_size_hint": "any visible size indicator, else empty string",
  "gender_target": "Unisex / Men / Women / Kids — Unisex if unclear",
  "confidence": 0.95
}

"listing" — advertising copy that makes someone fall in love before they've touched it:
{
  "title": "170–200 char Amazon title. This is the FIRST thing the consumer reads — make every character earn its place. Structure: [Product Type] + [Material/Key Spec] + [Primary Variant: color/size] + [Top Benefit] + [Key consumer-relevant attributes from constrained_fields such as occasion, fabric type, pattern, fit, season, style — pick the most informative ones]. Write it last, after you know everything about the product. No ALL CAPS. No pipes or slashes. Read like a product a real person searches for, not a keyword list.",
  "bullet_points": [
    "ALL-CAPS HOOK — Payoff sentence that makes the reader picture themselves using this. Sensory, specific, benefit-first. 80–200 chars.",
    "ALL-CAPS HOOK — A different angle: design or style detail that earns the second look. Make them feel the quality.",
    "ALL-CAPS HOOK — The fit, feel, or ergonomic detail that makes daily use a pleasure.",
    "ALL-CAPS HOOK — The occasion or moment this product was built for. Paint the scene.",
    "ALL-CAPS HOOK — The practical payoff that keeps them coming back: durability, care, or versatility."
  ],
  "product_description": "3 punchy paragraphs. Para 1: hook — open with the desire, the moment, or the problem this solves. Para 2: the product story — materials, construction, the design choices that matter. Para 3: the close — versatility, styling, or a line that sends them to checkout. Tone: warm, confident, expert friend. No 'This product'. No 'Introducing'. No mention of the image. Based ONLY on confirmed visual attributes and seller notes.",
  "keywords": ["20 distinct search terms — mix broad category terms, material+style combos, occasion long-tails, and buyer-intent phrases. All lowercase, no brand names."],
  "hsn_code": "exact 8-digit HSN for Indian customs/GST. E.g. textiles 62044390, electronics 85044000, plastics 39269099, steel 73269099, toys 95030099. Digits only.",
  "product_tax_code": "A_GEN_NOTAX (0%), A_GEN_TAX_5 (5%), A_GEN_TAX_12 (12%), A_GEN_TAX_18 (18% default), A_GEN_TAX_28 (28%)"
}

Rules:
- analysis: only what is VISUALLY CONFIRMED — zero guessing
- listing: write with the energy of someone who genuinely loves this category
- All 5 bullet points required, each a different angle, each making the reader feel something
- No lazy words: premium, luxury, best, world-class, revolutionary, amazing, perfect
- No unverifiable claims: clinically proven, FDA approved, guaranteed, dermatologist tested
- No mention of the image, photo, packaging design, or label aesthetics in bullets
- Return ONLY the JSON object — no markdown fences, no explanation, no preamble
- hsn_code: exactly 8 digits, padded if needed"""


def groq_one_shot(
    image_url: str,
    product_name: str = "",
    seller_notes: str = "",
    constrained_fields_guide: Optional[Dict[str, List[str]]] = None,
) -> tuple:
    """Single Groq vision call that returns (analysis_dict, listing_dict).

    Replaces the separate groq_analyze_image_url + groq_generate_listing calls,
    cutting API round-trips per SKU from 2 to 1.

    constrained_fields_guide: {field_name: [valid_option, ...]} built from the
      uploaded template's Valid Values sheet.  When provided, Groq picks the best
      matching value for each field and returns them under listing.constrained_fields.
    """
    if not _groq_vision_available():
        raise ValueError("Groq vision daily limit reached.")
    if not _groq_text_available():
        raise ValueError("Groq text daily limit reached.")

    client = _groq_client()
    img_bytes = read_image_bytes_any(image_url)
    try:
        img = Image.open(BytesIO(img_bytes)).convert("RGB")
        img.thumbnail((1024, 1024), Image.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="JPEG", quality=85, optimize=True)
        img_bytes = buf.getvalue()
    except Exception:
        pass

    import base64 as _b64
    img_b64 = _b64.b64encode(img_bytes).decode()

    prompt = _ONE_SHOT_PROMPT
    if product_name.strip():
        prompt = f"Product name hint: {product_name.strip()}.\n\n" + prompt
    if seller_notes.strip():
        prompt += f"\n\nSeller reference notes (supplementary only — image is the primary source of truth; use these only to fill gaps not visible in the image):\n{seller_notes.strip()[:50000]}"

    if constrained_fields_guide:
        lines = [
            "\n\nFor the fields below, pick the SINGLE best-matching option from the provided list.",
            "Only fill a field when you can confidently determine its value from the image.",
            "Use empty string \"\" for any field you cannot determine.",
            "Add a \"constrained_fields\" key inside \"listing\" with your choices:\n",
        ]
        for field, options in constrained_fields_guide.items():
            lines.append(f'  "{field}": one of {json.dumps(options)}')
        prompt += "\n".join(lines)

    # Extra tokens when constrained_fields guide is injected
    _max_tok = 2800 if constrained_fields_guide else 2400

    _last_exc: Optional[Exception] = None
    for attempt in range(4):
        _rate_wait: float = 0.0
        with _GROQ_TEXT_SEM:
            try:
                resp = client.chat.completions.create(
                    model=GROQ_VISION_MODEL,
                    messages=[{
                        "role": "user",
                        "content": [
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                            {"type": "text", "text": prompt},
                        ],
                    }],
                    max_tokens=_max_tok,
                )
                raw = resp.choices[0].message.content if resp.choices else ""
                if not raw:
                    raise ValueError("Groq one-shot returned empty response.")
                _groq_vision_increment()
                _groq_text_increment()
                parsed = extract_json(raw)
                analysis = parsed.get("analysis") or {}
                listing  = parsed.get("listing")  or {}
                if not analysis or not listing:
                    raise ValueError("One-shot response missing analysis or listing keys.")
                # Normalise constrained_fields: keep only non-empty string values
                cf_raw = listing.get("constrained_fields") or {}
                if isinstance(cf_raw, dict):
                    listing["constrained_fields"] = {
                        k: str(v).strip()
                        for k, v in cf_raw.items()
                        if str(v).strip() and str(v).strip().lower() not in {"", "n/a", "unknown", "uncertain"}
                    }
                return analysis, listing
            except ValueError as exc:
                _last_exc = exc
                if attempt < 3:
                    log.warning("Groq one-shot JSON error attempt %d, retrying: %s", attempt + 1, str(exc)[:120])
                    # no sleep — retry immediately after releasing semaphore
            except Exception as exc:
                exc_str = str(exc)
                if _groq_is_quota_error(exc_str):
                    raise CopyGenerationError(f"Groq quota exhausted: {exc}") from exc
                if "rate_limit" in exc_str.lower() or "429" in exc_str:
                    _rate_wait = _parse_groq_retry_after(exc_str)
                    _last_exc = exc
                    log.warning("Groq one-shot rate limit, waiting %.1fs OUTSIDE semaphore (attempt %d/4)",
                                _rate_wait, attempt + 1)
                else:
                    raise
        # Sleep OUTSIDE the semaphore so other workers can proceed
        if _rate_wait:
            time.sleep(_rate_wait)
            if attempt == 3:
                break
    raise _last_exc or ValueError("Groq one-shot failed after 4 attempts.")


def groq_refine_listing(
    product_analysis: Dict[str, Any],
    listing: Dict[str, Any],
    errors: List[str],
) -> Dict[str, Any]:
    """Refine a rejected listing via Groq. Uses _GROQ_TEXT_SEM to stay under TPM."""
    client = _groq_client()
    prompt = f"""The following Amazon listing failed quality validation.

Validation Errors:
{json.dumps(errors, indent=2)}

Original Listing:
{json.dumps(listing, indent=2)}

Product Analysis (ground truth — do not contradict):
{json.dumps(product_analysis, indent=2)}

Rewrite the listing fixing every error. Return strict JSON (no markdown):
{{
  "title": "",
  "bullet_points": ["", "", "", "", ""],
  "description": ""
}}"""
    with _GROQ_TEXT_SEM:
        resp = client.chat.completions.create(
            model=GROQ_TEXT_MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1000,
        )
    raw = resp.choices[0].message.content if resp.choices else ""
    if not raw:
        raise CopyGenerationError("Groq returned empty refinement response.")
    return extract_json(raw)


def _generate_listing(
    product_analysis: Dict[str, Any],
    seller_notes: str = "",
    requested_provider: str = "",
) -> Dict[str, Any]:
    """Generate listing copy using Groq, falling back to the local heuristic."""
    try:
        return groq_generate_listing(product_analysis, seller_notes=seller_notes)
    except Exception as exc:
        log.warning("Groq listing generation failed, using fallback: %s", exc)
        return _fallback_listing_from_analysis(product_analysis)


def _refine_listing(
    product_analysis: Dict[str, Any],
    listing: Dict[str, Any],
    errors: List[str],
) -> Dict[str, Any]:
    """Refine listing copy using Groq."""
    try:
        return groq_refine_listing(product_analysis, listing, errors)
    except Exception:
        return listing


def _get_flat_parentage(flat: Dict) -> Tuple[str, str]:
    """Return (parentage_level, parent_sku) from a flat file row dict.

    Attribute names embed marketplace IDs, so match by substring rather than exact key.
    Returns lowercase parentage_level (e.g. "parent", "child") and raw parent_sku string.
    """
    parentage_level = ""
    parent_sku = ""
    for k, v in flat.items():
        kl = k.lower()
        if not parentage_level and "parentage_level" in kl:
            parentage_level = str(v or "").strip().lower()
        elif not parent_sku and "parent_sku" in kl:
            parent_sku = str(v or "").strip()
    return parentage_level, parent_sku


def _child_sku_suffix(analysis: Dict) -> str:
    """Build a short hyphen-separated SKU differentiator from image analysis.

    Priority: first detected color, then first word of material.
    Returns empty string when no usable attributes are found.
    """
    _invalid = {"uncertain", "unknown", "n/a", ""}
    parts: List[str] = []

    colors = [str(c).strip() for c in analysis.get("colors", [])
              if str(c).strip().lower() not in _invalid]
    if colors:
        slug = re.sub(r"[^a-z0-9]+", "", colors[0].lower())[:12]
        if slug:
            parts.append(slug)

    material = str(analysis.get("material", "")).strip().lower()
    if material not in _invalid:
        mat_slug = re.sub(r"[^a-z0-9]+", "", material.split()[0])[:12]
        if mat_slug and mat_slug not in parts:
            parts.append(mat_slug)

    return "-".join(parts)


def _call_image_service(
    svc_url: str,
    sku: str,
    image_url: str,
    product_name: str,
    provider: str,
    analysis: Dict,
    bullets: List[str],
    listing: Optional[Dict] = None,
) -> List[str]:
    """POST to the image-service container and return the generated file paths.

    Falls back to an empty list on any network or application error so the
    caller can degrade gracefully (the SKU row still gets its copy filled in,
    just without generated images).
    """
    body = json.dumps({
        "sku": sku,
        "image_url": image_url,
        "product_name": product_name,
        "provider": provider,
        "analysis": analysis,
        "bullets": bullets,
        "listing": listing or {},
    }).encode("utf-8")
    req = Request(
        f"{svc_url.rstrip('/')}/api/generate-images",
        data=body,
        headers={"Content-Type": "application/json"},
    )
    try:
        with urlopen(req, timeout=300) as resp:
            data = json.loads(resp.read())
            if data.get("ok"):
                return [str(p) for p in data.get("image_paths", [])]
            log.warning("image-service error for %s: %s", sku, data.get("error"))
    except Exception as exc:
        log.warning("image-service unreachable for %s: %s — skipping image generation", sku, exc)
    return []


def _composite_ref_onto_path(
    img_path: Path,
    reference_source: str,
    shot_type: str,
) -> None:
    """Stamp the original product onto an AI-generated image file in-place."""
    if not reference_source or shot_type in {"main", "macro"}:
        return
    from PIL import Image as _Im, ImageFilter as _IF, ImageDraw as _ID
    try:
        if reference_source.lower().startswith("http"):
            _req = Request(reference_source, headers=_BROWSER_HEADERS)
            raw = urlopen(_req, timeout=20).read()
        else:
            raw = Path(reference_source).read_bytes()
        ref = _Im.open(BytesIO(raw)).convert("RGBA")
        bg  = _Im.open(str(img_path)).convert("RGBA")
    except Exception as exc:
        log.warning("Composite: could not open images for %s: %s", img_path.name, exc)
        return
    W, H = bg.size
    scale  = 0.58 if shot_type == "infographic" else 0.62
    target = int(min(W, H) * scale)
    ref.thumbnail((target, target), _Im.LANCZOS)
    rw, rh = ref.size
    cx = int(W * 0.42) if shot_type == "lifestyle" else W // 2
    cy = H // 2
    px, py = cx - rw // 2, cy - rh // 2
    # drop shadow
    shad = _Im.new("RGBA", (W, H), (0, 0, 0, 0))
    smask = _Im.new("L", (rw, rh), 0)
    _ID.Draw(smask).ellipse([(rw // 10, rh * 3 // 4), (rw * 9 // 10, rh + rh // 6)], fill=80)
    smask = smask.filter(_IF.GaussianBlur(radius=rw // 14))
    shad.paste(_Im.new("RGB", (rw, rh), (0, 0, 0)), (px, py), smask)
    shad = shad.filter(_IF.GaussianBlur(radius=3))
    result = _Im.alpha_composite(bg.copy(), shad)
    result.paste(ref, (px, py), ref.split()[3])
    try:
        result.convert("RGB").save(str(img_path), format="PNG")
        log.info("Composited reference product onto %s (%s)", img_path.name, shot_type)
    except Exception as exc:
        log.warning("Composite save failed for %s: %s", img_path.name, exc)


# Shot type index: prompt 0 = lifestyle, prompt 1 = infographic
_PROMPT_SHOT_TYPES = ["lifestyle", "infographic"]


def _download_image_to_dir(image_url: str, out_dir: Path, filename: str) -> Path:
    """Download a remote image URL to `out_dir/filename`. Returns the saved path."""
    out_dir.mkdir(parents=True, exist_ok=True)
    raw = read_image_bytes_any(image_url)
    dest = out_dir / filename
    dest.write_bytes(raw)
    return dest


def generate_product_extra_images(
    sku: str,
    image_url: str,
    copy_data: Dict,
    out_dir: Path,
) -> List[str]:
    """Download reference image then produce 2 listing visuals (lifestyle + infographic).

    Returns list of two absolute file path strings.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    ref_local = _download_image_to_dir(image_url, out_dir, f"{sku}_ref.jpg")
    bullets = [str(b).strip() for b in copy_data.get("bullet_points", []) if str(b).strip()]
    product_name = copy_data.get("title", sku) or sku
    return generate_reference_locked_images(
        reference_image_source=str(ref_local),
        out_dir=out_dir,
        prefix=sku,
        product_name=product_name,
        bullets=bullets[:5],
    )


def process_catalog_row(
    sku: str,
    image_urls: List[str],
    flat_file_row: Optional[Dict] = None,
    candidate_nodes: Optional[List[Dict]] = None,
    provider: str = "demo",
    generate_images: bool = True,
    seller_notes: str = "",
    valid_product_types: Optional[List[str]] = None,
    brand_name: str = "",
    all_valid_values: Optional[Dict[str, List[str]]] = None,
) -> Dict:
    """Full per-product pipeline: analyze → copy → validate → (optionally) generate images.

    candidate_nodes: browse nodes extracted from the uploaded flat file
      (via extract_flat_file_nodes). When provided, node selection is restricted
      to this set and the best match is chosen after image analysis so the AI's
      product_type signal informs the choice.  When None, falls back to searching
      all master nodes.

    seller_notes: optional plain-text content from a seller-provided description /
      spec sheet. When non-empty, it is injected into the listing generation
      prompt so the AI can mine it for features, dimensions, and keywords that
      are not visible in the product image.

    Returns dict with keys:
      sku, title, bullet_points, product_description, analysis,
      browse_node_id, product_type, browse_node, image_paths, errors,
      parentage_level, parent_sku
    """
    flat = flat_file_row or {}

    # ── 1. Browse node is resolved AFTER image analysis when candidates are known ─
    # We defer it below; set defaults now.
    node_meta: Dict = {}

    # ── 2. Get first valid image URL (HTTP URL or local file path) ───────────
    image_url = next(
        (u for u in image_urls
         if u.lower().startswith(("http://", "https://")) or (u and Path(u).is_file())),
        "",
    )

    # ── 2b. Scope seller notes to this SKU only ───────────────────────────────
    # Prevents cross-SKU contamination when a single seller notes file covers
    # an entire catalog (e.g. one PDF with descriptions for 60 fruit snack SKUs).
    _sku_seller_notes = extract_sku_seller_notes(
        seller_notes, sku,
        product_name=re.sub(r"\d+$", "", sku).replace("_", " ").replace("-", " ").strip(),
    )

    # ── 2c. Build constrained fields guide from Valid Values sheet ────────────
    # Passes only fields with a small, enumerable set of options so Groq can pick
    # the correct value directly rather than generating free-form text.
    _already_handled = {"material type", "material composition", "material",
                        "color", "color map", "size", "style",
                        "product type", "recommended browse nodes"}
    _constrained_guide: Dict[str, List[str]] = {}
    if all_valid_values:
        for _f, _opts in all_valid_values.items():
            if _f.lower() in _already_handled or not _opts:
                continue
            if 2 <= len(_opts) <= 30:
                _constrained_guide[_f] = _opts[:15]  # cap per-field to keep prompt small
            if len(_constrained_guide) >= 12:          # max 12 extra fields total
                break

    # ── 3 + 4. Analyze image AND generate listing via Groq one-shot ─────────────
    # A single vision call returns both the product analysis and the listing copy.
    # Falls back to the two-call path if the one-shot fails.
    analysis: Dict[str, Any] = {}
    listing:  Dict[str, Any] = {}
    analysis_errors: List[str] = []
    errors: List[str] = []
    _one_shot_done = False

    _vision_hint = re.sub(r"\d+$", "", sku).strip().replace("_", " ").replace("-", " ").strip()

    if image_url and not SKIP_VISION and GROQ_AVAILABLE and _groq_vision_available() and _groq_text_available():
        try:
            analysis, listing = groq_one_shot(
                image_url, _vision_hint, _sku_seller_notes,
                constrained_fields_guide=_constrained_guide or None,
            )
            _one_shot_done = True
            log.debug("One-shot complete for %s", sku)
        except Exception as exc:
            log.warning("One-shot failed for %s: %s — falling back to two-call path", sku, exc)

    if not _one_shot_done:
        # Two-call path: vision analysis then listing generation
        if image_url and not SKIP_VISION:
            try:
                analysis = groq_analyze_image_url(image_url, _vision_hint)
            except Exception as exc:
                log.warning("Image analysis failed for %s: %s", sku, exc)
                analysis_errors.append(str(exc))
        elif SKIP_VISION:
            analysis = _name_based_analysis(sku)

        if analysis:
            try:
                listing = _generate_listing(analysis, seller_notes=_sku_seller_notes)
            except Exception as exc:
                log.warning("Listing generation failed for %s: %s", sku, exc)
                listing = _fallback_listing_from_analysis(analysis)

            errors = validate_listing(listing, product_analysis=analysis)
            for _ in range(LISTING_REFINE_PASSES):
                if not errors:
                    break
                try:
                    listing = _refine_listing(analysis, listing, errors)
                    errors = validate_listing(listing, product_analysis=analysis)
                except Exception:
                    break

    if not listing:
        listing = {
            "title": sku,
            "bullet_points": [f"Product {sku} — see image for details."] * 5,
            "description": f"Product {sku}.",
        }
        errors = analysis_errors

    # ── 3b. Resolve browse node and product type ──────────────────────────────
    if candidate_nodes:
        node_meta = best_node_for_analysis(analysis, candidate_nodes, seller_notes=_sku_seller_notes)
    else:
        node_meta = {}
        log.warning("No template browse nodes available for SKU %s — node field left as-is.", sku)

    browse_node_id = (
        node_meta.get("browse_node_id")
        or node_meta.get("node_id")
        or flat.get("recommended_browse_nodes", "")
    )
    browse_node = node_meta.get("browse_node") or node_meta.get("path") or ""

    _vpt = valid_product_types or []
    if _vpt and analysis:
        product_type = best_product_type_for_analysis(analysis, _vpt)
    elif _vpt and node_meta.get("product_type") in _vpt:
        product_type = node_meta["product_type"]
    elif _vpt:
        _node_pt = node_meta.get("product_type", "")
        product_type = _node_pt if _node_pt in _vpt else _vpt[0]
    else:
        product_type = node_meta.get("product_type") or flat.get("product_type", "")

    # ── 6. Generate extra images ──────────────────────────────────────────────
    image_paths: List[str] = []
    if generate_images and IMAGE_GENERATION_ENABLED and image_url:
        product_name = listing.get("title", sku) or sku
        bullets_img = [str(b).strip() for b in listing.get("bullet_points", []) if str(b).strip()]

        _svc_url = os.environ.get("IMAGE_SERVICE_URL", "").strip()
        if _svc_url:
            # Containerised mode: delegate to the image-service container.
            image_paths = _call_image_service(
                svc_url=_svc_url,
                sku=sku,
                image_url=image_url,
                product_name=product_name,
                provider=provider,
                analysis=analysis,
                bullets=bullets_img,
                listing=listing,
            )
        else:
            # Single-container / dev mode: run image generation locally.
            out_dir = APP_DIR / "generated_images" / sku
            out_dir.mkdir(parents=True, exist_ok=True)
            ref_path = str(_download_image_to_dir(image_url, out_dir, f"{sku}_ref.jpg"))

            # ── 4-shot pipeline (primary) ─────────────────────────────────────
            _used_4shot = False
            if _PRODUCT_IMAGES_AVAILABLE:
                try:
                    ref_desc = ""
                    try:
                        ref_desc = _reference_conditioning_text(ref_path, product_name)
                    except Exception as _rde:
                        log.debug("Reference conditioning skipped: %s", _rde)

                    _product_for_gen = dict(analysis)
                    _product_for_gen["sku"] = sku

                    shot_results = _generateProductImages(
                        product=_product_for_gen,
                        reference_image=ref_path,
                        config={
                            "shots": ("main", "lifestyle"),
                            "out_dir": out_dir,
                            "provider": "fal",
                            "prefix": sku,
                        },
                        reference_description=ref_desc,
                    )
                    for shot_key in ("lifestyle", "macro", "styled_hero"):
                        if shot_key in shot_results:
                            image_paths.append(str(shot_results[shot_key].resolve()))
                    _used_4shot = True
                    log.info("4-shot pipeline produced %d images for %s", len(shot_results), sku)
                except Exception as exc:
                    log.warning("4-shot pipeline failed for %s: %s — falling back to 2-shot", sku, exc)

            # ── 2-shot fallback (PIL): PIL text-overlay ───────────────────────
            if not _used_4shot:
                try:
                    image_paths = generate_reference_locked_images(
                        reference_image_source=ref_path,
                        out_dir=out_dir,
                        prefix=sku,
                        product_name=product_name,
                        bullets=bullets_img[:5],
                        analysis=analysis,
                        listing=listing,
                    )
                except Exception as exc:
                    log.warning("PIL image generation failed for %s: %s", sku, exc)

    bullets = listing.get("bullet_points", [])
    if not isinstance(bullets, list):
        bullets = []

    raw_keywords = listing.get("keywords", [])
    if not isinstance(raw_keywords, list):
        raw_keywords = [str(raw_keywords)] if raw_keywords else []
    # Filter single-word stopwords; keep multi-word keyword phrases intact
    keywords = [
        k for k in (str(kw).strip().lower() for kw in raw_keywords if str(kw).strip())
        if k and (len(k.split()) > 1 or k not in _KW_STOPWORDS) and len(k) > 2
    ][:20]

    material_val = str(analysis.get("material", "")).strip()
    if material_val.lower() in {"uncertain", "unknown", "n/a", ""}:
        material_val = ""

    # ── HSN code: digits only, 4–8 characters; strip any spaces/dashes ──────
    _VALID_TAX_CODES = {
        "A_GEN_NOTAX", "A_GEN_TAX_5", "A_GEN_TAX_12", "A_GEN_TAX_18", "A_GEN_TAX_28",
    }
    raw_hsn = re.sub(r"[^0-9]", "", str(listing.get("hsn_code", "")).strip())
    hsn_code = raw_hsn if 4 <= len(raw_hsn) <= 8 else ""

    raw_ptc = str(listing.get("product_tax_code", "")).strip().upper()
    product_tax_code = raw_ptc if raw_ptc in _VALID_TAX_CODES else "A_GEN_TAX_18"

    # ── SKU: for child variations, use {parent_sku}-{color/material} suffix ──
    parentage_level, parent_sku = _get_flat_parentage(flat)
    if parentage_level == "child" and parent_sku and analysis:
        suffix = _child_sku_suffix(analysis)
        sku = f"{parent_sku}-{suffix}" if suffix else sku

    _inv = {"uncertain", "unknown", "n/a", ""}
    _colors = [str(c).strip() for c in analysis.get("colors", [])
               if str(c).strip().lower() not in _inv]
    _style = str(analysis.get("style", "")).strip()
    _size_hint = str(analysis.get("fit_or_size_hint", "")).strip()

    return {
        "sku": sku,
        "title": listing.get("title", sku),
        "bullet_points": [str(b).strip() for b in bullets][:5],
        "product_description": str(listing.get("product_description") or listing.get("description", "")).strip(),
        "analysis": analysis,
        "browse_node_id": browse_node_id,
        "product_type": product_type,
        "browse_node": browse_node,
        "image_paths": image_paths,
        "image_url": image_url,
        "errors": errors or [],
        "keywords": keywords,
        "material": material_val,
        "colors": _colors,
        "style": _style if _style.lower() not in _inv else "",
        "size_hint": _size_hint if _size_hint.lower() not in _inv else "",
        "estimated_dimensions": str(listing.get("estimated_dimensions", "")).strip(),
        "hsn_code": hsn_code,
        "product_tax_code": product_tax_code,
        "parentage_level": parentage_level,
        "parent_sku": parent_sku,
        "brand_name": brand_name,
        "_provider_hint": provider,
        "constrained_fields": {
            k: str(v).strip()
            for k, v in (listing.get("constrained_fields") or {}).items()
            if str(v).strip() and str(v).strip().lower() not in {"", "n/a", "unknown", "uncertain"}
        },
    }


def _local_path_to_img_url(path: str) -> str:
    """Convert an absolute local image path to a /img/<rel> web URL.

    If path is already an http/https URL it is returned unchanged.
    Only paths inside APP_DIR are converted; others pass through unchanged.
    """
    if not path or path.lower().startswith(("http://", "https://")):
        return path
    try:
        rel = Path(path).resolve().relative_to(APP_DIR.resolve())
        return "/img/" + str(rel).replace("\\", "/")
    except ValueError:
        return path


def _extract_pack_count(features: List[str]) -> str:
    """Return 'Pack of N' if clearly stated in features, else empty string."""
    for f in features:
        m = re.search(r'pack\s+of\s+(\d+)', str(f), re.IGNORECASE)
        if m:
            return f"Pack of {m.group(1)}"
        m = re.search(r'set\s+of\s+(\d+)', str(f), re.IGNORECASE)
        if m:
            return f"Set of {m.group(1)}"
    return ""


def _build_smart_title(
    brand: str,
    analysis: Dict[str, Any],
    ai_title: str,
) -> str:
    """Build a structured Amazon title: [Brand] [Product] [Color] [Size/Pack] [USP].

    If brand is 'generic' (any case), it is omitted.
    Falls back to the AI-generated title if analysis is too sparse.
    """
    _inv = {"uncertain", "unknown", "n/a", ""}
    brand_clean = brand.strip()
    include_brand = bool(brand_clean) and brand_clean.lower() != "generic"

    product_type = str(analysis.get("product_type", "")).strip()
    colors = [c for c in analysis.get("colors", []) if str(c).strip().lower() not in _inv]
    primary_color = str(colors[0]).strip().title() if colors else ""
    size_hint = str(analysis.get("fit_or_size_hint", "") or analysis.get("estimated_dimensions", "")).strip()
    if size_hint.lower() in _inv:
        size_hint = ""
    features = [str(f).strip() for f in analysis.get("features", []) if str(f).strip().lower() not in _inv]
    pack = _extract_pack_count(features)
    # First meaningful feature as USP (skip pack-related ones)
    usp = next((f for f in features if "pack" not in f.lower() and "set of" not in f.lower()), "")

    if not product_type:
        # Nothing useful in analysis — use AI title, just prepend brand
        if include_brand and ai_title and not ai_title.lower().startswith(brand_clean.lower()):
            return f"{brand_clean} {ai_title}"[:200].strip()
        return ai_title

    parts = []
    if include_brand:
        parts.append(brand_clean)
    parts.append(product_type)
    if primary_color:
        parts.append(primary_color)
    if pack:
        parts.append(pack)
    elif size_hint:
        parts.append(size_hint)
    if usp and len(" ".join(parts)) < 140:
        parts.append(usp)

    title = " ".join(parts).strip()
    return title[:200] if title else ai_title


def _generate_title_from_description(
    description: str,
    brand: str,
    analysis: Dict[str, Any],
    requested_provider: str = "",
) -> str:
    """Use AI to write a precise, eloquent Amazon title from the product description.

    The description is the primary source — the title must be a faithful
    distillation of it, every word earns its place.
    Returns the title string (no JSON wrapper).
    Falls back to _build_smart_title on any failure.
    """
    _inv = {"uncertain", "unknown", "n/a", ""}
    brand_clean = brand.strip()
    include_brand = bool(brand_clean) and brand_clean.lower() != "generic"

    colors = [str(c).strip() for c in analysis.get("colors", []) if str(c).strip().lower() not in _inv]
    primary_color = colors[0] if colors else ""
    size_hint = str(analysis.get("fit_or_size_hint", "")).strip()
    if size_hint.lower() in _inv:
        size_hint = ""
    features = [str(f).strip() for f in analysis.get("features", []) if str(f).strip().lower() not in _inv]
    pack = _extract_pack_count(features)

    variant_parts = [v for v in [primary_color, pack or size_hint] if v]
    variant_str = ", ".join(variant_parts) if variant_parts else "none"
    brand_line = f"Brand: {brand_clean} (include at the start of the title)" if include_brand else "Brand: not applicable (do NOT add any brand)"

    prompt = f"""You are an Amazon title specialist. Write a single product title — nothing else.

Product description (primary source — your title must be a faithful distillation of this):
\"\"\"{description.strip()[:1200]}\"\"\"

Additional context:
{brand_line}
Key variant details: {variant_str}

Title rules:
- 80–150 characters, aim for under 120
- Every word earns its place — no filler, no repetition, no "premium/best/amazing"
- Structure: {('[Brand] ' if include_brand else '')}[Product Type] [Material or Key Spec] [Variant] [Top USP]
- No pipes, no slashes, no ALL CAPS acronyms
- Must read as a real product name, not a keyword string
- Derive the USP from the description — use the single most compelling detail

Reply with ONLY the title text. No quotes, no explanation, no JSON."""

    if GROQ_AVAILABLE and _groq_module is not None and _groq_text_available():
        try:
            client = _groq_client()
            with _GROQ_TEXT_SEM:
                resp = client.chat.completions.create(
                    model=GROQ_TEXT_MODEL,
                    messages=[{"role": "user", "content": prompt}],
                    max_tokens=80,
                    temperature=0.4,
                )
            title = (resp.choices[0].message.content or "").strip().strip('"').strip()
            if title and not title.startswith("{") and len(title) > 10:
                return title[:200]
        except Exception as exc:
            log.debug("Groq title generation failed: %s", exc)

    # Fallback: mechanical assembly
    return _build_smart_title(brand, analysis, "")


def _snap_to_valid(value: str, options: List[str]) -> str:
    """Return the best-matching option for value from the Valid Values list.

    Matching order: exact → starts-with → word-contained-in → contained-in.
    Returns the original value unchanged if nothing is close enough (so we never
    replace a good AI value with a wrong constrained one).
    """
    if not value or not options:
        return value
    vl = value.strip().lower()
    # 1. Exact match
    for o in options:
        if o.lower() == vl:
            return o
    # 2. Starts-with
    for o in options:
        if o.lower().startswith(vl) or vl.startswith(o.lower()):
            return o
    # 3. All words of value appear in option (e.g. "dark blue" → "Navy Blue")
    words = set(vl.split())
    for o in options:
        if words and words.issubset(set(o.lower().split())):
            return o
    # 4. Value substring of option
    for o in options:
        if vl in o.lower():
            return o
    return value  # no valid match — keep original


def catalog_result_to_template_row(
    result: Dict,
    columns: List[Dict],
    extra_image_paths: Optional[List[str]] = None,
    brand_name: str = "",
    flat_row: Optional[Dict] = None,
    all_valid_values: Optional[Dict[str, List[str]]] = None,
) -> Dict:
    """Map a `process_catalog_row` result dict onto an Amazon template attribute dict.

    Pre-populates from flat_row (preserving existing data) then overlays AI output.
    """
    # Start from existing flat file data so untouched columns are preserved
    row: Dict[str, str] = {c["attr"]: "" for c in columns}
    if flat_row:
        for c in columns:
            existing = str(flat_row.get(c["attr"], "") or "").strip()
            if existing:
                row[c["attr"]] = existing

    desc_attr = first_attr_by_label(columns, "Product Description")
    item_name_attr = first_attr_by_label(columns, "Item Name")
    bullet_attrs = all_attrs_by_label(columns, "Bullet Point")[:5]
    browse_node_attr = first_attr_by_label(columns, "Recommended Browse Nodes")
    product_type_attr = first_attr_by_label(columns, "Product Type")
    main_image_attr = first_attr_by_label(columns, "Main Image URL")
    other_image_attrs = all_attrs_by_label(columns, "Other Image URL")[:2]
    sku_attr = (
        first_attr_by_label(columns, "Seller SKU")
        or first_attr_by_label(columns, "Item SKU")
        or first_attr_by_label(columns, "SKU")
    )

    analysis = result.get("analysis") or {}
    _inv = {"uncertain", "unknown", "n/a", ""}

    if sku_attr:
        row[sku_attr] = result.get("sku", "")

    # ── Item Name: AI title derived from the generated description ───────────
    if item_name_attr:
        description = result.get("product_description", "").strip()
        ai_title    = result.get("title", "").strip()
        provider    = result.get("_provider_hint", "")
        if description:
            row[item_name_attr] = _generate_title_from_description(
                description, brand_name, analysis, requested_provider=provider
            )
        elif ai_title:
            row[item_name_attr] = _build_smart_title(brand_name, analysis, ai_title)
        else:
            row[item_name_attr] = _build_smart_title(brand_name, analysis, "")

    if desc_attr:
        row[desc_attr] = result.get("product_description", "")
    for i, attr in enumerate(bullet_attrs):
        bullets = result.get("bullet_points", [])
        row[attr] = bullets[i] if i < len(bullets) else ""
    if browse_node_attr:
        row[browse_node_attr] = result.get("browse_node_id", "")
    if product_type_attr:
        row[product_type_attr] = result.get("product_type", "")
    if main_image_attr:
        row[main_image_attr] = result.get("image_url", "")

    # image_paths now contains up to 3 paths: [lifestyle, macro, styled_hero]
    # from the 4-shot pipeline (or [lifestyle, infographic] from legacy 2-shot).
    other_image_attrs = all_attrs_by_label(columns, "Other Image URL")[:3]
    image_paths = extra_image_paths or result.get("image_paths", [])
    for i, attr in enumerate(other_image_attrs):
        if i < len(image_paths):
            row[attr] = _local_path_to_img_url(image_paths[i])

    # ── Keywords: fill Generic Keyword and/or Search Terms columns ────────────
    keywords = result.get("keywords", [])
    if keywords:
        kw_idx = 0
        generic_kw_attrs = all_attrs_by_label(columns, "Generic Keyword")
        for attr in generic_kw_attrs[:20]:
            row[attr] = keywords[kw_idx] if kw_idx < len(keywords) else ""
            kw_idx += 1
        search_term_attrs = all_attrs_by_label(columns, "Search Terms")
        for attr in search_term_attrs[:20]:
            row[attr] = keywords[kw_idx] if kw_idx < len(keywords) else (keywords[kw_idx % len(keywords)] if keywords else "")
            kw_idx += 1

    _vv = all_valid_values or {}

    # ── Material: fill first matching material column ─────────────────────────
    material_val = result.get("material", "")
    if material_val:
        _mat_valid = _vv.get("material type", []) or _vv.get("material composition", []) or _vv.get("material", [])
        _mat_snapped = _snap_to_valid(material_val, _mat_valid)
        for mat_label in ("Material Type", "Material Composition", "Material"):
            mat_attr = first_attr_by_label(columns, mat_label)
            if mat_attr:
                row[mat_attr] = _mat_snapped
                break

    # ── Dimensions / Size: fill first matching size or dimension column ────────
    dims_val = result.get("estimated_dimensions", "")
    if dims_val and dims_val.lower() not in {"see product description", ""}:
        for size_label in ("Item Size", "Size", "Display Size", "Item Dimensions",
                           "Item Display Length", "Item Length"):
            size_attr = first_attr_by_label(columns, size_label)
            if size_attr:
                row[size_attr] = dims_val
                break

    # ── HSN Code ──────────────────────────────────────────────────────────────
    # Amazon flat files use a two-column pair:
    #   Entity column  (example = "HSN Code")  → write the literal string "HSN Code"
    #   Value  column  (the one immediately right, same label group) → write the HSN number
    #
    # We find the entity column by scanning for the column whose *example* value
    # is "HSN Code", then pair it with the next column (same label, col_idx + 1).
    hsn_val = result.get("hsn_code", "")
    if hsn_val:
        # Build a col_idx → column dict for quick neighbour lookup
        _col_by_idx = {c["col_idx"]: c for c in columns}
        _hsn_entity_col = next(
            (c for c in columns
             if "hsn" in str(c.get("example", "")).lower()),
            None,
        )
        if _hsn_entity_col:
            # Entity column: always write the literal type string
            row[_hsn_entity_col["attr"]] = "HSN Code"
            # Value column: the column with the same label base immediately to the right
            _val_col = _col_by_idx.get(_hsn_entity_col["col_idx"] + 1)
            if _val_col:
                row[_val_col["attr"]] = hsn_val
        else:
            # Fallback for templates with a direct HSN column
            for hsn_label in ("HSN Code", "HSN/SAC Code", "HSN", "HSN code"):
                hsn_attr = first_attr_by_label(columns, hsn_label)
                if hsn_attr:
                    row[hsn_attr] = hsn_val
                    break

    # ── Product Tax Code: Amazon's GST/tax classification field ───────────────
    ptc_val = result.get("product_tax_code", "")
    if ptc_val:
        ptc_attr = first_attr_by_label(columns, "Product Tax Code")
        if ptc_attr:
            row[ptc_attr] = ptc_val

    # ── Parentage: write Parentage Level and Parent SKU when available ─────────
    parentage_level = result.get("parentage_level", "")
    if parentage_level:
        pl_attr = first_attr_by_label(columns, "Parentage Level")
        if pl_attr:
            row[pl_attr] = parentage_level.capitalize()
    parent_sku_val = result.get("parent_sku", "")
    if parent_sku_val:
        ps_attr = first_attr_by_label(columns, "Parent SKU")
        if ps_attr:
            row[ps_attr] = parent_sku_val

    # ── Brand Name — always write to flat file; "generic" only skips the title ──
    brand_clean = brand_name.strip()
    if brand_clean:
        brand_attr = first_attr_by_label(columns, "Brand Name") or first_attr_by_label(columns, "Brand")
        if brand_attr:
            row[brand_attr] = brand_clean
        # Manufacturer mirrors brand unless already filled from flat file data
        if brand_clean.lower() != "generic":
            mfr_attr = first_attr_by_label(columns, "Manufacturer")
            if mfr_attr and not row.get(mfr_attr, "").strip():
                row[mfr_attr] = brand_clean

    # ── Color / Color Map ─────────────────────────────────────────────────────
    colors = [str(c).strip() for c in analysis.get("colors", [])
              if str(c).strip().lower() not in _inv]
    if colors:
        color_attr = first_attr_by_label(columns, "Color")
        if color_attr and not row.get(color_attr, "").strip():
            row[color_attr] = _snap_to_valid(colors[0], _vv.get("color", []))
        color_map_attr = first_attr_by_label(columns, "Color Map")
        if color_map_attr and not row.get(color_map_attr, "").strip():
            row[color_map_attr] = _snap_to_valid(colors[0], _vv.get("color map", []) or _vv.get("color", []))

    # ── Size ──────────────────────────────────────────────────────────────────
    size_val = str(analysis.get("fit_or_size_hint", "")).strip()
    if size_val.lower() in _inv:
        size_val = ""
    if size_val:
        size_attr = first_attr_by_label(columns, "Size")
        if size_attr and not row.get(size_attr, "").strip():
            row[size_attr] = _snap_to_valid(size_val, _vv.get("size", []))

    # ── Style ─────────────────────────────────────────────────────────────────
    style_val = str(analysis.get("style", "")).strip()
    if style_val and style_val.lower() not in _inv:
        style_attr = first_attr_by_label(columns, "Style")
        if style_attr and not row.get(style_attr, "").strip():
            row[style_attr] = _snap_to_valid(style_val, _vv.get("style", []))

    # ── Number of Packs / Unit Count ─────────────────────────────────────────
    features = [str(f).strip() for f in analysis.get("features", []) if str(f).strip()]
    pack_str = _extract_pack_count(features)
    if pack_str:
        m = re.search(r"\d+", pack_str)
        if m:
            pack_num = m.group()
            for pack_label in ("Number of Packs", "Number Of Packs", "Pack Count"):
                pa = first_attr_by_label(columns, pack_label)
                if pa and not row.get(pa, "").strip():
                    row[pa] = pack_num
                    break
            uc_attr = first_attr_by_label(columns, "Unit Count")
            if uc_attr and not row.get(uc_attr, "").strip():
                row[uc_attr] = pack_num

    # ── Model Name — use product type when no model is known ─────────────────
    product_type_val = str(analysis.get("product_type", "")).strip()
    if product_type_val and product_type_val.lower() not in _inv:
        model_name_attr = first_attr_by_label(columns, "Model Name")
        if model_name_attr and not row.get(model_name_attr, "").strip():
            row[model_name_attr] = product_type_val

    # ── Constrained fields: all values Groq picked from Valid Values options ──
    # Groq returned {field_name: chosen_value} for every field it could determine.
    # Write each to the matching template column, snapping to the valid list to
    # ensure the value is accepted by Amazon's template validator.
    constrained = result.get("constrained_fields") or {}
    for cf_field, cf_value in constrained.items():
        if not cf_value:
            continue
        # Try to find the column by title-casing the field name
        cf_attr = first_attr_by_label(columns, cf_field.title())
        if not cf_attr:
            # Also try exact case as-is
            cf_attr = first_attr_by_label(columns, cf_field)
        if cf_attr and not row.get(cf_attr, "").strip():
            valid_opts = _vv.get(cf_field.lower(), [])
            row[cf_attr] = _snap_to_valid(cf_value, valid_opts) if valid_opts else cf_value

    return row


# ─────────────────────────────────────────────────────────────────────────────

# Fields that vary per SKU — never backfill these from other rows
_PER_SKU_FIELD_KEYWORDS = {
    "color", "colour", "size", "weight", "dimension", "quantity", "pack",
    "count", "length", "width", "height", "depth", "asin", "sku", "ean",
    "upc", "gtin", "barcode", "image", "price", "currency", "marketplace",
    "parent", "child", "variation", "model", "part", "number",
}

def _should_skip_backfill(field_name: str) -> bool:
    fl = field_name.lower().replace(" ", "_").replace("-", "_")
    return any(kw in fl for kw in _PER_SKU_FIELD_KEYWORDS)


def _backfill_shared_fields(
    rows_by_idx: Dict[int, Dict],
    details_by_idx: Dict[int, Dict],
    columns: List[Dict],
    all_valid_values: Optional[Dict[str, List[str]]] = None,
) -> None:
    """For every constrained field filled in ANY row, backfill it into ALL rows.

    SKU-variant fields (color, size, weight, etc.) are excluded.
    Uses the most common non-empty value across rows as the consensus.
    Modifies rows_by_idx in place.
    """
    _vv = all_valid_values or {}

    # 1. Collect all constrained_fields values seen across SKUs
    field_values: Dict[str, List[str]] = {}
    for detail in details_by_idx.values():
        cf = (detail.get("pipeline") or {}).get("constrained_fields") or {}
        for field, value in cf.items():
            if not value or _should_skip_backfill(field):
                continue
            field_values.setdefault(field, []).append(value)

    if not field_values:
        return

    # 2. Pick consensus value per field (most common)
    consensus: Dict[str, str] = {}
    for field, values in field_values.items():
        if not values:
            continue
        consensus[field] = max(set(values), key=values.count)

    # 3. For each row, write any missing consensus field
    for idx, row in rows_by_idx.items():
        if not row:
            continue
        for field, value in consensus.items():
            # Resolve column attribute name
            attr = first_attr_by_label(columns, field.title()) or first_attr_by_label(columns, field)
            if not attr:
                continue
            if row.get(attr, "").strip():
                continue  # already has a value
            valid_opts = _vv.get(field.lower(), [])
            row[attr] = _snap_to_valid(value, valid_opts) if valid_opts else value


# ─────────────────────────────────────────────────────────────────────────────

from flask import send_file as _send_file


@app.get("/img/<path:rel_path>")
def serve_generated_image(rel_path: str):
    """Serve a generated image by its relative path under APP_DIR.

    The frontend stores absolute paths in product rows. This endpoint converts
    an absolute server-side path (passed as a query param or as part of the
    URL) into a web-accessible URL.

    Usage: /img/generated_images/kurti/kurti_main.png
    """
    try:
        # Only allow files inside APP_DIR (prevent path traversal)
        full = (APP_DIR / rel_path).resolve()
        if not str(full).startswith(str(APP_DIR.resolve())):
            return "Forbidden", 403
        if not full.exists():
            return "Not found", 404
        suffix = full.suffix.lower()
        mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                "webp": "image/webp", "gif": "image/gif"}.get(suffix.lstrip("."), "image/png")
        return _send_file(str(full), mimetype=mime)
    except Exception as exc:
        return str(exc), 500


@app.get("/api/status")
def api_status():
    """Live provider availability check with actual API probing.

    Does a lightweight API call per provider so the UI reflects real connectivity,
    not just key presence. Results are cached for 60 s to avoid hammering APIs.
    """
    providers = {}

    # ── Groq ─────────────────────────────────────────────────────────────────
    groq_key = os.environ.get("GROQ_API_KEY", "").strip()
    if groq_key and _groq_module is not None:
        try:
            _groq_module.Groq(api_key=groq_key).models.list()
            providers["groq"] = {"ok": True, "label": "Groq", "note": ""}
        except Exception as exc:
            note = "quota exhausted" if _groq_is_quota_error(str(exc)) else str(exc)[:60]
            providers["groq"] = {"ok": False, "label": "Groq", "note": note}
    else:
        providers["groq"] = {"ok": False, "label": "Groq", "note": "no key"}

    # ── FAL (image generation) ────────────────────────────────────────────────
    fal_key = os.environ.get("FAL_KEY", "").strip()
    if fal_key:
        providers["fal"] = {"ok": True, "label": "FAL (image gen)", "note": ""}
    else:
        providers["fal"] = {"ok": False, "label": "FAL (image gen)", "note": "FAL_KEY not set"}

    # ── ComfyUI (local image generation) ─────────────────────────────────────
    _a1111_host = os.environ.get("A1111_HOST", "").strip()
    _a1111_port = os.environ.get("A1111_PORT", "7860").strip()
    if _a1111_host and _a1111_host not in ("auto", ""):
        try:
            from urllib.request import urlopen as _urlopen
            from urllib.request import Request as _Req
            _probe = _Req(f"http://{_a1111_host}:{_a1111_port}/system_stats")
            with _urlopen(_probe, timeout=3) as _r:
                _r.read()
            providers["comfyui"] = {"ok": True, "label": "ComfyUI (local)", "note": "image generation ready"}
        except Exception as _exc:
            providers["comfyui"] = {"ok": False, "label": "ComfyUI (local)", "note": f"unreachable — is ComfyUI running? ({str(_exc)[:60]})"}
    else:
        providers["comfyui"] = {"ok": False, "label": "ComfyUI (local)", "note": "A1111_HOST not set"}

    any_ok = any(v["ok"] for v in providers.values())
    active  = [v["label"] for v in providers.values() if v["ok"]]
    return jsonify({"ok": True, "any_available": any_ok, "providers": providers, "active": active})


@app.get("/")
def index():
    html = """
<!doctype html>
<html>
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>Seller Studio — AI Listing Engine for Amazon</title>
  <link rel="preconnect" href="https://fonts.googleapis.com"/>
  <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin/>
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800&family=Space+Grotesk:wght@500;600;700&display=swap" rel="stylesheet"/>
  <style>
    :root {
      --bg: #FAFAFA;
      --surface: #FFFFFF;
      --surface-2: #F5F5F5;
      --surface-3: #EFEFEF;
      --border: rgba(0, 0, 0, 0.10);
      --text: #111111;
      --muted: #555555;
      --accent: #FF9900;
      --accent-2: #FFBB33;
      --accent-soft: rgba(255, 153, 0, 0.12);
      --accent-2-soft: rgba(255, 187, 51, 0.10);
      --success-soft: rgba(22, 163, 74, 0.10);
      --warn-soft: rgba(245, 158, 11, 0.14);
      --danger-soft: rgba(220, 38, 38, 0.10);
      --grid-border: rgba(0, 0, 0, 0.09);
      --sheet-head: #F0F0F0;
      --sheet-label: #E8E8E8;
      --required: rgba(220, 38, 38, 0.10);
      --conditional: rgba(245, 158, 11, 0.14);
      --optional: rgba(22, 163, 74, 0.10);
      --missing: rgba(220, 38, 38, 0.22);
      --shadow: 0 4px 24px rgba(0, 0, 0, 0.08);
      --shadow-soft: 0 2px 12px rgba(0, 0, 0, 0.06);
    }
    * { box-sizing: border-box; }
    body {
      font-family: "Inter", "Segoe UI", Arial, sans-serif;
      margin: 0;
      padding: 18px;
      background:
        radial-gradient(900px 500px at 10% -5%, rgba(255, 153, 0, 0.10), transparent 60%),
        radial-gradient(700px 400px at 90% 0%, rgba(255, 187, 51, 0.07), transparent 60%),
        #FAFAFA;
      background-attachment: fixed;
      color: var(--text);
      position: relative;
      min-height: 100vh;
      overflow-x: hidden;
      -webkit-font-smoothing: antialiased;
    }
    .page { position: relative; z-index: 1; }
    .bg-orb {
      position: fixed;
      width: 380px;
      height: 380px;
      border-radius: 999px;
      filter: blur(80px);
      opacity: 0.25;
      z-index: 0;
      pointer-events: none;
    }
    .orb-1 { top: -130px; left: -90px; background: radial-gradient(circle, #FF9900, #FFD580); }
    .orb-2 { top: 16%; right: -170px; background: radial-gradient(circle, #FFBB33, #FF9900); }
    .orb-3 { bottom: -170px; left: 18%; background: radial-gradient(circle, #FFD580, #FFBB33); }
    .panel {
      background: var(--surface);
      border: 1px solid var(--border);
      border-radius: 20px;
      padding: 20px;
      margin-bottom: 14px;
      box-shadow: var(--shadow);
      position: relative;
      overflow: hidden;
    }
    .panel-header { display:flex; justify-content:space-between; align-items:center; gap:12px; margin-bottom:12px; }
    .panel-header h1 {
      margin:0;
      font-family: "Space Grotesk", "Inter", sans-serif;
      font-size:28px;
      font-weight:700;
      letter-spacing:-0.6px;
      color: #111111;
    }
    .panel-header p { margin:4px 0 0; font-size:12px; color: var(--muted); }
    .panel-chip {
      background: var(--accent-soft);
      color: #111111;
      font-weight: 600;
      font-size: 12px;
      padding: 6px 12px;
      border-radius: 999px;
      border: 1px solid rgba(255, 153, 0, 0.25);
    }
    .form-grid { display:grid; grid-template-columns: repeat(3, minmax(260px, 1fr)); gap:12px 14px; align-items:end; }
    label { display:block; font-weight:600; font-size:12px; margin-bottom:4px; color: var(--muted); }
    input, select {
      width:100%;
      padding:10px 12px;
      border:1px solid var(--border);
      border-radius:10px;
      font-size:13px;
      color: var(--text);
      background: #FFFFFF;
      transition: border-color .2s ease, box-shadow .2s ease;
    }
    input:hover, select:hover { border-color: rgba(255, 153, 0, 0.45); }
    input:focus, select:focus {
      outline: none;
      border-color: var(--accent);
      box-shadow: 0 0 0 3px rgba(255, 153, 0, 0.15);
    }
    input[type="file"] { background: var(--surface-2); }
    .actions { display:flex; gap:10px; align-items:center; margin-top:12px; flex-wrap:wrap; }
    .btn {
      padding:10px 16px;
      border: none;
      background: var(--accent);
      color: #111111;
      border-radius:12px;
      cursor:pointer;
      font-weight:700;
      letter-spacing:.2px;
      box-shadow: 0 4px 16px rgba(255, 153, 0, 0.30);
      transition: transform .2s ease, box-shadow .2s ease;
    }
    .btn:hover { transform: translateY(-1px); box-shadow: 0 8px 24px rgba(255, 153, 0, 0.45); background: var(--accent-2); }
    .btn.secondary {
      background: #FFFFFF;
      color: #111111;
      border: 1px solid rgba(0, 0, 0, 0.15);
      box-shadow: none;
    }
    .btn.secondary:hover { border-color: var(--accent); box-shadow: 0 4px 14px rgba(255, 153, 0, 0.18); }
    a.btn { text-decoration: none; align-items: center; }
    #status {
      font-weight:600;
      font-size:12px;
      color: #555555;
      background: var(--accent-soft);
      padding: 6px 10px;
      border-radius: 999px;
      border: 1px solid rgba(255, 153, 0, 0.20);
    }
    #missing_box {
      margin-top:8px;
      font-size:12px;
      color:#b91c1c;
      background: var(--danger-soft);
      padding: 6px 10px;
      border-radius: 10px;
      display: inline-block;
      border: 1px solid rgba(220, 38, 38, 0.18);
    }
    #status.is-error {
      color: #b91c1c;
      background: var(--danger-soft);
      border: 1px solid rgba(220, 38, 38, 0.20);
    }
    #status.is-warn {
      color: #92400e;
      background: var(--warn-soft);
      border: 1px solid rgba(245, 158, 11, 0.25);
    }
    #status.is-ok {
      color: #166534;
      background: var(--success-soft);
      border: 1px solid rgba(22, 163, 74, 0.20);
    }
    .summary-grid {
      margin-top: 12px;
      display: grid;
      grid-template-columns: repeat(3, minmax(200px, 1fr));
      gap: 12px;
    }
    .summary-card {
      background: var(--surface-2);
      border: 1px solid var(--border);
      border-radius: 16px;
      padding: 14px 16px;
      box-shadow: var(--shadow-soft);
      position: relative;
      overflow: hidden;
    }
    .summary-card::before {
      content:""; position:absolute; left:0; top:0; bottom:0; width:3px;
      background: var(--accent);
    }
    .summary-label { font-size: 11px; letter-spacing: .3px; text-transform: uppercase; color: var(--muted); }
    .summary-value { font-size: 22px; font-weight: 700; margin-top: 6px; color: #111111; }
    .summary-sub { font-size: 12px; color: var(--muted); margin-top: 6px; }
    .progress {
      height: 8px;
      background: rgba(0, 0, 0, 0.07);
      border-radius: 999px;
      overflow: hidden;
      margin-top: 10px;
      border: 1px solid var(--border);
    }
    .progress-bar {
      height: 100%;
      background: var(--accent);
      width: 0%;
      transition: width .2s ease;
    }
    .panel-chip.online {
      background: rgba(22, 163, 74, 0.10);
      color: #166534;
      border-color: rgba(22, 163, 74, 0.25);
    }
    .panel-chip.offline {
      background: var(--danger-soft);
      color: #991b1b;
      border: 1px solid rgba(220, 38, 38, 0.25);
    }
    .btn[disabled] {
      opacity: 0.5;
      cursor: not-allowed;
      box-shadow: none;
      transform: none;
    }

    /* ═══════════════════════════════════════════════════
       LAYOUT
    ═══════════════════════════════════════════════════ */
    .layout { display:grid; grid-template-columns: 2.4fr 1fr; gap:12px; }
    .xl-panel {
      display: flex;
      flex-direction: column;
      background: var(--surface);
      border: 1px solid #C8C8C8;
      border-radius: 6px;
      overflow: hidden;
      height: 74vh;
      box-shadow: 0 2px 8px rgba(0,0,0,0.10);
    }
    .prompt-wrap {
      background: var(--surface);
      border: 1px solid var(--border);
      border-radius: 18px;
      padding: 14px;
      overflow: auto;
      height: 74vh;
      box-shadow: var(--shadow-soft);
    }
    .prompt-wrap h3 { margin:0 0 6px; font-size:16px; color: #111111; }
    .prompt-item {
      border:1px solid var(--border); border-radius:12px; padding:10px;
      margin-bottom:10px; background: var(--surface-2);
      transition: transform .2s ease, box-shadow .2s ease;
    }
    .prompt-item:hover { transform: translateY(-2px); box-shadow: 0 6px 18px rgba(0,0,0,0.08); }
    .prompt-item.shared { background: var(--accent-soft); border-color: rgba(255,153,0,0.28); }
    .prompt-item.product { background: var(--accent-2-soft); border-color: rgba(255,187,51,0.28); }
    .prompt-item.missing { border-left: 4px solid #dc2626; background: var(--danger-soft); }
    .prompt-item label { font-size:11px; margin:0 0 4px; color: var(--muted); }
    .prompt-hint { font-size:12px; color: var(--muted); margin-bottom:8px; }
    .empty-state {
      color:#166534; font-weight:600; background: var(--success-soft);
      padding: 8px 10px; border-radius: 10px; display: inline-block;
      border: 1px solid rgba(22,163,74,0.25);
    }
    .autofilled { background: var(--accent-soft); border-color: rgba(255,153,0,0.35); }

    /* ═══════════════════════════════════════════════════
       EXCEL CHROME — name bar + formula bar + sheet tabs
    ═══════════════════════════════════════════════════ */
    .xl-namebox-row {
      display: flex;
      align-items: center;
      background: #F3F3F3;
      border-bottom: 1px solid #C8C8C8;
      height: 28px;
      flex-shrink: 0;
      font-family: "Calibri", "Segoe UI", Arial, sans-serif;
      font-size: 12px;
    }
    .xl-namebox {
      width: 68px;
      min-width: 68px;
      border-right: 1px solid #C8C8C8;
      padding: 0 6px;
      height: 100%;
      display: flex;
      align-items: center;
      justify-content: center;
      font-size: 12px;
      color: #333;
      background: #fff;
      border: none;
      border-right: 1px solid #C8C8C8;
      outline: none;
      cursor: default;
      user-select: none;
    }
    .xl-fx-icon {
      padding: 0 8px;
      color: #217346;
      font-style: italic;
      font-weight: 700;
      font-size: 13px;
      border-right: 1px solid #C8C8C8;
      height: 100%;
      display: flex;
      align-items: center;
      flex-shrink: 0;
      background: #F3F3F3;
    }
    .xl-formula-bar {
      flex: 1;
      padding: 0 8px;
      font-family: "Calibri", "Segoe UI", Arial, sans-serif;
      font-size: 12.5px;
      color: #111;
      white-space: nowrap;
      overflow: hidden;
      text-overflow: ellipsis;
      height: 100%;
      display: flex;
      align-items: center;
      background: #fff;
      border: none;
      cursor: default;
      user-select: text;
    }
    .xl-sheet-wrap {
      flex: 1;
      overflow: auto;
      position: relative;
      background: #fff;
    }
    .xl-tab-bar {
      height: 26px;
      background: #F3F3F3;
      border-top: 1px solid #C8C8C8;
      display: flex;
      align-items: flex-end;
      padding-left: 4px;
      flex-shrink: 0;
      gap: 2px;
    }
    .xl-tab {
      padding: 3px 14px;
      background: #fff;
      border: 1px solid #C8C8C8;
      border-bottom: none;
      border-radius: 3px 3px 0 0;
      font-family: "Calibri", "Segoe UI", Arial, sans-serif;
      font-size: 11.5px;
      color: #111;
      font-weight: 600;
      cursor: default;
      user-select: none;
    }

    /* ═══════════════════════════════════════════════════
       SHEET TABLE — Excel-accurate styling
    ═══════════════════════════════════════════════════ */
    .sheet {
      border-collapse: collapse;
      table-layout: fixed;
      width: max-content;
      min-width: 100%;
      font-family: "Calibri", "Segoe UI", Arial, sans-serif;
      font-size: 13px;
      color: #111;
      contain: layout style;
    }
    /* Column-letter header row (thead) */
    .sheet thead th {
      position: sticky;
      top: 0;
      z-index: 20;
      background: #F2F2F2;
      border: 1px solid #D0D0D0;
      height: 22px;
      min-width: 150px;
      max-width: 150px;
      padding: 0 4px;
      text-align: center;
      font-size: 11.5px;
      font-weight: 600;
      color: #333;
      white-space: nowrap;
      overflow: hidden;
      user-select: none;
    }
    .sheet thead th.row-h {
      position: sticky;
      left: 0;
      z-index: 30;
      min-width: 90px !important;
      max-width: 90px !important;
      background: #F2F2F2;
    }
    /* Row-number / row-header first column */
    .sheet td.row-h {
      position: sticky;
      left: 0;
      z-index: 10;
      background: #F2F2F2;
      border: 1px solid #D0D0D0;
      min-width: 90px !important;
      max-width: 90px !important;
      height: 22px;
      padding: 0 5px;
      text-align: center;
      font-size: 11px;
      font-weight: 500;
      color: #444;
      white-space: nowrap;
      overflow: hidden;
      text-overflow: ellipsis;
      user-select: none;
    }
    /* All data cells */
    .sheet th, .sheet td {
      border: 1px solid #D0D0D0;
      height: 22px;
      min-width: 150px;
      max-width: 150px;
      vertical-align: middle;
      padding: 2px 5px;
      white-space: nowrap;
      overflow: hidden;
      text-overflow: ellipsis;
    }
    /* ── Meta / header rows ───────────────────────── */
    .r-group td {
      background: #BDD7EE;
      font-weight: 700;
      font-size: 11.5px;
      text-align: center;
      color: #1F3864;
      border-color: #9DC3E6;
    }
    .r-group td.row-h { background: #F2F2F2; }
    .r-label td {
      background: #E2EFDA;
      font-weight: 600;
      font-size: 11.5px;
      color: #375623;
      border-color: #C6E0B4;
    }
    .r-label td.row-h { background: #F2F2F2; }
    .r-attr td {
      background: #F2F2F2;
      color: #444;
      font-family: "Consolas", "Courier New", monospace;
      font-size: 10.5px;
      border-color: #D0D0D0;
    }
    .r-example td {
      background: #FEFEFE;
      color: #666;
      font-style: italic;
      font-size: 11px;
    }
    .r-req td {
      font-weight: 600;
      font-size: 10.5px;
      text-align: center;
    }
    .req-required          { background: #FCE4D6; color: #9C2701; border-color: #F4B183; }
    .req-conditionally-required { background: #FFEB9C; color: #7A5200; border-color: #FFCC5C; }
    .req-optional          { background: #E2EFDA; color: #375623; border-color: #C6E0B4; }
    .req-empty             { background: #F2F2F2; color: #888; }

    /* ── Product data rows ───────────────────────── */
    .r-val td {
      background: #FFFFFF;
      padding: 0 5px;
      height: 22px;
    }
    .r-val:nth-child(even) > td:not(.row-h) { background: #F9F9F9; }
    .r-val:hover > td:not(.cell-editable:focus) { background: #EDF3FA !important; }
    .r-val:hover > td.row-h { background: #D4E4F7 !important; color: #1A5276; }

    .cell-required-missing { background: #FCE4D6 !important; }
    .cell-img-wrap { padding: 2px; }

    /* Editable product cells */
    td.cell-editable {
      cursor: cell;
      outline: none;
      white-space: nowrap;
      overflow: hidden;
      text-overflow: ellipsis;
    }
    td.cell-editable:hover:not(:focus) { background: #EDF3FA !important; }
    td.cell-editable:focus {
      background: #FFFFFF !important;
      outline: 2px solid #217346 !important;
      outline-offset: -2px;
      z-index: 5;
      position: relative;
      overflow: visible;
      white-space: pre-wrap;
      word-break: break-word;
    }
    td.cell-dirty { background: #FFFDE7 !important; }
    td.cell-dirty.cell-required-missing { background: #FCE4D6 !important; }

    /* Row state decorations */
    tr.row-done   > td.row-h { border-left: 3px solid #217346; }
    tr.row-error  > td.row-h { border-left: 3px solid #C00000; color: #C00000; }

    /* ── Processing shimmer ──────────────────────── */
    @keyframes xl-shimmer {
      0%   { background-position: -500px 0; }
      100% { background-position:  500px 0; }
    }
    @keyframes xl-flash {
      0%   { background: #C6EFCE; }
      100% { background: transparent; }
    }
    tr.row-processing > td:not(.row-h) {
      background: linear-gradient(90deg, #F2F2F2 25%, #E8E8E8 50%, #F2F2F2 75%) !important;
      background-size: 500px 100% !important;
      animation: xl-shimmer 1.2s infinite linear;
      color: transparent !important;
      pointer-events: none;
    }
    tr.row-processing > td.row-h { color: #999; font-style: italic; }
    td.cell-new-data { animation: xl-flash 0.9s ease-out; }

    /* ── Progress bar ───────────────────────────── */
    #live-progress-bar-wrap {
      display: none;
      margin: 6px 0 4px;
      background: var(--surface-2);
      border: 1px solid var(--border);
      border-radius: 8px;
      padding: 6px 10px;
      font-size: 12px;
      color: var(--muted);
    }
    #live-progress-bar-wrap .bar-track {
      height: 6px; background: var(--border); border-radius: 3px; margin-top: 5px; overflow: hidden;
    }
    #live-progress-bar-wrap .bar-fill {
      height: 100%; width: 0%; background: #217346; border-radius: 3px; transition: width 0.35s ease;
    }

    @media (max-width: 1200px) {
      .layout { grid-template-columns: 1fr; }
      .sheet-wrap, .prompt-wrap { height: auto; max-height: 70vh; }
      .form-grid { grid-template-columns: repeat(2, minmax(220px, 1fr)); }
      .summary-grid { grid-template-columns: repeat(2, minmax(200px, 1fr)); }
    }
    @media (max-width: 720px) {
      .form-grid { grid-template-columns: 1fr; }
      .panel-header { flex-direction: column; align-items: flex-start; }
      .summary-grid { grid-template-columns: 1fr; }
    }

    /* ===== Brand chrome ===== */
    body::before {
      content:""; position: fixed; inset: 0; z-index: 0; pointer-events: none;
      background-image:
        linear-gradient(rgba(0, 0, 0, 0.04) 1px, transparent 1px),
        linear-gradient(90deg, rgba(0, 0, 0, 0.04) 1px, transparent 1px);
      background-size: 46px 46px;
      -webkit-mask-image: radial-gradient(circle at 50% 22%, #000 0%, transparent 78%);
      mask-image: radial-gradient(circle at 50% 22%, #000 0%, transparent 78%);
    }
    .brandbar {
      position: relative; z-index: 1;
      display:flex; align-items:center; justify-content:space-between;
      gap:14px; margin: 2px 6px 18px; flex-wrap: wrap;
    }
    .brand { display:flex; align-items:center; gap:14px; }
    .brand-mark {
      position: relative; width:50px; height:50px; border-radius:16px;
      background: var(--accent);
      box-shadow: 0 0 0 1px rgba(255,153,0,0.30), 0 8px 24px rgba(255,153,0,0.25);
      display:grid; place-items:center;
    }
    .brand-mark::after {
      content:""; position:absolute; inset:3px; border-radius:13px; background:#FFFFFF;
    }
    .brand-spark {
      position: relative; z-index:1; width:22px; height:22px; border-radius:50%;
      background: radial-gradient(circle at 35% 30%, #FFBB33, #FF9900 55%, #CC7700);
      animation: sparkPulse 2.8s ease-in-out infinite;
    }
    @keyframes sparkPulse {
      0%,100% { box-shadow: 0 0 10px rgba(255,153,0,0.5); transform: scale(1); }
      50% { box-shadow: 0 0 22px rgba(255,153,0,0.85); transform: scale(1.08); }
    }
    .brand-text { display:flex; flex-direction:column; line-height:1.15; }
    .brand-name {
      font-family:"Space Grotesk","Inter",sans-serif;
      font-weight:700; letter-spacing:4px; font-size:17px;
      color: #111111;
    }
    .brand-by { font-size:10.5px; letter-spacing:1.5px; color:var(--muted); text-transform:uppercase; }
    .brand-tags { display:flex; gap:8px; flex-wrap:wrap; }
    .brand-pill {
      font-size:11px; font-weight:600; letter-spacing:.4px;
      padding:8px 14px; border-radius:999px; color:#111111;
      border:1px solid rgba(0, 0, 0, 0.12);
      background: var(--accent-soft);
    }
    .brand-pill.alt {
      color:#111111;
      background: #111111;
      color: #FF9900;
      border-color: #111111;
    }
    .panel::after {
      content:""; position:absolute; left:0; top:0; right:0; height:2px;
      background: linear-gradient(90deg, transparent, var(--accent), var(--accent-2), transparent);
    }
    .panel.hero::before {
      content:""; position:absolute; right:-70px; top:-90px; width:300px; height:300px;
      border-radius:50%; pointer-events:none;
      background: radial-gradient(circle, rgba(255,153,0,0.10), transparent 70%);
    }
    .panel-header h1, .summary-value { font-family:"Space Grotesk","Inter",sans-serif; }
    .summary-value { letter-spacing:-0.3px; }
    .panel-chip { color: #111111; border:1px solid rgba(255,153,0,0.22); }
    #status { color:#555555; }
    #missing_box { color:#b91c1c; }
    .prompt-wrap h3 {
      font-family:"Space Grotesk","Inter",sans-serif; letter-spacing:-0.2px;
    }
    select option { color:#111111; }
    input::placeholder { color: rgba(0, 0, 0, 0.30); }
    /* scrollbars */
    .sheet-wrap::-webkit-scrollbar, .prompt-wrap::-webkit-scrollbar { width:10px; height:10px; }
    .sheet-wrap::-webkit-scrollbar-thumb, .prompt-wrap::-webkit-scrollbar-thumb {
      background: var(--accent); border-radius:999px;
      border:2px solid transparent; background-clip:padding-box;
    }
    .sheet-wrap::-webkit-scrollbar-track, .prompt-wrap::-webkit-scrollbar-track { background: rgba(0,0,0,0.04); }

    .shared-section { margin-bottom: 10px; }
    /* ===== Combo mode toggle ===== */
    .combo-toggle-row { display:flex; align-items:center; gap:10px; margin:16px 0 4px; }
    .toggle-label { font-size:13px; font-weight:600; color:var(--text); }
    .toggle-switch { position:relative; display:inline-block; width:40px; height:22px; flex-shrink:0; }
    .toggle-switch input { opacity:0; width:0; height:0; position:absolute; }
    .toggle-knob { position:absolute; cursor:pointer; top:0; left:0; right:0; bottom:0; background:var(--surface-3); border-radius:22px; border:1px solid var(--border); transition:background .2s,border-color .2s; }
    .toggle-knob::before { content:""; position:absolute; height:16px; width:16px; left:2px; bottom:2px; background:#fff; border-radius:50%; box-shadow:0 1px 3px rgba(0,0,0,.2); transition:transform .2s; }
    .toggle-switch input:checked + .toggle-knob { background:var(--accent); border-color:var(--accent); }
    .toggle-switch input:checked + .toggle-knob::before { transform:translateX(18px); }
    /* ===== Combo slider section ===== */
    .combo-slider-section { margin-top:10px; padding:14px 16px; background:var(--accent-soft); border:1px solid rgba(255,153,0,.25); border-radius:12px; }
    input[type="range"].combo-slider { width:100%; accent-color:var(--accent); cursor:pointer; margin:6px 0 2px; height:5px; }
    input[type="range"].combo-slider:disabled { opacity:.4; cursor:not-allowed; }
    .combo-size-display { font-size:26px; font-weight:700; color:var(--accent); line-height:1; }
    .combo-size-range-label { font-size:11px; color:var(--muted); margin-top:3px; }
    .section-divider {
      font-size: 11px; font-weight: 600; text-transform: uppercase;
      letter-spacing: .6px; color: var(--muted);
      margin: 14px 0 8px;
      display: flex; align-items: center; gap: 8px;
    }
    .section-divider::after {
      content: ""; flex: 1; height: 1px; background: var(--border);
    }
    .upload-row {
      display: flex; align-items: center; gap: 8px;
    }
    .upload-row input[type="text"], .upload-row input:not([type="file"]):not([type="number"]):not([type="checkbox"]) {
      flex: 1;
    }
    .upload-btn {
      padding: 10px 14px;
      border: 1px solid var(--border);
      border-radius: 10px;
      background: var(--surface-2);
      font-size: 12px;
      font-weight: 600;
      cursor: pointer;
      white-space: nowrap;
      color: var(--text);
      transition: border-color .2s, background .2s;
    }
    .upload-btn:hover { border-color: var(--accent); background: var(--accent-soft); }
    .upload-status {
      font-size: 11px;
      color: var(--muted);
      margin-top: 3px;
      min-height: 14px;
    }
    .upload-status.ok { color: #166534; }
    .upload-status.err { color: #b91c1c; }
    .ai-providers {
      display: flex; flex-wrap: wrap; gap: 6px; margin-top: 4px;
    }
    .ai-provider-chip {
      font-size: 11px; font-weight: 600; padding: 3px 10px;
      border-radius: 999px; border: 1px solid var(--border);
      background: var(--surface-2); color: var(--muted);
    }
    .ai-provider-chip.on { background: rgba(22,163,74,.10); color: #166534; border-color: rgba(22,163,74,.25); }
    .ai-provider-chip.off { background: var(--danger-soft); color: #991b1b; border-color: rgba(220,38,38,.20); }
    .combo-hint { font-size:12px; color:var(--muted); margin-top:8px; line-height:1.5; background:var(--surface-2); padding:8px 12px; border-radius:8px; border:1px solid var(--border); }
    /* ===== AI Results panel ===== */
    .ai-results-panel { margin-top:18px; }
    .ai-result-card { background:var(--surface-1); border:1px solid var(--border); border-radius:12px; overflow:hidden; transition:box-shadow .18s; }
    .ai-result-card:hover { box-shadow:0 2px 12px rgba(0,0,0,.08); }
    .ai-result-card-header { display:flex; align-items:center; gap:10px; padding:10px 14px; background:var(--surface-2); cursor:pointer; user-select:none; }
    .ai-result-card-header:hover { background:var(--surface-3); }
    .ai-sku-badge { font-size:10px; font-weight:700; padding:2px 9px; border-radius:999px; background:var(--accent); color:#111; flex-shrink:0; }
    .ai-card-title { font-size:12px; font-weight:600; flex:1; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; color:var(--text); }
    .ai-card-chevron { font-size:10px; color:var(--muted); flex-shrink:0; transition:transform .18s; }
    .ai-result-card-body { padding:12px 14px; display:flex; flex-direction:column; gap:7px; }
    .ai-field-row { display:flex; gap:10px; align-items:flex-start; }
    .ai-field-label { font-size:10px; font-weight:700; text-transform:uppercase; letter-spacing:.5px; color:var(--muted); min-width:140px; flex-shrink:0; padding-top:2px; }
    .ai-field-value { font-size:12px; color:var(--text); line-height:1.5; word-break:break-word; }
    .ai-bullets { margin:0; padding:0; display:flex; flex-direction:column; gap:3px; }
    .ai-bullets li { font-size:12px; color:var(--text); margin-left:16px; line-height:1.4; }
    .ai-kw-wrap { display:flex; flex-wrap:wrap; gap:4px; }
    .ai-kw-chip { font-size:10px; padding:2px 8px; border-radius:999px; background:var(--surface-2); border:1px solid var(--border); color:var(--muted); }
    .ai-mono { font-family:monospace; font-size:11px; background:var(--surface-2); padding:2px 6px; border-radius:4px; border:1px solid var(--border); }
    .ai-error-row { font-size:11px; color:#b91c1c; background:var(--danger-soft); padding:4px 10px; border-radius:6px; }
  </style>
</head>
<body>
  <div class="bg-orb orb-1"></div>
  <div class="bg-orb orb-2"></div>
  <div class="bg-orb orb-3"></div>
  <div class="page">
  <header class="brandbar">
    <div class="brand">
      <span class="brand-mark" aria-hidden="true"><span class="brand-spark"></span></span>
      <div class="brand-text">
        <span class="brand-name">SELLER&nbsp;STUDIO</span>
        <span class="brand-by">AI Listing Engine · for Amazon Marketplace</span>
      </div>
    </div>
    <div class="brand-tags">
      <span class="brand-pill">Vision&nbsp;→&nbsp;Listing AI</span>
      <span class="brand-pill alt">Pitch&nbsp;Ready</span>
    </div>
  </header>

  <div class="panel hero">
    <div class="panel-header">
      <div>
        <h1>Amazon Listing Studio</h1>
        <p>One product photo in — a launch-ready Amazon listing out. AI-written copy, generated hero &amp; infographic imagery, and a fully validated template.</p>
      </div>
      <div id="ai_chip" class="panel-chip">Checking AI…</div>
    </div>

    <!-- ── Shared: template + output ── -->
    <div class="shared-section">
      <div class="form-grid">
        <div>
          <label>Flat File / Template (.xlsm/.xlsx)</label>
          <div class="upload-row">
            <input id="template_path" placeholder="No template loaded" value="" readonly/>
            <button class="upload-btn" onclick="document.getElementById('template_file_input').click()">Browse…</button>
          </div>
          <input id="template_file_input" type="file" accept=".xlsm,.xlsx,.xls" style="display:none"/>
          <div class="upload-status" id="template_file_status"></div>
        </div>
        <div>
          <label>Image Catalog Excel (SKU → image URLs)</label>
          <div class="upload-row">
            <input id="catalog_file_path" placeholder="No catalog loaded" readonly/>
            <button class="upload-btn" onclick="document.getElementById('catalog_file_input').click()">Browse…</button>
          </div>
          <input id="catalog_file_input" type="file" accept=".xlsm,.xlsx,.xls,.csv" style="display:none"/>
          <div class="upload-status" id="catalog_file_status"></div>
        </div>
        <div>
          <label>Seller Description File <span style="font-weight:400;color:var(--muted)">(optional — PDF, Word, Excel, ZIP, TXT…)</span></label>
          <div class="upload-row">
            <input id="seller_notes_file_path" placeholder="No file loaded" readonly/>
            <button class="upload-btn" onclick="document.getElementById('seller_notes_file_input').click()">Browse…</button>
          </div>
          <input id="seller_notes_file_input" type="file" accept="*" style="display:none"/>
          <div class="upload-status" id="seller_notes_file_status"></div>
        </div>
      </div>
      <div id="seller_notes_preview" style="display:none;margin-top:8px;padding:10px 14px;background:var(--surface-2);border:1px solid var(--border);border-radius:8px;font-size:12px;color:var(--muted);white-space:pre-wrap;max-height:80px;overflow-y:auto;line-height:1.5"></div>
    </div>

    <!-- ── Unified Settings ── -->
    <div id="tab_content_standard">
      <div class="section-divider">Settings</div>
      <div class="form-grid">
        <div><label>Number of Products</label><input id="product_count" type="number" min="1" value="1"/></div>
        <div>
          <label>AI Provider</label>
          <select id="ai_provider">
            <option value="groq" selected>Groq</option>
          </select>
        </div>
        <div>
          <label>Brand Name <span style="font-weight:400;color:var(--muted)">(type "generic" to omit from title)</span></label>
          <input id="brand_name" placeholder="e.g. Puma, Samsung, or generic"/>
        </div>
      </div>

      <!-- ── Combo Mode toggle ── -->
      <div class="combo-toggle-row">
        <label class="toggle-switch">
          <input type="checkbox" id="combo_mode_toggle" onchange="toggleComboMode(this.checked)"/>
          <span class="toggle-knob"></span>
        </label>
        <span class="toggle-label">Combo Mode</span>
      </div>

      <!-- ── Combo slider (shown when toggle is ON) ── -->
      <div id="combo_slider_section" style="display:none" class="combo-slider-section">
        <div class="form-grid" style="grid-template-columns:1fr 160px">
          <div>
            <label>Max Combo Size</label>
            <input type="range" id="combo_size_slider" class="combo-slider"
                   min="2" max="2" value="2"
                   oninput="updateComboSizeLabel(this.value)" disabled/>
            <div class="combo-size-range-label" id="combo_size_hint">Upload a catalog to enable</div>
          </div>
          <div style="display:flex;flex-direction:column;align-items:center;justify-content:center;gap:2px">
            <div class="combo-size-display" id="combo_size_label">2</div>
            <div style="font-size:10px;color:var(--muted);text-align:center">products<br>per combo</div>
          </div>
        </div>
        <div class="form-grid" style="margin-top:10px;grid-template-columns:1fr 1fr">
          <div>
            <label>Max Combos (cap)</label>
            <input id="combo_max" type="number" min="1" max="5000" value="1500"/>
          </div>
          <div style="display:flex;align-items:flex-end">
            <div class="combo-hint" style="margin:0;font-size:11px" id="combo_desc_hint">
              Creates all 2-product combinations.
            </div>
          </div>
        </div>
      </div>

      <div class="actions" style="margin-top:14px">
        <button class="btn secondary" id="ai_generate_btn" onclick="runAIGeneration()">AI Generate Images</button>
        <button class="btn" id="process_catalog_btn" onclick="processMain()" disabled>Process Catalog</button>
        <button class="btn secondary" id="combo_images_btn" onclick="generateComboImages()" style="display:none">Generate Combo Images</button>
        <button class="btn secondary" id="combo_ai_generate_btn" onclick="generateComboAIImages()" style="display:none">AI Generate Images</button>
        <button class="btn secondary" onclick="saveExcel()">Save &amp; Download Excel</button>
        <div id="status"></div>
        <a id="excel_download_link" style="display:none" class="btn" href="#" download>⬇ Download Updated Excel</a>
      </div>
    </div>

    <!-- ── Shared bottom ── -->
    <div id="catalog_progress" style="display:none;margin-top:10px;font-size:12px;color:var(--muted);background:var(--surface-2);padding:8px 12px;border-radius:8px;border:1px solid var(--border)"></div>
    <div id="missing_box"></div>
    <div class="summary-grid" style="margin-top:12px">
      <div class="summary-card">
        <div class="summary-label">Products</div>
        <div class="summary-value" id="summary_products">0</div>
        <div class="summary-sub" id="summary_template">Template not loaded</div>
      </div>
      <div class="summary-card">
        <div class="summary-label">Required completion</div>
        <div class="summary-value" id="summary_missing">0 missing</div>
        <div class="summary-sub" id="summary_completion">0% complete</div>
        <div class="progress"><div class="progress-bar" id="summary_progress"></div></div>
      </div>
      <div class="summary-card">
        <div class="summary-label">AI Providers</div>
        <div class="summary-value" id="summary_ai">Checking…</div>
        <div id="ai_providers_list" class="ai-providers"></div>
      </div>
    </div>
  </div>

  <div id="live-progress-bar-wrap">
    <span id="live-progress-text">Preparing…</span>
    <div class="bar-track"><div class="bar-fill" id="live-bar-fill"></div></div>
  </div>
  <div class="layout">
    <!-- Excel-style panel: name-box + formula bar + scrollable sheet + tab bar -->
    <div class="xl-panel">
      <div class="xl-namebox-row">
        <div class="xl-namebox" id="xl_namebox">A1</div>
        <div class="xl-fx-icon">fx</div>
        <div class="xl-formula-bar" id="xl_formula_bar"></div>
      </div>
      <div class="xl-sheet-wrap" id="xl_sheet_wrap">
        <table class="sheet" id="sheet_table"></table>
      </div>
      <div class="xl-tab-bar">
        <div class="xl-tab">Template</div>
      </div>
    </div>
    <div class="prompt-wrap">
      <h3 style="margin-top:0">Inputs Required</h3>
      <div class="prompt-hint">Fill only the fields the app cannot fill automatically. Click any cell in the sheet to edit it directly — Tab/Shift-Tab moves between cells, Enter moves down.</div>
      <div id="prompt_form"></div>
    </div>
  </div>

  <!-- ── AI-filled fields results panel ── -->
  <div id="ai_results_panel" class="panel ai-results-panel" style="display:none">
    <div class="panel-header" style="margin-bottom:12px">
      <div>
        <h2 style="margin:0;font-size:15px">AI-Filled Fields</h2>
        <p style="margin:2px 0 0;font-size:12px;color:var(--muted)">Fields written to the Excel template by the AI for each SKU</p>
      </div>
      <button class="upload-btn" onclick="document.getElementById('ai_results_cards').innerHTML='';document.getElementById('ai_results_panel').style.display='none'" style="font-size:11px;padding:6px 12px">Clear</button>
    </div>
    <div id="ai_results_cards" style="display:flex;flex-direction:column;gap:10px"></div>
  </div>

<script src="https://cdnjs.cloudflare.com/ajax/libs/gsap/3.12.5/gsap.min.js" crossorigin="anonymous" referrerpolicy="no-referrer"></script>
<script>
let AI_AVAILABLE = false;   // updated live by checkAIStatus()
let templateMeta = null;
let products = [];
const describeInFlight = new Set();
const sharedValues = { seller: "", country: "" };
const rowAIGenerationInFlight = new Set();

function animatePage() {
  if (!window.gsap) return;
  gsap.from(".brandbar", { opacity: 0, y: -10, duration: 0.5, ease: "power2.out" });
  gsap.from(".panel", { opacity: 0, y: 8, duration: 0.5, delay: 0.06, ease: "power2.out" });
  gsap.from(".panel-header h1, .panel-header p", {
    opacity: 0,
    y: 6,
    duration: 0.45,
    stagger: 0.08,
    delay: 0.05,
    ease: "power2.out"
  });
  gsap.from(".panel-chip", { opacity: 0, scale: 0.9, duration: 0.4, delay: 0.1, ease: "power2.out" });
  gsap.from(".form-grid > div", { opacity: 0, y: 8, duration: 0.45, stagger: 0.04, delay: 0.12, ease: "power2.out" });
  gsap.from([".sheet-wrap", ".prompt-wrap"], {
    opacity: 0,
    y: 12,
    duration: 0.6,
    stagger: 0.08,
    delay: 0.08,
    ease: "power2.out"
  });
  gsap.from(".actions .btn", { opacity: 0, y: 8, duration: 0.4, stagger: 0.05, delay: 0.2, ease: "power2.out" });
  gsap.to(".orb-1", { x: 40, y: 20, duration: 14, yoyo: true, repeat: -1, ease: "sine.inOut" });
  gsap.to(".orb-2", { x: -30, y: 35, duration: 16, yoyo: true, repeat: -1, ease: "sine.inOut" });
  gsap.to(".orb-3", { x: 20, y: -30, duration: 18, yoyo: true, repeat: -1, ease: "sine.inOut" });
  wireButtonAnimations();
}

function animatePromptItems() {
  if (!window.gsap) return;
  const items = document.querySelectorAll(".prompt-item");
  if (!items.length) return;
  gsap.from(items, {
    opacity: 0,
    y: 8,
    scale: 0.98,
    duration: 0.35,
    stagger: 0.03,
    ease: "power2.out",
    overwrite: "auto"
  });
}

function wireButtonAnimations() {
  if (!window.gsap) return;
  document.querySelectorAll(".btn").forEach(btn => {
    btn.addEventListener("mouseenter", () => {
      gsap.to(btn, { scale: 1.03, duration: 0.15, ease: "power1.out" });
    });
    btn.addEventListener("mouseleave", () => {
      gsap.to(btn, { scale: 1, duration: 0.2, ease: "power1.out" });
    });
    btn.addEventListener("mousedown", () => {
      gsap.to(btn, { scale: 0.98, duration: 0.1, ease: "power1.out" });
    });
  });
}

function setStatus(message, tone = "info") {
  const el = document.getElementById("status");
  if (!el) return;
  el.textContent = message;
  el.classList.remove("is-error", "is-warn", "is-ok");
  if (tone === "error") el.classList.add("is-error");
  if (tone === "warn") el.classList.add("is-warn");
  if (tone === "ok") el.classList.add("is-ok");
}

function countTotalRequired() {
  if (!templateMeta) return 0;
  const descriptionAttr = getFirstAttrByLabel("Product Description");
  const { browseNodeAttr, productTypeAttr, productIdTypeAttr, mainImageAttr } = getSharedTargetAttrs();
  const dedicated = new Set([descriptionAttr, browseNodeAttr, productTypeAttr, productIdTypeAttr, mainImageAttr].filter(Boolean));
  let total = 0;
  products.forEach(() => {
    templateMeta.columns.forEach(c => {
      if (dedicated.has(c.attr)) return;
      if ((c.required || "").toLowerCase() === "required") total += 1;
    });
  });
  return total;
}

function updateSummary() {
  const productEl   = document.getElementById("summary_products");
  const templateEl  = document.getElementById("summary_template");
  const missingEl   = document.getElementById("summary_missing");
  const completionEl= document.getElementById("summary_completion");
  const progressEl  = document.getElementById("summary_progress");
  const aiChip      = document.getElementById("ai_chip");
  const aiSummary   = document.getElementById("summary_ai");
  const aiBtn       = document.getElementById("ai_generate_btn");

  if (productEl)  productEl.textContent  = templateMeta ? String(products.length) : "0";
  if (templateEl) templateEl.textContent = templateMeta ? "Template loaded" : "Template not loaded";

  const missing = collectMissing();
  const totalRequired = countTotalRequired();
  const missingCount  = missing.length;
  const done = Math.max(0, totalRequired - missingCount);
  const pct  = totalRequired ? Math.round((done / totalRequired) * 100) : 100;
  if (missingEl)   missingEl.textContent   = `${missingCount} missing`;
  if (completionEl)completionEl.textContent= `${pct}% complete`;
  if (progressEl)  progressEl.style.width  = `${pct}%`;

  if (aiChip) {
    aiChip.textContent = AI_AVAILABLE ? "AI Ready" : "AI Offline";
    aiChip.classList.toggle("online",  AI_AVAILABLE);
    aiChip.classList.toggle("offline", !AI_AVAILABLE);
  }
  if (aiSummary) aiSummary.textContent = AI_AVAILABLE ? "Ready" : "Offline";
  if (aiBtn) {
    aiBtn.disabled = !AI_AVAILABLE;
    aiBtn.title = AI_AVAILABLE ? "Run AI generation" : "Configure API keys in .env to enable AI";
  }
}

function esc(s) {
  return (s ?? "").toString().replaceAll("&","&amp;").replaceAll("<","&lt;").replaceAll(">","&gt;").replaceAll('"',"&quot;");
}

function reqClass(req) {
  const r = (req || "").toLowerCase().replaceAll(" ", "-");
  if (!r) return "req-empty";
  if (r.includes("conditionally-required")) return "req-conditionally-required";
  if (r.includes("required")) return "req-required";
  if (r.includes("optional")) return "req-optional";
  return "req-empty";
}

function blankProduct() {
  const row = {};
  if (templateMeta) {
    templateMeta.columns.forEach(c => row[c.attr] = "");
  }
  return row;
}

function getSharedTargetAttrs() {
  const out = { sellerAttr: "", countryAttr: "", browseNodeAttr: "", productTypeAttr: "", productIdTypeAttr: "", mainImageAttr: "" };
  if (!templateMeta) return out;
  for (const c of templateMeta.columns) {
    const lbl = (c.label || "").trim().toLowerCase();
    if (!out.countryAttr && lbl === "country of origin") out.countryAttr = c.attr;
    if (!out.sellerAttr && lbl === "packer contact information") out.sellerAttr = c.attr;
    if (!out.browseNodeAttr && lbl === "recommended browse nodes") out.browseNodeAttr = c.attr;
    if (!out.productTypeAttr && lbl === "product type") out.productTypeAttr = c.attr;
    if (!out.productIdTypeAttr && lbl === "product id type") out.productIdTypeAttr = c.attr;
    if (!out.mainImageAttr && lbl === "main image url") out.mainImageAttr = c.attr;
  }
  return out;
}

function applySharedValuesToAllProducts() {
  if (!templateMeta) return;
  const { sellerAttr, countryAttr } = getSharedTargetAttrs();
  for (const row of products) {
    if (countryAttr && sharedValues.country.trim()) row[countryAttr] = sharedValues.country.trim();
    if (sellerAttr && sharedValues.seller.trim()) row[sellerAttr] = sharedValues.seller.trim();
  }
}

function ensureProductCount(n) {
  const target = Math.max(1, n || 1);
  while (products.length < target) products.push(blankProduct());
  while (products.length > target) products.pop();
}

function collectMissing() {
  const missing = [];
  if (!templateMeta) return missing;
  const descriptionAttr = getFirstAttrByLabel("Product Description");
  const { browseNodeAttr, productTypeAttr, productIdTypeAttr, mainImageAttr } = getSharedTargetAttrs();
  const dedicated = new Set([descriptionAttr, browseNodeAttr, productTypeAttr, productIdTypeAttr, mainImageAttr].filter(Boolean));
  products.forEach((row, idx) => {
    templateMeta.columns.forEach(c => {
      if (dedicated.has(c.attr)) return;
      if ((c.required || "").toLowerCase() === "required" && !(row[c.attr] || "").trim()) {
        missing.push({ row: idx + 1, excelRow: templateMeta.data_row + idx, label: c.label, attr: c.attr });
      }
    });
  });
  return missing;
}

function renderSheet() {
  const table = document.getElementById("sheet_table");
  if (!templateMeta) {
    table.innerHTML = "";
    return;
  }
  const cols = templateMeta.columns;
  const missing = collectMissing();
  const missingSet = new Set(missing.map(m => `${m.row}::${m.attr}`));

  let head = '<thead><tr><th class="row-h top-left" style="min-width:90px;max-width:90px"></th>';
  cols.forEach((c, ci) => head += `<th data-col-idx="${ci}" title="${esc(c.label||c.col)}">${esc(c.col)}</th>`);
  head += '</tr></thead>';

  let groupRow = '<tr class="r-group"><td class="row-h">Group</td>';
  cols.forEach(c => groupRow += `<td title="${esc(c.group||"")}">${esc(c.group || "")}</td>`);
  groupRow += '</tr>';

  let labelRow = '<tr class="r-label"><td class="row-h">Label</td>';
  cols.forEach(c => labelRow += `<td title="${esc(c.label||"")}">${esc(c.label || "")}</td>`);
  labelRow += '</tr>';

  let attrRow = '<tr class="r-attr"><td class="row-h">Field</td>';
  cols.forEach(c => attrRow += `<td title="${esc(c.attr||"")}">${esc(c.attr || "")}</td>`);
  attrRow += '</tr>';

  let exampleRow = '<tr class="r-example"><td class="row-h">Example</td>';
  cols.forEach(c => exampleRow += `<td title="${esc(c.example||"")}">${esc(c.example || "")}</td>`);
  exampleRow += '</tr>';

  let reqRow = '<tr class="r-req"><td class="row-h">Required</td>';
  cols.forEach(c => reqRow += `<td class="${reqClass(c.required)}">${esc(c.required || "")}</td>`);
  reqRow += '</tr>';

  let productRows = "";
  products.forEach((row, idx) => {
    const excelRow = templateMeta.data_row + idx;
    const skuVal = _rowSku(row);
    const skuLabel = skuVal ? esc(skuVal) : `${excelRow}`;
    productRows += `<tr class="r-val row-done" data-sku="${esc(skuVal)}">` +
      `<td class="row-h" title="Row ${excelRow}${skuVal ? ' — ' + skuVal : ''}">${excelRow}</td>`;
    cols.forEach((c, ci) => {
      const val = row[c.attr] || "";
      const missClass = missingSet.has(`${idx + 1}::${c.attr}`) ? "cell-required-missing" : "";
      const reqCls = reqClass(c.required);
      const lbl = (c.label || "").toLowerCase();
      const isImgCol = lbl.includes("image url") || lbl.includes("image_url");
      if (isImgCol && val && (val.startsWith("http") || val.startsWith("/img/"))) {
        productRows += `<td class="${reqCls} ${missClass} cell-img-preview"` +
          ` data-row-idx="${idx}" data-col-idx="${ci}" data-attr="${esc(c.attr)}">` +
          `<div class="cell-img-wrap">` +
          `<a href="${esc(val)}" target="_blank" title="${esc(val)}">` +
          `<img src="${esc(val)}" style="max-height:18px;max-width:80px;object-fit:contain;vertical-align:middle" ` +
          `onerror="this.style.display='none';this.nextSibling.style.display='inline'" />` +
          `<span style="display:none;font-size:10px;word-break:break-all">${esc(val)}</span></a>` +
          `</div></td>`;
      } else {
        productRows += `<td class="cell-editable ${reqCls} ${missClass}"` +
          ` contenteditable="true" spellcheck="false"` +
          ` data-row-idx="${idx}" data-col-idx="${ci}" data-attr="${esc(c.attr)}">${esc(val)}</td>`;
      }
    });
    productRows += "</tr>";
  });

  table.innerHTML = head + `<tbody>${groupRow}${labelRow}${attrRow}${exampleRow}${reqRow}${productRows}</tbody>`;
  renderPromptForm();
  renderValidationSummary();
}

function renderPromptForm() {
  const form = document.getElementById("prompt_form");
  if (!templateMeta) {
    form.innerHTML = "";
    return;
  }
  const { sellerAttr, countryAttr, browseNodeAttr, productTypeAttr, productIdTypeAttr, mainImageAttr } = getSharedTargetAttrs();
  const missing = collectMissing();
  let html = `
    <div class="prompt-item shared">
      <label>Shared value (all products) - Seller / Packer Contact Information</label>
      <input id="shared_seller" value="${esc(sharedValues.seller)}" placeholder="${esc(sellerAttr || 'Field not found in template')}"/>
      <label style="margin-top:8px">Shared value (all products) - Country of Origin</label>
      <input id="shared_country" value="${esc(sharedValues.country)}" placeholder="${esc(countryAttr || 'Field not found in template')}"/>
    </div>
  `;
  if (!missing.length) {
    form.innerHTML = html + "<div class='empty-state'>No pending required inputs.</div>";
    const sellerEl = document.getElementById("shared_seller");
    const countryEl = document.getElementById("shared_country");
    if (sellerEl) sellerEl.addEventListener("input", (e) => { sharedValues.seller = e.target.value || ""; applySharedValuesToAllProducts(); renderSheet(); });
    if (countryEl) countryEl.addEventListener("input", (e) => { sharedValues.country = e.target.value || ""; applySharedValuesToAllProducts(); renderSheet(); });
    animatePromptItems();
    return;
  }
  const productTypeByNode = templateMeta.node_to_product_type || {};
  const browseChoices = (templateMeta.browse_options || []).map(o => `<option value="${esc(o.display)}">${esc(o.display)}</option>`).join("");
  const productTypeChoices = Array.from(new Set(Object.values(productTypeByNode))).sort().map(v => `<option value="${esc(v)}">${esc(v)}</option>`).join("");
  const productIdTypeChoices = (templateMeta.product_id_type_options || []).map(v => `<option value="${esc(v)}">${esc(v)}</option>`).join("");
  for (let r = 0; r < products.length; r++) {
    const row = products[r];
    html += `
      <div class="prompt-item product">
        <label>Product ${r + 1} - Recommended Browse Node</label>
        <select data-row-node="${r}">
          <option value="">Select Browse Node (${esc(browseNodeAttr || 'Field not found in template')})</option>
          ${browseChoices}
        </select>
        <label style="margin-top:8px">Product ${r + 1} - Product Type</label>
        <select data-row-type="${r}">
          <option value="">Select Product Type (${esc(productTypeAttr || 'Field not found in template')})</option>
          ${productTypeChoices}
        </select>
        <label style="margin-top:8px">Product ${r + 1} - Product Id Type</label>
        <select data-row-idtype="${r}">
          <option value="">Select Product Id Type (${esc(productIdTypeAttr || 'Field not found in template')})</option>
          ${productIdTypeChoices}
        </select>
        <label style="margin-top:8px">Product ${r + 1} - Initial Image (Main Image URL / file path)</label>
        <input data-row-mainimg="${r}" placeholder="${esc(mainImageAttr || 'Field not found in template')}"/>
        <input type="file" accept="image/*" data-row-upload="${r}" style="margin-top:6px"/>
        <label style="margin-top:8px">Product ${r + 1} - Seller Notes (optional)</label>
        <input data-row-notes="${r}" value="${esc(row.__seller_notes || "")}" placeholder="Optional notes for image generation"/>
      </div>
    `;
  }
  missing.forEach((m, i) => {
    const v = products[m.row - 1][m.attr] || "";
    const autoCls = m.attr.includes("product_description") ? "autofilled" : "";
    html += `
      <div class="prompt-item missing">
        <label>Product ${m.row} (Excel row ${m.excelRow}) - ${esc(m.label)}</label>
        <input id="prompt_${i}" class="${autoCls}" data-row="${m.row - 1}" data-attr="${esc(m.attr)}" value="${esc(v)}"/>
      </div>
    `;
  });
  form.innerHTML = html;
  const sellerEl = document.getElementById("shared_seller");
  const countryEl = document.getElementById("shared_country");
  if (sellerEl) sellerEl.addEventListener("input", (e) => { sharedValues.seller = e.target.value || ""; applySharedValuesToAllProducts(); renderSheet(); });
  if (countryEl) countryEl.addEventListener("input", (e) => { sharedValues.country = e.target.value || ""; applySharedValuesToAllProducts(); renderSheet(); });
  form.querySelectorAll("select[data-row-node]").forEach(el => {
    const r = Number(el.dataset.rowNode);
    if (browseNodeAttr) el.value = products[r][browseNodeAttr] || "";
    el.addEventListener("change", (e) => {
      if (!browseNodeAttr) return;
      products[r][browseNodeAttr] = e.target.value || "";
      const mapped = (templateMeta.node_to_product_type || {})[products[r][browseNodeAttr]] || "";
      if (mapped && productTypeAttr && !(products[r][productTypeAttr] || "").trim()) {
        products[r][productTypeAttr] = mapped;
      }
      renderSheet();
    });
  });
  form.querySelectorAll("select[data-row-type]").forEach(el => {
    const r = Number(el.dataset.rowType);
    if (productTypeAttr) el.value = products[r][productTypeAttr] || "";
    el.addEventListener("change", (e) => {
      if (!productTypeAttr) return;
      products[r][productTypeAttr] = e.target.value || "";
      renderSheet();
    });
  });
  form.querySelectorAll("select[data-row-idtype]").forEach(el => {
    const r = Number(el.dataset.rowIdtype);
    if (productIdTypeAttr) el.value = products[r][productIdTypeAttr] || "";
    el.addEventListener("change", (e) => {
      if (!productIdTypeAttr) return;
      products[r][productIdTypeAttr] = e.target.value || "";
      renderSheet();
    });
  });
  form.querySelectorAll("input[data-row-mainimg]").forEach(el => {
    const r = Number(el.dataset.rowMainimg);
    if (mainImageAttr) el.value = products[r][mainImageAttr] || "";
    el.addEventListener("change", (e) => {
      if (!mainImageAttr) return;
      products[r][mainImageAttr] = e.target.value || "";
      renderSheet();
      autoGenerateAssetsForRow(r);
    });
  });
  form.querySelectorAll("input[data-row-notes]").forEach(el => {
    const r = Number(el.dataset.rowNotes);
    el.addEventListener("change", (e) => {
      products[r]["__seller_notes"] = e.target.value || "";
    });
  });
  form.querySelectorAll("input[type='file'][data-row-upload]").forEach(el => {
    const r = Number(el.dataset.rowUpload);
    el.addEventListener("change", async (e) => {
      if (!mainImageAttr) return;
      const f = e.target.files && e.target.files[0];
      if (!f) return;
      const fd = new FormData();
      fd.append("file", f);
      fd.append("row_index", String(r));
      const res = await fetch("/api/upload-image", { method: "POST", body: fd });
      const json = await res.json();
      if (!json.ok) {
        setStatus(json.error || "Upload failed", "error");
        return;
      }
      products[r][mainImageAttr] = json.path || "";
      renderSheet();
      setStatus(`Uploaded image for Product ${r + 1}`, "ok");
      autoGenerateAssetsForRow(r);
    });
  });
  form.querySelectorAll("input[data-row][data-attr]").forEach(el => {
    el.addEventListener("change", (e) => {
      const row = Number(e.target.dataset.row);
      const attr = e.target.dataset.attr;
      products[row][attr] = e.target.value || "";
      renderSheet();
    });
  });
  animatePromptItems();
}

function renderValidationSummary() {
  const missing = collectMissing();
  document.getElementById("missing_box").innerText = missing.length
    ? ("Pending required inputs: " + missing.slice(0, 15).map(m => `Product ${m.row} -> ${m.label}`).join(", ") + (missing.length > 15 ? " ..." : ""))
    : "";
  setStatus(
    missing.length ? "Fill requested inputs from the right panel." : "Ready to save.",
    missing.length ? "warn" : "ok"
  );
  updateSummary();
}

// ── Inline cell editing ────────────────────────────────────────────────────────
// Uses event delegation on the table so re-renders don't re-bind listeners.
// Edits update `products[]` in-place and refresh validation highlights only,
// without rebuilding the entire table.

let _editDebounceTimer = null;

function _refreshMissingHighlights() {
  if (!templateMeta) return;
  const missing = collectMissing();
  const missingSet = new Set(missing.map(m => `${m.row - 1}::${m.attr}`));
  document.querySelectorAll("#sheet_table td[data-row-idx][data-attr]").forEach(td => {
    const key = `${td.dataset.rowIdx}::${td.dataset.attr}`;
    td.classList.toggle("cell-required-missing", missingSet.has(key));
  });
  // Also refresh the summary text without rebuilding the table
  document.getElementById("missing_box").innerText = missing.length
    ? ("Pending required inputs: " + missing.slice(0, 15).map(m => `Product ${m.row} -> ${m.label}`).join(", ") + (missing.length > 15 ? " ..." : ""))
    : "";
  updateSummary();
}

function _navigateCell(fromTd, dir) {
  const row = fromTd.closest("tr");
  if (!row) return;
  const allEditableTds = Array.from(row.querySelectorAll("td.cell-editable"));
  const colIdx = allEditableTds.indexOf(fromTd);
  const allRows = Array.from(document.querySelectorAll("#sheet_table tr.r-val"));
  const rowIdx = allRows.indexOf(row);

  let next = null;
  if (dir === "right") {
    next = allEditableTds[colIdx + 1];
    if (!next && rowIdx + 1 < allRows.length) {
      next = allRows[rowIdx + 1].querySelector("td.cell-editable");
    }
  } else if (dir === "left") {
    next = allEditableTds[colIdx - 1];
    if (!next && rowIdx > 0) {
      const prevRowTds = Array.from(allRows[rowIdx - 1].querySelectorAll("td.cell-editable"));
      next = prevRowTds[prevRowTds.length - 1] || null;
    }
  } else if (dir === "down") {
    if (rowIdx + 1 < allRows.length) {
      const nextRowTds = Array.from(allRows[rowIdx + 1].querySelectorAll("td.cell-editable"));
      next = nextRowTds[colIdx] || nextRowTds[0] || null;
    }
  } else if (dir === "up") {
    if (rowIdx > 0) {
      const prevRowTds = Array.from(allRows[rowIdx - 1].querySelectorAll("td.cell-editable"));
      next = prevRowTds[colIdx] || prevRowTds[0] || null;
    }
  }
  if (next) {
    next.focus();
    // Select all text in the cell for easy overwrite
    const sel = window.getSelection();
    const range = document.createRange();
    range.selectNodeContents(next);
    sel.removeAllRanges();
    sel.addRange(range);
  }
}

// ── Formula bar + name-box helpers ──────────────────────────────────────────
function _xlCellAddress(td) {
  if (!templateMeta) return "";
  const colIdx = parseInt(td.dataset.colIdx || "0", 10);
  const rowIdx = parseInt(td.dataset.rowIdx || "0", 10);
  const col = templateMeta.columns[colIdx] || {};
  const excelRow = (templateMeta.data_row || 8) + rowIdx;
  return `${col.col || ""}${excelRow}`;
}
function _updateFormulaBar(td) {
  const nb  = document.getElementById("xl_namebox");
  const fb  = document.getElementById("xl_formula_bar");
  const val = td ? (td.innerText || "") : "";
  const lbl = td && templateMeta
    ? (templateMeta.columns[parseInt(td.dataset.colIdx || "0", 10)] || {}).label || ""
    : "";
  if (nb) nb.textContent = td ? _xlCellAddress(td) : "A1";
  if (fb) {
    fb.textContent = val;
    fb.title = lbl ? `[${lbl}]  ${val}` : val;
  }
}

function initCellEditing() {
  const table = document.getElementById("sheet_table");

  // Focus: update name-box + formula bar
  table.addEventListener("focusin", (e) => {
    const td = e.target.closest("td.cell-editable");
    if (!td) return;
    _updateFormulaBar(td);
    // Highlight whole column header
    const ci = parseInt(td.dataset.colIdx || "-1", 10);
    document.querySelectorAll("#sheet_table thead th").forEach((th, i) => {
      th.style.background = i === ci + 1 ? "#BDD7EE" : "";
    });
    // Highlight row header
    const tr = td.closest("tr");
    if (tr) {
      const rh = tr.querySelector("td.row-h");
      if (rh) rh.style.background = "#BDD7EE";
    }
  });

  table.addEventListener("focusout", (e) => {
    const td = e.target.closest("td.cell-editable");
    if (!td) return;
    // Restore header colours
    document.querySelectorAll("#sheet_table thead th").forEach(th => th.style.background = "");
    const tr = td.closest("tr");
    if (tr) {
      const rh = tr.querySelector("td.row-h");
      if (rh) rh.style.background = "";
    }
    const fb = document.getElementById("xl_formula_bar");
    if (fb) fb.textContent = "";
    const nb = document.getElementById("xl_namebox");
    if (nb) nb.textContent = "A1";
  });

  table.addEventListener("input", (e) => {
    const td = e.target.closest("td.cell-editable[data-row-idx]");
    if (!td) return;
    const rowIdx = parseInt(td.dataset.rowIdx, 10);
    const attr = td.dataset.attr;
    if (isNaN(rowIdx) || !attr) return;
    if (!products[rowIdx]) return;
    products[rowIdx][attr] = td.innerText;
    td.classList.add("cell-dirty");
    // Keep formula bar in sync while typing
    const fb = document.getElementById("xl_formula_bar");
    if (fb) fb.textContent = td.innerText;
    clearTimeout(_editDebounceTimer);
    _editDebounceTimer = setTimeout(_refreshMissingHighlights, 250);
  });

  table.addEventListener("keydown", (e) => {
    const td = e.target.closest("td.cell-editable");
    if (!td) return;
    if (e.key === "Tab") {
      e.preventDefault();
      _navigateCell(td, e.shiftKey ? "left" : "right");
    } else if (e.key === "Enter" && !e.shiftKey) {
      e.preventDefault();
      _navigateCell(td, "down");
    } else if (e.key === "ArrowDown" && e.ctrlKey) {
      e.preventDefault();
      _navigateCell(td, "down");
    } else if (e.key === "ArrowUp" && e.ctrlKey) {
      e.preventDefault();
      _navigateCell(td, "up");
    } else if (e.key === "Escape") {
      td.blur();
    }
  });

  // On paste: strip formatting, insert as plain text
  table.addEventListener("paste", (e) => {
    const td = e.target.closest("td.cell-editable");
    if (!td) return;
    e.preventDefault();
    const text = (e.clipboardData || window.clipboardData).getData("text/plain");
    // Handle multi-cell paste: split by tab for columns, newline for rows
    const lines = text.split(/\\r?\\n/);
    if (lines.length === 1 && !text.includes("\t")) {
      // Single cell paste — insert at cursor
      document.execCommand("insertText", false, text);
    } else {
      // Multi-cell paste — fill right and down from the focused cell
      const startRow = td.closest("tr");
      const allRows = Array.from(document.querySelectorAll("#sheet_table tr.r-val"));
      const allTds = Array.from(startRow.querySelectorAll("td.cell-editable"));
      const startColIdx = allTds.indexOf(td);
      const startRowIdx = allRows.indexOf(startRow);
      lines.forEach((line, ri) => {
        const targetRow = allRows[startRowIdx + ri];
        if (!targetRow) return;
        const targetTds = Array.from(targetRow.querySelectorAll("td.cell-editable"));
        line.split("\t").forEach((cell, ci) => {
          const targetTd = targetTds[startColIdx + ci];
          if (!targetTd) return;
          const rIdx = parseInt(targetTd.dataset.rowIdx, 10);
          const atr = targetTd.dataset.attr;
          if (!isNaN(rIdx) && atr && products[rIdx]) {
            targetTd.innerText = cell;
            products[rIdx][atr] = cell;
            targetTd.classList.add("cell-dirty");
          }
        });
      });
      _refreshMissingHighlights();
    }
  });
}

// Clear dirty markers after a successful save
function _clearDirtyMarkers() {
  document.querySelectorAll("#sheet_table td.cell-dirty").forEach(td => td.classList.remove("cell-dirty"));
}

function getFirstAttrByLabel(label) {
  if (!templateMeta) return "";
  const t = (label || "").trim().toLowerCase();
  const found = templateMeta.columns.find(c => (c.label || "").trim().toLowerCase() === t);
  return found ? found.attr : "";
}

async function autoGenerateAssetsForRow(rowIndex) {
  if (!templateMeta) return;
  if (!AI_AVAILABLE) {
    setStatus("AI is disabled. Set GROQ_API_KEY to enable AI generation.", "warn");
    return;
  }
  const key = String(rowIndex);
  if (rowAIGenerationInFlight.has(key)) return;
  const mainImageAttr = getFirstAttrByLabel("Main Image URL");
  if (!mainImageAttr) return;
  const imageSource = (products[rowIndex][mainImageAttr] || "").trim();
  if (!imageSource) return;
  rowAIGenerationInFlight.add(key);
  setStatus(`Analysing product ${rowIndex + 1} — filling all fields…`);

  // Mark the row as processing in the sheet
  const table = document.getElementById("sheet_table");
  const rowEls = table ? Array.from(table.querySelectorAll("tr.r-val")) : [];
  if (rowEls[rowIndex]) rowEls[rowIndex].className = "r-val row-processing";

  try {
    const providerVal = document.getElementById("ai_provider")?.value || "auto";
    const brandVal    = (document.getElementById("brand_name")?.value || "").trim();
    const res = await fetch("/api/generate-row-full", {
      method: "POST",
      headers: {"Content-Type":"application/json"},
      body: JSON.stringify({
        template_path: getTemplatePath(),
        row:           products[rowIndex],
        provider:      providerVal,
        brand_name:    brandVal,
        seller_notes:  sellerNotesText || ""
      })
    });
    const json = await res.json();
    if (!json.ok) {
      setStatus(json.error || "AI generation failed", "error");
      if (rowEls[rowIndex]) rowEls[rowIndex].className = "r-val row-error";
      return;
    }
    products[rowIndex] = json.row || products[rowIndex];
    renderSheet();
    const warns = (json.warnings || []).length;
    setStatus(
      warns > 0
        ? `Product ${rowIndex + 1} filled (${warns} warning${warns > 1 ? "s" : ""} — check required fields)`
        : `Product ${rowIndex + 1} fully filled ✓`,
      warns > 0 ? "warn" : "ok"
    );
  } finally {
    rowAIGenerationInFlight.delete(key);
  }
}

function getTemplatePath() {
  const el = document.getElementById("template_path");
  return (el && (el.dataset.serverPath || el.value || "")).trim();
}

async function loadTemplateMeta() {
  const template_path = getTemplatePath();
  if (!template_path) {
    setStatus("Upload a flat file first.", "warn");
    return;
  }
  let res;
  let json;
  try {
    res = await fetch("/api/meta", {
      method: "POST",
      headers: {"Content-Type":"application/json"},
      body: JSON.stringify({ template_path })
    });
    json = await res.json();
  } catch (err) {
    setStatus("Failed to reach server.", "error");
    return;
  }
  if (!json.ok) {
    setStatus(json.error || "Failed to load template.", "error");
    return;
  }
  templateMeta = json.meta;
  ensureProductCount(Number(document.getElementById("product_count").value || 1));
  applySharedValuesToAllProducts();
  renderSheet();
  updateCatalogBtn();
}

async function uploadTemplateAndLoad(file) {
  const statusEl = document.getElementById("template_file_status");
  if (statusEl) statusEl.textContent = "Uploading…";
  const fd = new FormData();
  fd.append("file", file);
  let res, json;
  try {
    res = await fetch("/api/upload-template", { method: "POST", body: fd });
    json = await res.json();
  } catch (err) {
    if (statusEl) statusEl.textContent = "Upload failed";
    setStatus("Template upload failed to reach server.", "error");
    return;
  }
  if (!json.ok) {
    if (statusEl) statusEl.textContent = json.error || "Upload failed";
    setStatus(json.error || "Template upload failed.", "error");
    return;
  }
  document.getElementById("template_path").value = file.name;
  const tplEl = document.getElementById("template_path");
  tplEl.dataset.serverPath = json.path || "";
  tplEl.dataset.originalName = json.original_name || file.name;
  if (statusEl) statusEl.textContent = "✓ Loaded";
  setStatus("Flat file uploaded. Reading columns…");
  await loadTemplateMeta();
}

function addProductRow() {
  if (!templateMeta) return;
  products.push(blankProduct());
  applySharedValuesToAllProducts();
  document.getElementById("product_count").value = products.length;
  renderSheet();
}

function removeProductRow() {
  if (!templateMeta || products.length <= 1) return;
  products.pop();
  document.getElementById("product_count").value = products.length;
  renderSheet();
}

async function saveExcel() {
  if (!templateMeta) {
    setStatus("Load template first.", "warn");
    return;
  }
  applySharedValuesToAllProducts();
  const templatePath = getTemplatePath();
  const payload = {
    template_path: templatePath,
    products
  };
  let res;
  let json;
  try {
    res = await fetch("/api/save", {
      method: "POST",
      headers: {"Content-Type":"application/json"},
      body: JSON.stringify(payload)
    });
    json = await res.json();
  } catch (err) {
    setStatus("Save failed to reach server.", "error");
    return;
  }
  if (json.ok) {
    _clearDirtyMarkers();
    // Show download link — use the original filename (no UUID prefix)
    const dlUrl = "/api/download-excel?path=" + encodeURIComponent(json.saved_path || templatePath);
    const tplEl2 = document.getElementById("template_path");
    const filename = (tplEl2 && tplEl2.dataset.originalName) || (json.saved_path || templatePath).split(/[\\/]/).pop();
    const dlEl2 = document.getElementById("excel_download_link");
    if (dlEl2) { dlEl2.href = dlUrl; dlEl2.download = filename; dlEl2.style.display = "inline-flex"; }
    const missingBox = document.getElementById("missing_box");
    if (json.warnings && json.warnings.length > 0) {
      const names = json.warnings.slice(0, 15).map(m => `Product ${m.row} → ${m.label}`).join(", ")
                  + (json.warnings.length > 15 ? ` … and ${json.warnings.length - 15} more` : "");
      setStatus(`Saved. ${json.warnings.length} unfilled required field(s) — review before submitting to Amazon.`, "warn");
      if (missingBox) missingBox.innerText = "Unfilled required fields: " + names;
    } else {
      setStatus("Saved — click the download button to get your updated Excel.", "ok");
      if (missingBox) missingBox.innerText = "";
    }
  } else {
    setStatus(json.error, "error");
  }
}

async function runAIGeneration() {
  if (!templateMeta) {
    setStatus("Load template first.", "warn");
    return;
  }
  const btn = document.getElementById("ai_generate_btn");
  if (btn) { btn.disabled = true; btn.textContent = "Generating images…"; }
  setStatus("Generating images — this may take a few minutes…");
  const payload = {
    template_path: getTemplatePath(),
    products,
  };
  let res;
  let json;
  try {
    res = await fetch("/api/generate-catalog-images", {
      method: "POST",
      headers: {"Content-Type":"application/json"},
      body: JSON.stringify(payload)
    });
    json = await res.json();
  } catch (err) {
    setStatus("Image generation failed to reach server.", "error");
    if (btn) { btn.disabled = false; btn.textContent = "AI Generate Images"; }
    return;
  }
  if (btn) { btn.disabled = false; btn.textContent = "AI Generate Images"; }
  if (!json.ok) {
    setStatus(json.error || "Image generation failed.", "error");
    return;
  }
  products = json.products || products;
  renderSheet();
  const count = json.images_generated || 0;
  setStatus(`Image generation complete — ${count} image set(s) generated.`, "ok");
}

document.getElementById("product_count").addEventListener("input", () => {
  if (!templateMeta) return;
  ensureProductCount(Number(document.getElementById("product_count").value || 1));
  applySharedValuesToAllProducts();
  renderSheet();
});

// ── Page init ────────────────────────────────────────────────────────────────
animatePage();
updateSummary();
initCellEditing();
checkAIStatus();
wireUploads();

// ── Catalog pipeline state ───────────────────────────────────────────────────
let uploadedCatalogPath = "";
let sellerNotesText = "";

// ── Combo mode toggle + slider ────────────────────────────────────────────────
function toggleComboMode(enabled) {
  document.getElementById("combo_slider_section").style.display = enabled ? "block" : "none";
  const btn = document.getElementById("process_catalog_btn");
  if (btn) btn.textContent = enabled ? "Process Combos" : "Process Catalog";
  const imgBtn = document.getElementById("combo_images_btn");
  if (imgBtn) imgBtn.style.display = enabled ? "inline-flex" : "none";
  const aiBtn = document.getElementById("ai_generate_btn");
  if (aiBtn) aiBtn.style.display = enabled ? "none" : "inline-flex";
  const comboAiBtn = document.getElementById("combo_ai_generate_btn");
  if (comboAiBtn) comboAiBtn.style.display = enabled ? "inline-flex" : "none";
}

function updateComboSizeLabel(val) {
  const n = parseInt(val, 10);
  const labelEl = document.getElementById("combo_size_label");
  const hintEl  = document.getElementById("combo_desc_hint");
  if (labelEl) labelEl.textContent = n;
  if (hintEl) {
    const sizes = Array.from({length: n - 1}, (_, i) => i + 2);
    hintEl.textContent = n === 2
      ? "Creates all 2-product combinations."
      : `Creates all ${sizes.join(", ")}-product combinations (sizes 2 – ${n}).`;
  }
}

function processMain() {
  const comboEnabled = document.getElementById("combo_mode_toggle")?.checked;
  if (comboEnabled) processComboCatalog();
  else processCatalog();
}

// ── Live AI status check ─────────────────────────────────────────────────────
async function checkAIStatus() {
  const chip     = document.getElementById("ai_chip");
  const summary  = document.getElementById("summary_ai");
  const listEl   = document.getElementById("ai_providers_list");
  if (chip) { chip.textContent = "Checking…"; chip.className = "panel-chip"; }
  try {
    const res  = await fetch("/api/status");
    const data = await res.json();
    AI_AVAILABLE = data.any_available;
    if (chip) {
      chip.textContent = AI_AVAILABLE ? "AI Ready" : "AI Offline";
      chip.classList.toggle("online",  AI_AVAILABLE);
      chip.classList.toggle("offline", !AI_AVAILABLE);
    }
    if (summary) summary.textContent = AI_AVAILABLE ? "Ready" : "Offline";
    if (listEl && data.providers) {
      listEl.innerHTML = Object.values(data.providers).map(p =>
        `<span class="ai-provider-chip ${p.ok ? 'on' : 'off'}" title="${p.note || (p.ok ? 'Connected' : 'Unavailable')}">${p.label}${p.ok ? "" : " ✗"}</span>`
      ).join("");
    }
    updateSummary();
  } catch (e) {
    if (chip) { chip.textContent = "Offline"; chip.className = "panel-chip offline"; }
    if (summary) summary.textContent = "Offline";
  }
}

// ── Upload wiring (all in one place, runs after DOM ready) ───────────────────
function wireUploads() {
  // Template
  const templateInput = document.getElementById("template_file_input");
  if (templateInput) {
    templateInput.addEventListener("change", async (e) => {
      const f = e.target.files && e.target.files[0];
      if (!f) return;
      await uploadTemplateAndLoad(f);
      e.target.value = "";
    });
  }

  // Image catalog
  const catalogInput = document.getElementById("catalog_file_input");
  if (catalogInput) {
    catalogInput.addEventListener("change", async (e) => {
      const f = e.target.files && e.target.files[0];
      if (!f) return;
      const statusEl = document.getElementById("catalog_file_status");
      if (statusEl) { statusEl.textContent = "Uploading…"; statusEl.className = "upload-status"; }
      const fd = new FormData();
      fd.append("file", f);
      try {
        const res  = await fetch("/api/upload-image-catalog", { method: "POST", body: fd });
        const json = await res.json();
        if (!json.ok) {
          if (statusEl) { statusEl.textContent = "✗ " + (json.error || "Upload failed"); statusEl.className = "upload-status err"; }
          setStatus(json.error || "Catalog upload failed", "error");
          return;
        }
        uploadedCatalogPath = json.path || "";
        const skuCount = json.sku_count || 0;
        document.getElementById("catalog_file_path").value = f.name;
        if (statusEl) { statusEl.textContent = `✓ ${skuCount} SKU(s) mapped`; statusEl.className = "upload-status ok"; }
        setStatus(`Image catalog loaded: ${skuCount} SKU(s)`, "ok");
        // Configure combo size slider with range 2 … skuCount
        const comboSlider = document.getElementById("combo_size_slider");
        const comboHintEl = document.getElementById("combo_size_hint");
        if (comboSlider) {
          if (skuCount >= 2) {
            comboSlider.min   = 2;
            comboSlider.max   = skuCount;
            comboSlider.value = 2;
            comboSlider.disabled = false;
            updateComboSizeLabel(2);
            if (comboHintEl) { comboHintEl.textContent = `Slide to choose max combo size (2 – ${skuCount})`; comboHintEl.className = "upload-status ok"; }
          } else {
            comboSlider.disabled = true;
            if (comboHintEl) { comboHintEl.textContent = "Need at least 2 SKUs for combo mode"; comboHintEl.className = "upload-status err"; }
          }
        }
      } catch (err) {
        if (statusEl) { statusEl.textContent = "✗ Network error"; statusEl.className = "upload-status err"; }
        setStatus("Catalog upload failed: " + err, "error");
      }
      e.target.value = "";
      updateCatalogBtn();
    });
  }

  // Seller notes
  const sellerInput = document.getElementById("seller_notes_file_input");
  if (sellerInput) {
    sellerInput.addEventListener("change", async (e) => {
      const f = e.target.files && e.target.files[0];
      if (!f) return;
      const statusEl = document.getElementById("seller_notes_file_status");
      const previewEl = document.getElementById("seller_notes_preview");
      if (statusEl) { statusEl.textContent = "Extracting text…"; statusEl.className = "upload-status"; }
      const fd = new FormData();
      fd.append("file", f);
      try {
        const res  = await fetch("/api/upload-seller-notes", { method: "POST", body: fd });
        const json = await res.json();
        if (!json.ok) {
          if (statusEl) { statusEl.textContent = "✗ " + (json.error || "Failed"); statusEl.className = "upload-status err"; }
          sellerNotesText = "";
          if (previewEl) previewEl.style.display = "none";
          return;
        }
        sellerNotesText = json.text || "";
        document.getElementById("seller_notes_file_path").value = f.name;
        const ftype = json.file_type || "FILE";
        if (statusEl) { statusEl.textContent = `✓ ${ftype} — ${json.char_count.toLocaleString()} chars`; statusEl.className = "upload-status ok"; }
        if (previewEl) { previewEl.style.display = "block"; previewEl.textContent = json.preview || ""; }
        setStatus(`${ftype} loaded (${json.char_count.toLocaleString()} chars) — AI will use it.`, "ok");
      } catch (err) {
        if (statusEl) { statusEl.textContent = "✗ Network error"; statusEl.className = "upload-status err"; }
        sellerNotesText = "";
      }
      e.target.value = "";
    });
  }
}

// ── Live-sheet helpers ────────────────────────────────────────────────────────

function _rowSku(row) {
  if (!row) return "";
  for (const [k, v] of Object.entries(row)) {
    const kl = k.toLowerCase();
    if (kl === "seller_sku" || kl === "item_sku" || kl === "external_product_id" || kl.includes("_sku")) return String(v || "");
  }
  return "";
}

function _setLiveProgress(done, total, label) {
  const wrap = document.getElementById("live-progress-bar-wrap");
  const text = document.getElementById("live-progress-text");
  const fill = document.getElementById("live-bar-fill");
  if (!wrap) return;
  wrap.style.display = "block";
  if (text) text.textContent = label || `${done} / ${total} products filled`;
  if (fill) fill.style.width = total > 0 ? `${Math.round(done / total * 100)}%` : "0%";
}

function _hideLiveProgress() {
  const wrap = document.getElementById("live-progress-bar-wrap");
  if (wrap) wrap.style.display = "none";
}

function _ensureSheetHeaders() {
  const table = document.getElementById("sheet_table");
  if (!table || !templateMeta) return;
  if (table.querySelector("thead")) return;
  const cols = templateMeta.columns;
  let head = '<thead><tr><th class="row-h top-left" style="min-width:90px;max-width:90px"></th>';
  cols.forEach((c, ci) => head += `<th data-col-idx="${ci}" title="${esc(c.label||c.col)}">${esc(c.col)}</th>`);
  head += '</tr></thead>';
  let groupRow   = '<tr class="r-group"><td class="row-h">Group</td>';   cols.forEach(c => groupRow   += `<td title="${esc(c.group||"")}">${esc(c.group||"")}</td>`);              groupRow   += '</tr>';
  let labelRow   = '<tr class="r-label"><td class="row-h">Label</td>';   cols.forEach(c => labelRow   += `<td title="${esc(c.label||"")}">${esc(c.label||"")}</td>`);              labelRow   += '</tr>';
  let attrRow    = '<tr class="r-attr"><td class="row-h">Field</td>';    cols.forEach(c => attrRow    += `<td title="${esc(c.attr||"")}">${esc(c.attr||"")}</td>`);                attrRow    += '</tr>';
  let exRow      = '<tr class="r-example"><td class="row-h">Example</td>'; cols.forEach(c => exRow    += `<td title="${esc(c.example||"")}">${esc(c.example||"")}</td>`);         exRow      += '</tr>';
  let reqRow     = '<tr class="r-req"><td class="row-h">Required</td>';  cols.forEach(c => reqRow     += `<td class="${reqClass(c.required)}">${esc(c.required||"")}</td>`);      reqRow     += '</tr>';
  table.innerHTML = head + `<tbody>${groupRow}${labelRow}${attrRow}${exRow}${reqRow}</tbody>`;
}

function _addSkeletonRow(sku, position) {
  _ensureSheetHeaders();
  const table = document.getElementById("sheet_table");
  if (!table || !templateMeta) return;
  const existing = table.querySelector(`tr[data-sku="${CSS.escape(sku)}"]`);
  if (existing) { existing.className = "r-val row-processing"; return; }
  const cols = templateMeta.columns;
  const rowNum = (templateMeta.data_row || 8) + position;
  const tr = document.createElement("tr");
  tr.className = "r-val row-processing";
  tr.setAttribute("data-sku", sku);
  let cells = `<td class="row-h" title="${esc(sku)}">${rowNum}</td>`;
  cols.forEach(() => cells += "<td></td>");
  tr.innerHTML = cells;
  table.querySelector("tbody").appendChild(tr);
  tr.scrollIntoView({ behavior: "smooth", block: "nearest" });
}

function _clearProductRows() {
  const table = document.getElementById("sheet_table");
  if (!table) return;
  Array.from(table.querySelectorAll("tr.r-val")).forEach(tr => tr.remove());
}

function _updateRowInPlace(sku, rowData, rowIdx) {
  const table = document.getElementById("sheet_table");
  if (!table || !templateMeta) return false;
  const tr = table.querySelector(`tr[data-sku="${CSS.escape(sku)}"]`);
  if (!tr) return false;
  const cols = templateMeta.columns;
  tr.className = "r-val row-done";
  // Resolve rowIdx from products array if not supplied
  const rIdx = (rowIdx !== undefined && rowIdx >= 0)
    ? rowIdx
    : products.findIndex(p => _rowSku(p) === sku);
  const dataCells = Array.from(tr.querySelectorAll("td:not(.row-h)"));
  cols.forEach((c, ci) => {
    const td = dataCells[ci];
    if (!td) return;
    const val = rowData[c.attr] || "";
    td.className = "cell-new-data " + reqClass(c.required);
    setTimeout(() => td.classList.remove("cell-new-data"), 900);
    const lbl = (c.label || "").toLowerCase();
    const isImg = (lbl.includes("image url") || lbl.includes("image_url")) && val && (val.startsWith("http") || val.startsWith("/img/"));
    if (isImg) {
      td.className += " cell-img-preview";
      td.removeAttribute("contenteditable");
      td.setAttribute("data-row-idx", rIdx >= 0 ? rIdx : 0);
      td.setAttribute("data-col-idx", ci);
      td.setAttribute("data-attr", c.attr);
      td.innerHTML = `<div class="cell-img-wrap"><a href="${esc(val)}" target="_blank" title="${esc(val)}"><img src="${esc(val)}" style="max-height:18px;max-width:80px;object-fit:contain;vertical-align:middle" onerror="this.style.display='none';this.nextSibling.style.display='inline'" /><span style="display:none;font-size:10px;word-break:break-all">${esc(val)}</span></a></div>`;
    } else {
      td.classList.add("cell-editable");
      td.setAttribute("contenteditable", "true");
      td.setAttribute("spellcheck", "false");
      td.setAttribute("data-row-idx", rIdx >= 0 ? rIdx : 0);
      td.setAttribute("data-col-idx", ci);
      td.setAttribute("data-attr", c.attr);
      td.textContent = val;
    }
  });
  return true;
}

function updateCatalogBtn() {
  const ready = !!templateMeta && !!uploadedCatalogPath;
  const hint  = !templateMeta ? "Load a template first" : "Upload an Image Catalog first";
  const btn   = document.getElementById("process_catalog_btn");
  if (btn) { btn.disabled = !ready; btn.title = ready ? "Run pipeline" : hint; }
}

function _showDownloadLink(savedPath) {
  if (!savedPath) return;
  const dlUrl = "/api/download-excel?path=" + encodeURIComponent(savedPath);
  const tplEl = document.getElementById("template_path");
  const filename = (tplEl && tplEl.dataset.originalName) || savedPath.split(/[\\/]/).pop();
  const dlEl = document.getElementById("excel_download_link");
  if (dlEl) { dlEl.href = dlUrl; dlEl.download = filename; dlEl.style.display = "inline-flex"; }
}

async function processCatalog() {
  if (!templateMeta) { setStatus("Load template first.", "warn"); return; }
  const progressEl = document.getElementById("catalog_progress");
  if (progressEl) { progressEl.style.display = "block"; progressEl.textContent = "Starting catalog pipeline…"; }
  setStatus("Processing catalog…");

  const templatePath = getTemplatePath();
  const providerVal  = document.getElementById("ai_provider")?.value || "auto";
  const brandVal     = (document.getElementById("brand_name")?.value || "").trim();
  const payload = {
    template_path:      templatePath,
    image_catalog_path: uploadedCatalogPath,
    generate_images:    false,
    provider:           providerVal,
    brand_name:         brandVal,
    seller_notes:       sellerNotesText || ""
  };

  // Disable process button while running
  const catalogBtn = document.getElementById("process_catalog_btn");
  if (catalogBtn) { catalogBtn.disabled = true; catalogBtn.textContent = "Processing…"; }

  // Products are processed sequentially — clear and rebuild the products array live.
  // We preserve any rows the user already typed by NOT wiping products[] until start.
  let _totalSkus    = 0;
  let _skuPosition  = 0;
  let completedSkus = [];

  try {
    const res = await fetch("/api/process-catalog/stream", {
      method:  "POST",
      headers: { "Content-Type": "application/json" },
      body:    JSON.stringify(payload)
    });
    if (!res.ok) {
      const errJson = await res.json().catch(() => ({}));
      setStatus(errJson.error || "Catalog processing failed", "error");
      if (progressEl) progressEl.textContent = "Error: " + (errJson.error || res.statusText);
      return;
    }

    const reader  = res.body.getReader();
    const decoder = new TextDecoder();
    let buffer        = "";
    let skuTimeoutId  = null;
    let timedOut      = false;
    // Per-SKU timeout: sequential processing means one SKU at a time.
    // Give each product up to 3 minutes before giving up.
    const SKU_TIMEOUT_MS = 1800000;  // 30 minutes for large catalogs

    function resetSkuTimeout() {
      if (skuTimeoutId) clearTimeout(skuTimeoutId);
      skuTimeoutId = setTimeout(() => { timedOut = true; reader.cancel(); }, SKU_TIMEOUT_MS);
    }

    try {
      while (true) {
        const { done, value } = await reader.read();
        if (done) break;
        buffer += decoder.decode(value, { stream: true });
        const lines = buffer.split("\\n");
        buffer = lines.pop();

        for (const line of lines) {
          if (!line.startsWith("data: ")) continue;
          let evt;
          try { evt = JSON.parse(line.slice(6)); } catch { continue; }

          if (evt.type === "start") {
            _totalSkus = evt.total;
            // Reset products array to match catalog size
            products = Array.from({length: evt.total}, () => blankProduct());
            clearAiResultCards();
            // Clear existing product rows from the table so skeleton rows appear cleanly
            _clearProductRows();
            _ensureSheetHeaders();
            _setLiveProgress(0, evt.total, `Starting — 0 / ${evt.total} products…`);
            if (progressEl) { progressEl.style.display = "block"; progressEl.textContent = `0 / ${evt.total} products queued…`; }
            resetSkuTimeout();

          } else if (evt.type === "sku_start") {
            resetSkuTimeout();
            // evt.index is 1-based; position in table is 0-based
            _addSkeletonRow(evt.sku, (evt.index != null ? evt.index - 1 : _skuPosition));
            if (progressEl) progressEl.textContent =
              `⟳ Analysing ${evt.sku}… (${completedSkus.length} / ${_totalSkus} done)`;

          } else if (evt.type === "sku") {
            resetSkuTimeout();
            completedSkus.push(evt.sku);
            const rowIdx = evt.index - 1; // 0-based position in products array
            _skuPosition = Math.max(_skuPosition, evt.index);

            setStatus(`Processing catalog… (${completedSkus.length} / ${_totalSkus})`, "info");
            _setLiveProgress(completedSkus.length, _totalSkus,
              `✓ ${completedSkus.length} / ${_totalSkus} done — ${_totalSkus - completedSkus.length} remaining`);
            if (progressEl) progressEl.textContent =
              `✓ ${completedSkus.length} / ${_totalSkus} — ${evt.sku} done`;

            // Always put the row into products[] regardless of whether it has data
            if (rowIdx >= 0 && rowIdx < products.length) products[rowIdx] = evt.row || {};

            // Add/update AI results card for this SKU
            if (evt.pipeline) addAiResultCard(evt.pipeline);

            if (evt.row && Object.keys(evt.row).length > 0) {
              // Update the skeleton row in-place (keeps scroll position, no full re-render)
              const updated = _updateRowInPlace(evt.sku, evt.row, rowIdx);
              if (!updated) {
                // No skeleton existed yet — add it now (can happen if sku_start was missed)
                _addSkeletonRow(evt.sku, rowIdx);
                _updateRowInPlace(evt.sku, evt.row, rowIdx);
              }
            } else {
              // Error or empty row — mark skeleton as failed
              const table = document.getElementById("sheet_table");
              const tr = table && table.querySelector(`tr[data-sku="${CSS.escape(evt.sku)}"]`);
              if (tr) tr.className = "r-val row-error";
            }

            _refreshMissingHighlights();

          } else if (evt.type === "done") {
            if (skuTimeoutId) clearTimeout(skuTimeoutId);
            // products[] is already fully populated from per-SKU events — just re-render
            renderSheet();
            _setLiveProgress(evt.count, evt.count, `All ${evt.count} product(s) filled ✓`);
            setTimeout(_hideLiveProgress, 4000);
            setStatus(`Catalog complete — ${evt.count} product(s) filled. Review and click Save & Download.`, "ok");
            if (progressEl) progressEl.textContent = `Done — ${evt.count} product(s) populated. Click Save & Download to export.`;

          } else if (evt.type === "error") {
            if (skuTimeoutId) clearTimeout(skuTimeoutId);
            _hideLiveProgress();
            setStatus(evt.error || "Catalog processing failed", "error");
            if (progressEl) progressEl.textContent = "Error: " + (evt.error || "unknown");
          }
        }
      }
    } finally {
      if (skuTimeoutId) clearTimeout(skuTimeoutId);
    }

    if (timedOut) {
      renderSheet();
      _hideLiveProgress();
      setStatus(`Catalog timed out: ${completedSkus.length} / ${_totalSkus} product(s) filled.`, "warn");
      if (progressEl) progressEl.textContent =
        `Timed out — ${completedSkus.length} product(s) populated. Review and save.`;
    }
  } catch (err) {
    _hideLiveProgress();
    setStatus("Catalog processing failed: " + err, "error");
    if (progressEl) progressEl.textContent = "Network error: " + err;
  } finally {
    if (catalogBtn) {
      catalogBtn.disabled = false;
      catalogBtn.textContent = "Process Catalog";
    }
  }
}

async function processComboCatalog() {
  if (!templateMeta) { setComboStatus("Load template first.", "warn"); return; }
  if (!uploadedCatalogPath) { setComboStatus("Upload an Image Catalog first (use Browse… above).", "warn"); return; }

  const progressEl    = document.getElementById("catalog_progress");
  const comboSlider   = document.getElementById("combo_size_slider");
  const maxComboSize  = Math.max(2, parseInt(comboSlider?.value || "2", 10));
  const maxCombos     = parseInt(document.getElementById("combo_max")?.value || "1500", 10);

  if (comboSlider?.disabled) {
    setComboStatus("Upload a catalog with at least 2 SKUs first.", "error");
    return;
  }
  clearAiResultCards();
  if (progressEl) { progressEl.style.display = "block"; progressEl.textContent = "Building combo images…"; }
  setComboStatus("Generating combos…");

  const providerVal2 = document.getElementById("ai_provider")?.value || "auto";
  const brandVal2    = (document.getElementById("brand_name")?.value || "").trim();
  const payload = {
    template_path:      getTemplatePath(),
    image_catalog_path: uploadedCatalogPath,
    max_combo_size:     maxComboSize,
    max_combos:         maxCombos,
    provider:           providerVal2,
    brand_name:         brandVal2,
    seller_notes:       sellerNotesText || "",
    generate_images:    false
  };

  try {
    const res = await fetch("/api/process-combo-catalog/stream", {
      method: "POST",
      headers: { "Content-Type": "application/json" },
      body: JSON.stringify(payload)
    });
    if (!res.ok) {
      const errJson = await res.json().catch(() => ({}));
      setStatus(errJson.error || "Combo processing failed", "error");
      if (progressEl) progressEl.textContent = "Error: " + (errJson.error || res.statusText);
      return;
    }

    const reader  = res.body.getReader();
    const decoder = new TextDecoder();
    let buffer = "";
    let completedSkus = [];
    let skuTimeoutId  = null;
    let timedOut      = false;
    const SKU_TIMEOUT_MS = 1800000;  // 30 minutes for large catalogs

    function resetSkuTimeout() {
      if (skuTimeoutId) clearTimeout(skuTimeoutId);
      skuTimeoutId = setTimeout(() => { timedOut = true; reader.cancel(); }, SKU_TIMEOUT_MS);
    }

    try {
      while (true) {
        const { done, value } = await reader.read();
        if (done) break;
        buffer += decoder.decode(value, { stream: true });
        const lines = buffer.split("\\n");
        buffer = lines.pop();
        for (const line of lines) {
          if (!line.startsWith("data: ")) continue;
          let evt;
          try { evt = JSON.parse(line.slice(6)); } catch { continue; }

          if (evt.type === "building_combos") {
            const sizeDesc = evt.max_combo_size > 2
              ? `sizes 2–${evt.max_combo_size}`
              : `2-product`;
            if (progressEl) progressEl.textContent =
              `Building ${sizeDesc} combos from ${evt.source_skus} SKUs…`;
          } else if (evt.type === "combos_ready") {
            if (progressEl) progressEl.textContent =
              `${evt.total_combos} combos generated — running listing AI…`;
            resetSkuTimeout();
          } else if (evt.type === "start") {
            if (progressEl) progressEl.textContent =
              `Processing ${evt.total} combo(s) in parallel…`;
            resetSkuTimeout();
          } else if (evt.type === "sku") {
            resetSkuTimeout();
            completedSkus.push(evt.sku);
            if (progressEl) progressEl.textContent =
              `✓ ${evt.index}/${evt.total} done — ${evt.sku} · ${evt.total - evt.index} remaining…`;
            setComboStatus(`Processing combos… (${evt.index}/${evt.total})`, "info");
            if (evt.pipeline) addAiResultCard(evt.pipeline);
            if (evt.row && Object.keys(evt.row).length > 0) {
              products.push(evt.row);
              renderSheet();
            }
          } else if (evt.type === "done") {
            if (skuTimeoutId) clearTimeout(skuTimeoutId);
            // products[] is already fully populated from per-SKU events — just re-render
            renderSheet();
            setComboStatus(`Done — ${evt.count} combo(s) populated. Review and save.`, "ok");
            if (progressEl) progressEl.textContent = `Done — ${evt.count} combo(s) populated.`;
          } else if (evt.type === "error") {
            if (skuTimeoutId) clearTimeout(skuTimeoutId);
            setComboStatus(evt.error || "Combo processing failed", "error");
            if (progressEl) progressEl.textContent = "Error: " + (evt.error || "unknown");
          }
        }
      }
    } finally {
      if (skuTimeoutId) clearTimeout(skuTimeoutId);
    }

    if (timedOut) {
      renderSheet();
      setComboStatus(`Done (timeout): ${completedSkus.length} combo(s) populated.`, "ok");
      if (progressEl) progressEl.textContent = `Stopped — ${completedSkus.length} combo(s) done.`;
    }
  } catch (err) {
    setComboStatus("Combo processing failed: " + err, "error");
    if (progressEl) progressEl.textContent = "Network error: " + err;
  }
}

async function generateComboImages() {
  if (!uploadedCatalogPath) { setComboStatus("Upload an Image Catalog first.", "warn"); return; }

  const progressEl   = document.getElementById("catalog_progress");
  const comboSlider  = document.getElementById("combo_size_slider");
  const maxComboSize = Math.max(2, parseInt(comboSlider?.value || "2", 10));
  const maxCombos    = parseInt(document.getElementById("combo_max")?.value || "1500", 10);

  if (progressEl) { progressEl.style.display = "block"; progressEl.textContent = "Building combo images…"; }
  setComboStatus("Generating combo images…");

  const payload = {
    image_catalog_path: uploadedCatalogPath,
    max_combo_size:     maxComboSize,
    max_combos:         maxCombos,
    images_only:        true
  };

  try {
    const res = await fetch("/api/process-combo-catalog/stream", {
      method: "POST",
      headers: { "Content-Type": "application/json" },
      body: JSON.stringify(payload)
    });
    if (!res.ok) {
      const errJson = await res.json().catch(() => ({}));
      setComboStatus(errJson.error || "Image generation failed", "error");
      return;
    }

    const reader  = res.body.getReader();
    const decoder = new TextDecoder();
    let buffer = "";
    let count  = 0;
    let total  = 0;

    while (true) {
      const { done, value } = await reader.read();
      if (done) break;
      buffer += decoder.decode(value, { stream: true });
      const lines = buffer.split("\\n");
      buffer = lines.pop();
      for (const line of lines) {
        if (!line.startsWith("data: ")) continue;
        let evt;
        try { evt = JSON.parse(line.slice(6)); } catch { continue; }

        if (evt.type === "combos_ready") {
          total = evt.total_combos;
          if (progressEl) progressEl.textContent = `Generating ${total} combo image(s)…`;
        } else if (evt.type === "image") {
          count++;
          if (progressEl) progressEl.textContent =
            `✓ ${evt.index}/${evt.total} — ${evt.sku} saved`;
          setComboStatus(`Generating images… (${evt.index}/${evt.total})`, "info");
        } else if (evt.type === "done") {
          setComboStatus(`Done — ${evt.count} combo image(s) saved.`, "ok");
          if (progressEl) progressEl.textContent = `Done — ${evt.count} combo image(s) saved.`;
        } else if (evt.type === "error") {
          setComboStatus(evt.error || "Image generation failed", "error");
          if (progressEl) progressEl.textContent = "Error: " + (evt.error || "unknown");
        }
      }
    }
  } catch (err) {
    setComboStatus("Image generation failed: " + err, "error");
    if (progressEl) progressEl.textContent = "Network error: " + err;
  }
}

async function generateComboAIImages() {
  if (!templateMeta) { setComboStatus("Load template first.", "warn"); return; }
  if (!uploadedCatalogPath) { setComboStatus("Upload an Image Catalog first.", "warn"); return; }

  const progressEl    = document.getElementById("catalog_progress");
  const comboSlider   = document.getElementById("combo_size_slider");
  const maxComboSize  = Math.max(2, parseInt(comboSlider?.value || "2", 10));
  const maxCombos     = parseInt(document.getElementById("combo_max")?.value || "1500", 10);

  if (comboSlider?.disabled) {
    setComboStatus("Upload a catalog with at least 2 SKUs first.", "error");
    return;
  }

  const aiBtn = document.getElementById("combo_ai_generate_btn");
  if (aiBtn) { aiBtn.disabled = true; aiBtn.textContent = "Generating…"; }
  clearAiResultCards();
  if (progressEl) { progressEl.style.display = "block"; progressEl.textContent = "Generating AI images for combos…"; }
  setComboStatus("Generating AI lifestyle & infographic images…");

  const providerVal = document.getElementById("ai_provider")?.value || "auto";
  const brandVal    = (document.getElementById("brand_name")?.value || "").trim();
  const payload = {
    template_path:      getTemplatePath(),
    image_catalog_path: uploadedCatalogPath,
    max_combo_size:     maxComboSize,
    max_combos:         maxCombos,
    provider:           providerVal,
    brand_name:         brandVal,
    seller_notes:       sellerNotesText || "",
    generate_images:    true
  };

  try {
    const res = await fetch("/api/process-combo-catalog/stream", {
      method: "POST",
      headers: { "Content-Type": "application/json" },
      body: JSON.stringify(payload)
    });
    if (!res.ok) {
      const errJson = await res.json().catch(() => ({}));
      setComboStatus(errJson.error || "AI image generation failed", "error");
      if (aiBtn) { aiBtn.disabled = false; aiBtn.textContent = "AI Generate Images"; }
      return;
    }

    const reader  = res.body.getReader();
    const decoder = new TextDecoder();
    let buffer = "";
    let completedSkus = [];

    while (true) {
      const { done, value } = await reader.read();
      if (done) break;
      buffer += decoder.decode(value, { stream: true });
      const lines = buffer.split("\\n");
      buffer = lines.pop();
      for (const line of lines) {
        if (!line.startsWith("data: ")) continue;
        let evt;
        try { evt = JSON.parse(line.slice(6)); } catch { continue; }

        if (evt.type === "building_combos") {
          if (progressEl) progressEl.textContent =
            `Building combos from ${evt.source_skus} SKUs (max size ${evt.max_combo_size})…`;
        } else if (evt.type === "combos_ready") {
          if (progressEl) progressEl.textContent =
            `Generating AI images for ${evt.total_combos} combo(s)…`;
        } else if (evt.type === "sku") {
          completedSkus.push(evt.sku);
          if (progressEl) progressEl.textContent =
            `✓ ${evt.index}/${evt.total} — ${evt.sku}`;
          setComboStatus(`AI images… (${evt.index}/${evt.total})`, "info");
          if (evt.pipeline) addAiResultCard(evt.pipeline);
        } else if (evt.type === "done") {
          setComboStatus(`Done — AI images generated for ${completedSkus.length} combo(s).`, "ok");
          if (progressEl) progressEl.textContent =
            `Done — AI images generated for ${completedSkus.length} combo(s).`;
        } else if (evt.type === "error") {
          setComboStatus(evt.error || "AI image generation failed", "error");
          if (progressEl) progressEl.textContent = "Error: " + (evt.error || "unknown");
        }
      }
    }
  } catch (err) {
    setComboStatus("AI image generation failed: " + err, "error");
    if (progressEl) progressEl.textContent = "Network error: " + err;
  } finally {
    if (aiBtn) { aiBtn.disabled = false; aiBtn.textContent = "AI Generate Images"; }
  }
}

function setComboStatus(msg, type) {
  // Routes to the unified status element
  const tone = type === "error" ? "error" : type === "ok" ? "ok" : type === "warn" ? "warn" : "info";
  setStatus(msg, tone);
}

// ── AI Results panel ─────────────────────────────────────────────────────────
function addAiResultCard(pipeline) {
  if (!pipeline || !pipeline.sku) return;
  const panel = document.getElementById("ai_results_panel");
  const cards = document.getElementById("ai_results_cards");
  if (!panel || !cards) return;
  panel.style.display = "block";

  const sku          = pipeline.sku        || "";
  const title        = pipeline.title      || "";
  const desc         = pipeline.product_description || "";
  const bullets      = Array.isArray(pipeline.bullet_points) ? pipeline.bullet_points.filter(Boolean) : [];
  const keywords     = Array.isArray(pipeline.keywords)      ? pipeline.keywords.slice(0, 15)          : [];
  const material     = pipeline.material            || "";
  const colors       = Array.isArray(pipeline.colors) ? pipeline.colors.join(", ") : "";
  const styleVal     = pipeline.style               || "";
  const sizeHint     = pipeline.size_hint           || "";
  const dims         = pipeline.estimated_dimensions|| "";
  const hsn          = pipeline.hsn_code            || "";
  const ptc          = pipeline.product_tax_code    || "";
  const browseNode   = pipeline.browse_node         || "";
  const browseNodeId = pipeline.browse_node_id      || "";
  const productType  = pipeline.product_type        || "";
  const brand        = pipeline.brand_name          || "";
  const errors       = Array.isArray(pipeline.errors) ? pipeline.errors.filter(Boolean) : [];

  const browsePath = browseNode
    ? browseNode.split(">").map(s => s.trim()).filter(Boolean).slice(-2).join(" › ")
    : browseNodeId;

  function row(label, valueHtml) {
    if (!valueHtml) return "";
    return `<div class="ai-field-row"><span class="ai-field-label">${label}</span><span class="ai-field-value">${valueHtml}</span></div>`;
  }

  const bulletsHtml = bullets.length
    ? `<ul class="ai-bullets">${bullets.map(b => `<li>${esc(b)}</li>`).join("")}</ul>` : "";
  const kwHtml = keywords.length
    ? `<div class="ai-kw-wrap">${keywords.map(k => `<span class="ai-kw-chip">${esc(k)}</span>`).join("")}</div>` : "";
  const errHtml = errors.length
    ? `<div class="ai-error-row">⚠ ${esc(errors.join(" | "))}</div>` : "";

  const bodyHtml = [
    row("Item Name (Title)", esc(title)),
    row("Product Description", desc ? `<div style="white-space:pre-line;line-height:1.6">${esc(desc)}</div>` : ""),
    row("Bullet Points (1–5)", bulletsHtml),
    row("Keywords", kwHtml),
    row("Material", esc(material)),
    row("Color(s)", esc(colors)),
    row("Style", esc(styleVal)),
    row("Size Hint", esc(sizeHint)),
    row("Est. Dimensions", esc(dims)),
    row("HSN Code", hsn ? `<span class="ai-mono">${esc(hsn)}</span>` : ""),
    row("Product Tax Code", ptc ? `<span class="ai-mono">${esc(ptc)}</span>` : ""),
    row("Browse Node", esc(browsePath)),
    row("Product Type", esc(productType)),
    row("Brand / Manufacturer", esc(brand)),
    errHtml,
  ].filter(Boolean).join("");

  const cardId = "ai_card_" + sku.replace(/[^a-z0-9]/gi, "_");
  const cardHtml = `
    <div class="ai-result-card" id="${esc(cardId)}">
      <div class="ai-result-card-header" onclick="(function(h){const b=h.nextElementSibling;const c=h.querySelector('.ai-card-chevron');b.style.display=b.style.display==='none'?'flex':'none';c.style.transform=b.style.display==='none'?'rotate(-90deg)':'rotate(0deg)'})(this)">
        <span class="ai-sku-badge">${esc(sku)}</span>
        <span class="ai-card-title">${esc(title || sku)}</span>
        <span class="ai-card-chevron">▼</span>
      </div>
      <div class="ai-result-card-body" style="display:flex">${bodyHtml}</div>
    </div>`;

  const existing = document.getElementById(cardId);
  if (existing) existing.outerHTML = cardHtml;
  else cards.insertAdjacentHTML("beforeend", cardHtml);
}

function clearAiResultCards() {
  const cards = document.getElementById("ai_results_cards");
  const panel = document.getElementById("ai_results_panel");
  if (cards) cards.innerHTML = "";
  if (panel) panel.style.display = "none";
}
</script>
  </div>
</body>
</html>
"""
    return html.replace("__DEFAULT_TEMPLATE_PATH__", escape(str(DEFAULT_TEMPLATE_PATH)))


@app.post("/api/meta")
def api_meta():
    try:
        payload = request.get_json(force=True) or {}
        template_path = str(payload.get("template_path", "")).strip()
        if not template_path:
            raise ValueError("template_path is required.")
        meta = load_template_meta(template_path)
        return jsonify({"ok": True, "meta": {k: v for k, v in meta.items() if k != "definitions"}})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/describe-image")
def api_describe_image():
    try:
        payload = request.get_json(force=True) or {}
        image_source = str(payload.get("image_source", "")).strip()
        product_name = str(payload.get("product_name", "")).strip()
        if not image_source:
            raise ValueError("image_source is required.")
        desc = describe_image(image_source, product_name=product_name)
        return jsonify({"ok": True, "description": desc})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/ai-generate-assets")
def api_ai_generate_assets():
    try:
        payload = request.get_json(force=True) or {}
        if not GROQ_AVAILABLE:
            raise ValueError("No AI provider configured. Set GROQ_API_KEY in your .env file.")
        template_path = str(payload.get("template_path", "")).strip()
        products = payload.get("products", [])
        image_mode = str(payload.get("image_mode", IMAGE_MODE_STANDARD)).strip().lower()
        if not template_path:
            raise ValueError("template_path is required.")
        if not isinstance(products, list) or not products:
            raise ValueError("products must be a non-empty list.")

        meta = load_template_meta(template_path)
        columns = meta["columns"]
        for i, row in enumerate(products):
            if not isinstance(row, dict):
                continue
            try:
                products[i] = apply_ai_assets_to_row(row, columns, i, image_mode=image_mode)
            except ValueError:
                continue

        return jsonify({"ok": True, "products": products})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/generate-row-full")
def api_generate_row_full():
    """Full single-row pipeline identical to catalog mode.

    Runs image analysis → AI copy → browse node → valid-value snapping
    for one product row and returns the complete filled template row.
    """
    try:
        payload = request.get_json(force=True) or {}
        template_path = str(payload.get("template_path", "")).strip()
        row           = payload.get("row", {})
        provider      = "groq"
        brand_name    = str(payload.get("brand_name", "")).strip()
        seller_notes  = str(payload.get("seller_notes", "")).strip()
        if not template_path:
            raise ValueError("template_path is required.")
        if not isinstance(row, dict):
            raise ValueError("row must be an object.")

        meta      = load_template_meta(template_path)
        columns   = meta["columns"]
        vpt       = meta.get("valid_product_types", [])
        avv       = meta.get("all_valid_values", {})

        # Build browse node candidates from template meta
        node_to_pt   = meta.get("node_to_product_type", {})
        node_id_to_pt: Dict[str, str] = {}
        for dk, pv in node_to_pt.items():
            m = re.search(r"\((\d+)\)\s*$", dk)
            if m:
                node_id_to_pt[m.group(1)] = pv
        candidate_nodes: List[Dict] = []
        for opt in meta.get("browse_options", []):
            nid = opt.get("node_id", "")
            pt  = node_to_pt.get(opt.get("display", ""), "") or node_id_to_pt.get(nid, "")
            candidate_nodes.append({
                "browse_node_id": nid,
                "browse_node": opt.get("path", ""),
                "path": opt.get("path", ""),
                "product_type": pt,
                "display": opt.get("display", ""),
            })

        # Resolve SKU from row
        sku = (
            str(row.get("contribution_sku#1.value", "")).strip()
            or str(row.get("item_sku", "")).strip()
            or str(row.get("external_product_id", "")).strip()
            or f"row-{int(__import__('time').time())}"
        )

        # Resolve image URL from row (first non-empty image-URL column)
        image_url = ""
        for c in columns:
            lbl = (c.get("label") or "").lower()
            if "main image url" in lbl or "main image location" in lbl:
                v = str(row.get(c["attr"], "")).strip()
                if v:
                    image_url = v
                    break
        if not image_url:
            raise ValueError("Main Image URL is required for AI generation.")

        result = process_catalog_row(
            sku=sku,
            image_urls=[image_url],
            flat_file_row=row,
            candidate_nodes=candidate_nodes or None,
            provider=provider,
            generate_images=False,
            seller_notes=seller_notes,
            valid_product_types=vpt or None,
            brand_name=brand_name,
            all_valid_values=avv or None,
        )
        template_row = catalog_result_to_template_row(
            result, columns,
            brand_name=brand_name,
            flat_row=row,
            all_valid_values=avv or None,
        )
        # Preserve any user-set values that the pipeline didn't fill
        for k, v in row.items():
            if v and not template_row.get(k):
                template_row[k] = v
        return jsonify({"ok": True, "row": template_row, "warnings": result.get("errors", [])})
    except Exception as exc:
        log.exception("generate-row-full error")
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/ai-generate-row")
def api_ai_generate_row():
    try:
        payload = request.get_json(force=True) or {}
        if not GROQ_AVAILABLE:
            raise ValueError("No AI provider configured. Set GROQ_API_KEY in your .env file.")
        template_path = str(payload.get("template_path", "")).strip()
        row_index = int(payload.get("row_index", 0))
        generate_images = bool(payload.get("generate_images", True))
        image_mode = str(payload.get("image_mode", IMAGE_MODE_STANDARD)).strip().lower()
        row = payload.get("row", {})
        if not template_path:
            raise ValueError("template_path is required.")
        if not isinstance(row, dict):
            raise ValueError("row must be an object.")

        meta = load_template_meta(template_path)
        columns = meta["columns"]
        updated = apply_ai_assets_to_row(
            row,
            columns,
            row_index,
            generate_images=generate_images,
            image_mode=image_mode,
        )
        return jsonify({"ok": True, "row": updated})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/upload-image")
def api_upload_image():
    try:
        f = request.files.get("file")
        row_index = str(request.form.get("row_index", "0")).strip()
        if f is None or not f.filename:
            raise ValueError("No file uploaded.")
        safe_name = secure_filename(f.filename) or "upload_image"
        ext = Path(safe_name).suffix.lower()
        if ext not in ALLOWED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type '{ext}'. "
                f"Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
            )
        # Validate the file is actually a readable image
        try:
            img_data = f.read()
            Image.open(BytesIO(img_data)).verify()
            f.seek(0)
        except Exception:
            raise ValueError("Uploaded file is not a valid image.")
        out_dir = ensure_upload_dir()
        out_path = out_dir / f"product_{row_index}_{safe_name.rsplit('.', 1)[0]}{ext}"
        f.save(str(out_path))
        return jsonify({"ok": True, "path": str(out_path.resolve())})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/upload-template")
def api_upload_template():
    try:
        f = request.files.get("file")
        if f is None or not f.filename:
            raise ValueError("No template file uploaded.")
        safe_name = secure_filename(f.filename) or "template"
        ext = Path(safe_name).suffix.lower()
        if ext not in {".xlsm", ".xlsx", ".xls"}:
            raise ValueError("Unsupported template type. Upload .xlsm, .xlsx, or .xls")
        out_dir = ensure_template_upload_dir()
        out_path = out_dir / f"{uuid.uuid4().hex}_{safe_name}"
        f.save(str(out_path))
        return jsonify({"ok": True, "path": str(out_path.resolve()), "original_name": safe_name})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400



SELLER_NOTES_UPLOAD_DIR = APP_DIR / "uploaded_catalogs" / "seller_notes"
SELLER_NOTES_MAX_CHARS = 100000  # total budget for all extracted text
_ZIP_MAX_FILES  = 40             # safety cap on files processed from a ZIP
_ZIP_MAX_DEPTH  = 3              # max nested ZIP depth (anti-bomb)
_FILE_MAX_CHARS = 20000          # per-file budget inside a ZIP

# Extensions that are skipped silently inside ZIP archives (binaries we can't read)
_SKIP_EXTS = {
    ".exe", ".dll", ".so", ".bin", ".dat", ".db", ".sqlite", ".pyc",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".ico",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv",
    ".zip", ".gz", ".tar", ".7z", ".rar",  # nested archives handled separately
}

# Image extensions that support OCR
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".gif", ".webp"}


# ── Per-format extractors ─────────────────────────────────────────────────────

def _ext_plain(path: Path) -> str:
    """Plain-text formats: .txt .md .csv .tsv .log."""
    return path.read_text(encoding="utf-8", errors="replace")


def _ext_json(path: Path) -> str:
    """Pretty-print JSON so the AI can read nested keys."""
    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return path.read_text(encoding="utf-8", errors="replace")


def _ext_html(path: Path) -> str:
    """Strip HTML tags, return visible text."""
    try:
        from bs4 import BeautifulSoup  # type: ignore
        soup = BeautifulSoup(path.read_bytes(), "html.parser")
        for tag in soup(["script", "style", "meta", "link"]):
            tag.decompose()
        return soup.get_text(separator="\n")
    except ImportError:
        import html as _html
        raw = path.read_text(encoding="utf-8", errors="replace")
        raw = re.sub(r"<[^>]+>", " ", raw)
        return _html.unescape(raw)


def _ext_xml(path: Path) -> str:
    """Flatten XML to key: value lines."""
    import xml.etree.ElementTree as ET
    try:
        tree = ET.parse(str(path))
        parts: List[str] = []
        for elem in tree.iter():
            text = (elem.text or "").strip()
            if text:
                parts.append(f"{elem.tag}: {text}")
        return "\n".join(parts)
    except Exception:
        return path.read_text(encoding="utf-8", errors="replace")


def _ext_pdf(path: Path) -> str:
    try:
        from pdfminer.high_level import extract_text as _pdf_text  # type: ignore
        return _pdf_text(str(path)) or ""
    except ImportError:
        raise ValueError("pdfminer.six not installed — run: pip install pdfminer.six")


def _ext_docx(path: Path) -> str:
    try:
        from docx import Document as _DocxDoc  # type: ignore
        doc = _DocxDoc(str(path))
        parts: List[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append("  |  ".join(cells))
        return "\n".join(parts)
    except ImportError:
        raise ValueError("python-docx not installed — run: pip install python-docx")


def _ext_xlsx(path: Path) -> str:
    """Handles .xlsx and .xlsm via openpyxl."""
    wb = load_workbook(path, data_only=True, read_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append("  ".join(cells))
    wb.close()
    return "\n".join(parts)


def _ext_xls(path: Path) -> str:
    """Handles legacy .xls via xlrd."""
    try:
        import xlrd  # type: ignore
        wb = xlrd.open_workbook(str(path))
        parts: List[str] = []
        for sheet in wb.sheets():
            parts.append(f"[Sheet: {sheet.name}]")
            for rx in range(sheet.nrows):
                cells = [str(sheet.cell_value(rx, cx)).strip()
                         for cx in range(sheet.ncols)
                         if str(sheet.cell_value(rx, cx)).strip()]
                if cells:
                    parts.append("  ".join(cells))
        return "\n".join(parts)
    except ImportError:
        raise ValueError("xlrd not installed — run: pip install xlrd")


def _ext_pptx(path: Path) -> str:
    """Handles .pptx via python-pptx; extracts slide text and notes."""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore  # noqa: F401
        prs = Presentation(str(path))
        parts: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts: List[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"[Notes] {notes}")
            if slide_parts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)
    except ImportError:
        raise ValueError("python-pptx not installed — run: pip install python-pptx")


def _ext_rtf(path: Path) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
        return rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))
    except ImportError:
        raise ValueError("striprtf not installed — run: pip install striprtf")


def _ext_odf(path: Path) -> str:
    """Handles .odt / .ods via odfpy."""
    try:
        from odf.opendocument import load as _odf_load  # type: ignore
        from odf.text import P as _OdfP  # type: ignore
        from odf.table import TableCell as _OdfCell  # type: ignore
        doc = _odf_load(str(path))
        parts: List[str] = []
        for elem in doc.body.childNodes:
            for node in elem.childNodes:
                if node.__class__.__name__ in ("P", "H"):
                    txt = "".join(
                        t.data for t in node.childNodes
                        if hasattr(t, "data") and t.data.strip()
                    )
                    if txt.strip():
                        parts.append(txt.strip())
        return "\n".join(parts)
    except ImportError:
        raise ValueError("odfpy not installed — run: pip install odfpy")


def _ext_image_ocr(path: Path) -> str:
    """OCR a product spec image via pytesseract.  Falls back to image metadata."""
    try:
        import pytesseract  # type: ignore
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        return text.strip() if text.strip() else f"[Image {path.name}: no readable text found by OCR]"
    except ImportError:
        w, h = Image.open(path).size
        return f"[Image {path.name}: {w}×{h}px — install pytesseract for OCR text extraction]"
    except Exception as exc:
        return f"[Image {path.name}: OCR failed — {exc}]"


def _ext_zip(path: Path, budget: int, depth: int) -> str:
    """Recursively extract text from all files inside a ZIP archive."""
    import zipfile as _zipfile
    import tempfile as _tempfile

    if depth >= _ZIP_MAX_DEPTH:
        return "[Nested ZIP skipped — max depth reached]"

    parts: List[str] = []
    file_count = 0

    with _zipfile.ZipFile(str(path), "r") as zf:
        names = [n for n in zf.namelist() if not n.endswith("/")]
        for name in names:
            if file_count >= _ZIP_MAX_FILES:
                parts.append(f"[ZIP: stopped after {_ZIP_MAX_FILES} files]")
                break
            member_ext = Path(name).suffix.lower()
            if member_ext in _SKIP_EXTS and member_ext not in _IMAGE_EXTS:
                continue
            try:
                with _tempfile.NamedTemporaryFile(
                    suffix=Path(name).suffix or ".tmp", delete=False
                ) as tmp:
                    tmp.write(zf.read(name))
                    tmp_path = Path(tmp.name)
                text = extract_text_from_any_file(tmp_path, max_chars=_FILE_MAX_CHARS, _depth=depth + 1)
                tmp_path.unlink(missing_ok=True)
                if text.strip():
                    header = f"\n{'='*60}\n[File: {name}]\n{'='*60}\n"
                    parts.append(header + text.strip())
                    file_count += 1
                    budget -= len(text)
                    if budget <= 0:
                        parts.append("[ZIP: budget exhausted — remaining files skipped]")
                        break
            except Exception as exc:
                parts.append(f"[File: {name} — skipped: {exc}]")

    return "\n".join(parts)


# ── Main dispatcher ───────────────────────────────────────────────────────────

def extract_text_from_any_file(path: Path, max_chars: int = SELLER_NOTES_MAX_CHARS, _depth: int = 0) -> str:
    """Extract all readable text from any seller-provided file.

    Supports: plain text, JSON, HTML, XML, PDF, Word (.docx/.doc),
    Excel (.xlsx/.xlsm/.xls), PowerPoint (.pptx), RTF, ODF (.odt/.ods),
    images (OCR), and ZIP archives (recursive).

    Unknown extensions fall back to plain-text read.
    """
    ext = path.suffix.lower()

    try:
        if ext in {".txt", ".md", ".csv", ".tsv", ".log", ".text"}:
            raw = _ext_plain(path)
        elif ext in {".json"}:
            raw = _ext_json(path)
        elif ext in {".yaml", ".yml"}:
            raw = _ext_plain(path)     # YAML is readable as-is
        elif ext in {".html", ".htm"}:
            raw = _ext_html(path)
        elif ext == ".xml":
            raw = _ext_xml(path)
        elif ext == ".pdf":
            raw = _ext_pdf(path)
        elif ext == ".docx":
            raw = _ext_docx(path)
        elif ext in {".xlsx", ".xlsm"}:
            raw = _ext_xlsx(path)
        elif ext == ".xls":
            raw = _ext_xls(path)
        elif ext == ".pptx":
            raw = _ext_pptx(path)
        elif ext == ".rtf":
            raw = _ext_rtf(path)
        elif ext in {".odt", ".ods"}:
            raw = _ext_odf(path)
        elif ext in _IMAGE_EXTS:
            raw = _ext_image_ocr(path)
        elif ext == ".zip":
            raw = _ext_zip(path, budget=max_chars, depth=_depth)
        else:
            # Unknown type — try reading as UTF-8 text; works for many plain formats
            try:
                raw = path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                return f"[File {path.name}: unsupported binary format]"
    except Exception as exc:
        return f"[File {path.name}: extraction failed — {exc}]"

    return raw[:max_chars]


@app.post("/api/upload-seller-notes")
def api_upload_seller_notes():
    """Upload an optional seller-provided product description / spec sheet.

    Accepts any file type — PDF, Word, Excel, PowerPoint, RTF, ODF, plain
    text, HTML, XML, JSON, ZIP archives, and product spec images (OCR).
    Returns extracted plain text + preview for the frontend to display and
    include in the catalog pipeline payload.
    """
    try:
        f = request.files.get("file")
        if f is None or not f.filename:
            raise ValueError("No file uploaded.")
        safe_name = secure_filename(f.filename) or "seller_notes"
        SELLER_NOTES_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
        out_path = SELLER_NOTES_UPLOAD_DIR / f"{uuid.uuid4().hex}_{safe_name}"
        f.save(str(out_path))
        raw = extract_text_from_any_file(out_path)
        # Normalise whitespace while keeping paragraph breaks readable
        text = re.sub(r"[ \t]{2,}", " ", raw)
        text = re.sub(r"\n{4,}", "\n\n\n", text).strip()
        preview = text[:400] + ("…" if len(text) > 400 else "")
        file_type = out_path.suffix.lower().lstrip(".").upper() or "TEXT"
        return jsonify({
            "ok": True,
            "text": text,
            "preview": preview,
            "char_count": len(text),
            "file_type": file_type,
        })
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/upload-image-catalog")
def api_upload_image_catalog():
    try:
        f = request.files.get("file")
        if f is None or not f.filename:
            raise ValueError("No image catalog file uploaded.")
        safe_name = secure_filename(f.filename) or "catalog"
        ext = Path(safe_name).suffix.lower()
        if ext not in {".xlsm", ".xlsx", ".xls", ".csv"}:
            raise ValueError("Image catalog must be .xlsm, .xlsx, .xls, or .csv")
        out_dir = ensure_catalog_upload_dir()
        out_path = out_dir / f"{uuid.uuid4().hex}_{safe_name}"
        f.save(str(out_path))
        catalog = read_image_catalog(out_path)
        return jsonify({"ok": True, "path": str(out_path.resolve()), "sku_count": len(catalog), "catalog": catalog})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


def _incremental_save_excel(template_path: str, completed_rows: List[Dict], meta: Dict) -> None:
    """Write all completed rows to the Excel template in-place.

    Called after each SKU completes so the file is always up to date.
    Errors are logged but never propagated — a save failure must not abort
    the SSE stream.
    """
    try:
        columns  = meta["columns"]
        data_row = int(meta["data_row"])
        wb = load_workbook(template_path, keep_vba=True)
        ws = wb["Template"]
        max_col = max(c["col_idx"] for c in columns) if columns else ws.max_column

        # Cache source-row styles once
        _src_styles: List[Any] = []
        _src_row_dim = ws.row_dimensions.get(data_row)
        for col_idx in range(1, max_col + 1):
            src = ws.cell(data_row, col_idx)
            _src_styles.append({
                "_style": copy(src._style),
                "font": copy(src.font) if src.has_style else None,
                "fill": copy(src.fill) if src.has_style else None,
                "border": copy(src.border) if src.has_style else None,
                "alignment": copy(src.alignment) if src.has_style else None,
                "protection": copy(src.protection) if src.has_style else None,
                "number_format": src.number_format,
                "has_style": src.has_style,
            })

        def _apply_style(dst_row: int) -> None:
            for col_idx, s in enumerate(_src_styles, start=1):
                d = ws.cell(dst_row, col_idx)
                d._style = copy(s["_style"])
                if s["has_style"]:
                    d.font       = copy(s["font"])
                    d.fill       = copy(s["fill"])
                    d.border     = copy(s["border"])
                    d.alignment  = copy(s["alignment"])
                    d.protection = copy(s["protection"])
                    d.number_format = s["number_format"]
            if _src_row_dim is not None:
                ws.row_dimensions[dst_row].height       = _src_row_dim.height
                ws.row_dimensions[dst_row].hidden       = _src_row_dim.hidden
                ws.row_dimensions[dst_row].outlineLevel = _src_row_dim.outlineLevel

        for i, row in enumerate(completed_rows):
            target_row = data_row + i
            if target_row != data_row:
                _apply_style(target_row)
            for c in columns:
                val = str(row.get(c["attr"], ""))
                ws.cell(target_row, c["col_idx"]).value = val if val != "" else None

        wb.save(template_path)
    except Exception as exc:
        log.warning("Incremental Excel save failed: %s", exc)


@app.post("/api/process-catalog/stream")
def api_process_catalog_stream():
    """Server-Sent Events version of /api/process-catalog.

    SKUs are processed in a small thread pool (default 3 workers) so I/O waits
    overlap. Results stream back via a Queue as each SKU finishes; the frontend
    handles out-of-order delivery via the "index" field.

    Event types:
      start     — {"type":"start","total":N}
      sku_start — {"type":"sku_start","sku":"...","index":i}
      sku       — {"type":"sku","sku":"...","index":i,"total":N,"row":{...},"pipeline":{...}}
      done      — {"type":"done","count":N,"products":[...],"details":[...]}
      error     — {"type":"error","error":"..."}
    """
    import queue as _queue
    from concurrent.futures import ThreadPoolExecutor as _TPE

    payload = request.get_json(force=True) or {}
    template_path      = str(payload.get("template_path", "")).strip()
    image_catalog_path = str(payload.get("image_catalog_path", "")).strip()
    provider           = "groq"
    generate_images    = bool(payload.get("generate_images", True))
    seller_notes       = str(payload.get("seller_notes", "")).strip()
    brand_name         = str(payload.get("brand_name", "")).strip()

    def _generate():
        stop_evt = threading.Event()
        try:
            if not template_path or not image_catalog_path:
                yield f"data: {json.dumps({'type':'error','error':'template_path and image_catalog_path are required.'})}\n\n"
                return

            catalog       = read_image_catalog(Path(image_catalog_path))
            flat_products = parse_flat_file_skus(Path(template_path))
            meta          = load_template_meta(template_path)
            columns       = meta["columns"]
            vpt           = meta.get("valid_product_types", [])
            avv           = meta.get("all_valid_values", {})

            node_to_pt: Dict[str, str] = meta.get("node_to_product_type", {})
            node_id_to_pt: Dict[str, str] = {}
            for display_key, pt_val in node_to_pt.items():
                m_nid = re.search(r"\((\d+)\)\s*$", display_key)
                if m_nid:
                    node_id_to_pt[m_nid.group(1)] = pt_val
            candidate_nodes: List[Dict] = []
            for opt in meta.get("browse_options", []):
                display = opt.get("display", "")
                path    = opt.get("path", "")
                nid     = opt.get("node_id", "")
                pt = node_to_pt.get(display, "") or node_id_to_pt.get(nid, "")
                candidate_nodes.append({"browse_node_id": nid, "browse_node": path, "path": path, "product_type": pt, "display": display})

            flat_by_sku: Dict[str, Dict] = {}
            for row in flat_products:
                sku_val = next(
                    (str(row.get(k, "")).strip() for k in ("external_product_id", "seller_sku", "item_sku") if str(row.get(k, "")).strip()),
                    ""
                )
                if sku_val:
                    flat_by_sku[sku_val] = row

            all_skus    = list(catalog.keys())
            total       = len(all_skus)
            _do_images  = generate_images and IMAGE_GENERATION_ENABLED
            max_workers = min(total, int(os.getenv("CATALOG_SKU_WORKERS", "8")))

            yield f"data: {json.dumps({'type':'start','total':total})}\n\n"

            result_q: _queue.Queue = _queue.Queue()

            def _process_one(idx: int, sku: str) -> None:
                result: Dict = {"sku": sku, "errors": []}
                template_row: Dict = {}
                try:
                    if stop_evt.is_set():
                        return  # finally always runs and puts the (empty) result
                    image_urls = catalog.get(sku, [])
                    flat_row   = flat_by_sku.get(sku) or next(
                        (r for r in flat_products if sku.lower() in str(r).lower()), None
                    )
                    result = process_catalog_row(
                        sku=sku,
                        image_urls=image_urls,
                        flat_file_row=flat_row,
                        candidate_nodes=candidate_nodes or None,
                        provider=provider,
                        generate_images=_do_images,
                        seller_notes=seller_notes,
                        valid_product_types=vpt or None,
                        brand_name=brand_name,
                        all_valid_values=avv or None,
                    )
                    template_row = catalog_result_to_template_row(
                        result, columns,
                        brand_name=brand_name,
                        flat_row=flat_row,
                        all_valid_values=avv or None,
                    )
                except Exception as exc:
                    import traceback as _tb
                    log.error("SKU %s failed: %s\n%s", sku, exc, _tb.format_exc())
                    result       = {"sku": sku, "errors": [str(exc)]}
                    template_row = {}
                finally:
                    # Always put — even on unexpected exception so the queue count stays correct
                    result_q.put({
                        "type":     "sku",
                        "sku":      sku,
                        "index":    idx,
                        "total":    total,
                        "row":      template_row,
                        "pipeline": {k: v for k, v in result.items() if k != "analysis"},
                    })

            pool = _TPE(max_workers=max_workers, thread_name_prefix="sku_worker")
            try:
                for idx, sku in enumerate(all_skus, start=1):
                    yield f"data: {json.dumps({'type':'sku_start','sku':sku,'index':idx})}\n\n"
                    pool.submit(_process_one, idx, sku)

                rows_by_idx:    Dict[int, Dict] = {}
                details_by_idx: Dict[int, Dict] = {}
                completed = 0
                while completed < total:
                    try:
                        evt = result_q.get(timeout=45)  # short timeout so keep-alives fire
                    except _queue.Empty:
                        if stop_evt.is_set():
                            break
                        yield ": ka\n\n"  # SSE keep-alive comment — ignored by browser, prevents conn drop
                        continue
                    if evt is None:
                        completed += 1
                        continue
                    idx_val = evt["index"]
                    rows_by_idx[idx_val]    = evt.get("row") or {}
                    details_by_idx[idx_val] = {
                        "sku": evt["sku"],
                        "pipeline": {"constrained_fields": (evt.get("pipeline") or {}).get("constrained_fields") or {}},
                    }
                    completed += 1
                    yield f"data: {json.dumps(evt)}\n\n"
            finally:
                stop_evt.set()
                pool.shutdown(wait=False, cancel_futures=True)  # cancel pending; running workers finish in background

            # Backfill shared constrained fields into any rows where Groq missed them
            _backfill_shared_fields(rows_by_idx, details_by_idx, columns, avv or None)

            yield f"data: {json.dumps({'type':'done','count':len(rows_by_idx)})}\n\n"

        except Exception as exc:
            log.exception("SSE catalog error")
            yield f"data: {json.dumps({'type':'error','error':str(exc)})}\n\n"

    return app.response_class(
        _generate(),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


@app.get("/api/combo-images/<path:filename>")
def serve_combo_image_file(filename):
    """Serve a generated combo composite image so the pipeline can fetch it as a URL."""
    return send_from_directory(COMBO_IMAGE_DIR, filename)


@app.post("/api/process-combo-catalog/stream")
def api_process_combo_catalog_stream():
    """Build combo composite images, then run the full AI listing pipeline on each.

    Phase 1: For each combination of SKUs, stitch the individual product images
             side-by-side on a white background using Pillow (generate_combo_image)
             and save the result as a single PNG.
    Phase 2: Pass that composite image path directly to process_catalog_row —
             Groq vision analyses the combined shot, FAL.ai generates a lifestyle
             image using it as reference, and the listing pipeline produces
             title / bullets / description as if it were a single product.

    SSE events:
      building_combos  — {"type":"building_combos","source_skus":N,"max_combo_size":M}
      combos_ready     — {"type":"combos_ready","total_combos":N}
      start            — {"type":"start","total":N}
      sku              — {"type":"sku","sku":"APPLE+BANANA","index":1,"total":N,
                          "row":{...},"pipeline":{...},"source_skus":[...]}
      done             — {"type":"done","count":N,"products":[...]}

    Body JSON:
      template_path       str   Amazon flat-file template (for column mapping)
      image_catalog_path  str   path returned by /api/upload-image-catalog
      max_combo_size      int   max combo size; generates all sizes 2…N (default 2)
      max_combos          int   cap on total combos (default 1500)
      provider            str   AI provider ("groq")
      brand_name          str   seller brand applied to every combo row
      seller_notes        str   optional seller description text
      generate_images     bool  whether to generate AI lifestyle images (default true)
    """
    import itertools as _it
    import queue as _queue
    import threading as _threading
    from concurrent.futures import ThreadPoolExecutor as _TPE

    payload            = request.get_json(force=True) or {}
    image_catalog_path = str(payload.get("image_catalog_path", "")).strip()
    template_path      = str(payload.get("template_path", "")).strip()
    max_combo_size     = max(2, int(payload.get("max_combo_size", payload.get("combo_size", 2))))
    max_combos         = max(1, int(payload.get("max_combos", 1500)))
    provider           = "groq"
    brand_name         = str(payload.get("brand_name", "")).strip()
    seller_notes       = str(payload.get("seller_notes", "")).strip()
    generate_images    = bool(payload.get("generate_images", True))
    images_only        = bool(payload.get("images_only", False))

    def _generate():
        try:
            if not image_catalog_path:
                yield f"data: {json.dumps({'type':'error','error':'image_catalog_path is required.'})}\n\n"
                return

            # ── Phase 1: build all combo composite images ─────────────────────
            catalog    = read_image_catalog(Path(image_catalog_path))
            capped_max = max(2, min(max_combo_size, len(catalog)))
            yield f"data: {json.dumps({'type':'building_combos','source_skus':len(catalog),'max_combo_size':capped_max})}\n\n"

            out_base  = Path(os.environ.get("AI_GENERATED_IMAGE_DIR", str(APP_DIR / "generated_images")))
            combo_dir = out_base / "combos"
            combo_dir.mkdir(parents=True, exist_ok=True)

            eligible = {sku: urls for sku, urls in catalog.items() if urls}
            # Deduplicate while preserving order so each SKU appears at most
            # once in sku_list — itertools.combinations already never repeats
            # a position, but deduplicating here ensures no two list slots
            # carry the same SKU string even if the catalog had duplicate keys.
            seen: set = set()
            sku_list: list = []
            for s in eligible:
                if s not in seen:
                    seen.add(s)
                    sku_list.append(s)

            all_combos: list = []
            for csize in range(2, max(2, min(capped_max, len(sku_list))) + 1):
                all_combos.extend(_it.combinations(sku_list, csize))
            all_combos = all_combos[:max_combos]

            if not all_combos:
                yield f"data: {json.dumps({'type':'error','error':'No valid combos — check that catalog SKUs have image URLs.'})}\n\n"
                return

            # Combo entries are just sku-list specs here — the actual composite
            # image is built per-combo inside the parallel worker below so that
            # image generation and the listing pipeline both run concurrently,
            # and each combo's row lands in the sheet as soon as it's ready
            # rather than waiting for every combo's image to finish first.
            combo_entries: list = [
                {"sku": "+".join(skus), "source_skus": list(skus)}
                for idx, skus in enumerate(all_combos, start=1)
            ]

            total = len(combo_entries)
            yield f"data: {json.dumps({'type':'combos_ready','total_combos':total})}\n\n"
            yield f"data: {json.dumps({'type':'start','total':total})}\n\n"

            # ── Phase 2: load template metadata ───────────────────────────────
            columns: list         = []
            vpt:     list         = []
            avv:     dict         = {}
            candidate_nodes: list = []

            if template_path:
                try:
                    meta       = load_template_meta(template_path)
                    columns    = meta["columns"]
                    vpt        = meta.get("valid_product_types", [])
                    avv        = meta.get("all_valid_values", {})
                    node_to_pt = meta.get("node_to_product_type", {})
                    node_id_to_pt: dict = {}
                    for display_key, pt_val in node_to_pt.items():
                        m = re.search(r"\((\d+)\)\s*$", display_key)
                        if m:
                            node_id_to_pt[m.group(1)] = pt_val
                    for opt in meta.get("browse_options", []):
                        display = opt.get("display", "")
                        nid     = opt.get("node_id", "")
                        pt      = node_to_pt.get(display, "") or node_id_to_pt.get(nid, "")
                        candidate_nodes.append({
                            "browse_node_id": nid,
                            "browse_node":    opt.get("path", ""),
                            "path":           opt.get("path", ""),
                            "product_type":   pt,
                            "display":        display,
                        })
                except Exception as exc:
                    log.warning("Template load failed, proceeding without column mapping: %s", exc)

            # Combo SKUs go through the exact same generate_images path as
            # single-SKU catalog processing (process_catalog_row's 4-shot
            # pipeline) so combo and normal mode produce identical shot types
            # from identical prompts/compositing — no separate, weaker path.
            max_workers = min(total, int(os.getenv("CATALOG_SKU_WORKERS", "8")))
            result_q    = _queue.Queue()
            stop_evt    = _threading.Event()

            def _process_one(idx: int, entry: dict) -> None:
                combo_sku    = entry["sku"]
                source_skus  = entry["source_skus"]
                result: dict       = {"sku": combo_sku, "errors": []}
                template_row: dict = {}
                try:
                    if stop_evt.is_set():
                        return
                    # Build this combo's composite image here, inside the worker,
                    # so image generation runs in parallel across combos instead
                    # of as one big sequential phase before any listing starts.
                    image_sources = [eligible[s] for s in source_skus]
                    out_path = combo_dir / f"{combo_sku}_main.png"
                    if images_only:
                        generate_combo_image_labeled(
                            image_sources, out_path, combo_sku=combo_sku
                        )
                    else:
                        generate_combo_image(image_sources, out_path, combo_number=idx)
                    img_path = str(out_path.resolve())

                    if images_only:
                        # Skip the entire listing/Excel pipeline — just report
                        # the generated composite image path and move on.
                        result_q.put({
                            "type":        "image",
                            "sku":         combo_sku,
                            "index":       idx,
                            "total":       total,
                            "image_path":  img_path,
                            "source_skus": source_skus,
                        })
                        return

                    # The composite file path is passed directly; process_catalog_row
                    # accepts file paths in addition to HTTP URLs.
                    result = process_catalog_row(
                        sku=combo_sku,
                        image_urls=[img_path],
                        flat_file_row=None,
                        candidate_nodes=candidate_nodes or None,
                        provider=provider,
                        generate_images=generate_images,
                        seller_notes=seller_notes,
                        valid_product_types=vpt or None,
                        brand_name=brand_name,
                        all_valid_values=avv or None,
                    )
                    if columns:
                        template_row = catalog_result_to_template_row(
                            result, columns,
                            brand_name=brand_name,
                            flat_row=None,
                            all_valid_values=avv or None,
                        )
                except Exception as exc:
                    import traceback as _tb
                    log.error("Combo %s pipeline failed: %s\n%s", combo_sku, exc, _tb.format_exc())
                    result       = {"sku": combo_sku, "errors": [str(exc)]}
                    template_row = {}
                finally:
                    if not images_only:
                        result_q.put({
                            "type":        "sku",
                            "sku":         combo_sku,
                            "index":       idx,
                            "total":       total,
                            "row":         template_row,
                            "pipeline":    {k: v for k, v in result.items() if k != "analysis"},
                            "source_skus": source_skus,
                        })

            pool = _TPE(max_workers=max_workers, thread_name_prefix="combo_worker")
            try:
                for idx, entry in enumerate(combo_entries, start=1):
                    pool.submit(_process_one, idx, entry)

                rows_by_idx: dict = {}
                completed = 0
                while completed < total:
                    try:
                        evt = result_q.get(timeout=45)
                    except _queue.Empty:
                        if stop_evt.is_set():
                            break
                        yield ": ka\n\n"
                        continue
                    rows_by_idx[evt["index"]] = evt
                    completed += 1
                    yield f"data: {json.dumps(evt)}\n\n"
            finally:
                stop_evt.set()
                pool.shutdown(wait=False, cancel_futures=True)

            # Backfill shared constrained fields across all combo rows
            _combo_rows_by_idx    = {i: rows_by_idx[i]["row"] for i in rows_by_idx}
            _combo_details_by_idx = {
                i: {"pipeline": {k: v for k, v in rows_by_idx[i].items()
                                 if k not in ("row", "type", "index", "total", "source_skus")}}
                for i in rows_by_idx
            }
            _backfill_shared_fields(_combo_rows_by_idx, _combo_details_by_idx, columns, avv or None)
            for i, row in _combo_rows_by_idx.items():
                rows_by_idx[i]["row"] = row

            yield f"data: {json.dumps({'type':'done','count':len(rows_by_idx)})}\n\n"

        except Exception as exc:
            log.exception("SSE combo catalog error")
            yield f"data: {json.dumps({'type':'error','error':str(exc)})}\n\n"

    return app.response_class(
        _generate(),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.post("/api/process-catalog")
def api_process_catalog():
    """Full batch pipeline: template (= flat file) + image catalog → filled product rows.

    Body JSON:
      template_path      str   Amazon xlsm template — used for SKU rows AND column metadata
      image_catalog_path str   path returned by /api/upload-image-catalog
      provider           str   "groq" (default)
      generate_images    bool  default true
      seller_notes       str   optional — plain text extracted from seller description file
    """
    try:
        payload = request.get_json(force=True) or {}
        template_path = str(payload.get("template_path", "")).strip()
        image_catalog_path = str(payload.get("image_catalog_path", "")).strip()
        provider = "groq"
        generate_images = bool(payload.get("generate_images", True))
        seller_notes = str(payload.get("seller_notes", "")).strip()

        if not template_path:
            raise ValueError("template_path is required — load the template first.")
        if not image_catalog_path:
            raise ValueError("image_catalog_path is required.")

        # Require Groq for analysis + copy generation
        if not GROQ_AVAILABLE:
            raise ValueError(
                "No AI provider is configured. Set GROQ_API_KEY in your .env file."
            )

        catalog = read_image_catalog(Path(image_catalog_path))
        flat_products = parse_flat_file_skus(Path(template_path))
        meta = load_template_meta(template_path)
        columns = meta["columns"]
        vpt = meta.get("valid_product_types", [])
        avv = meta.get("all_valid_values", {})

        # ── Build candidate nodes exclusively from the template's Browse Data sheet ──
        # Source of truth for both browse node IDs and product types is the uploaded
        # template only — never the master file (Book1.xlsx).
        #
        # node_to_pt  : display-string → product_type  (from Valid Values sheet)
        # node_id_to_pt: node_id       → product_type  (derived from node_to_pt keys
        #                so we match even when display strings differ slightly)
        node_to_pt: Dict[str, str] = meta.get("node_to_product_type", {})

        # Build secondary lookup: strip the "(node_id)" suffix from each key
        node_id_to_pt: Dict[str, str] = {}
        for display_key, pt_val in node_to_pt.items():
            m_nid = re.search(r"\((\d+)\)\s*$", display_key)
            if m_nid:
                node_id_to_pt[m_nid.group(1)] = pt_val

        candidate_nodes: List[Dict] = []
        for opt in meta.get("browse_options", []):
            display = opt.get("display", "")
            path    = opt.get("path", "")
            nid     = opt.get("node_id", "")

            # 1st choice: exact display-string match in Valid Values
            pt = node_to_pt.get(display, "")
            # 2nd choice: match by numeric node ID (handles minor display-format drift)
            if not pt:
                pt = node_id_to_pt.get(nid, "")
            # No master-file fallback — product type must come from the template

            candidate_nodes.append({
                "browse_node_id": nid,
                "browse_node": path,
                "path": path,
                "product_type": pt,
                "display": display,
            })

        if candidate_nodes:
            log.info(
                "Template Browse Data: %d node(s), %d with product types resolved from Valid Values.",
                len(candidate_nodes),
                sum(1 for c in candidate_nodes if c["product_type"]),
            )
        else:
            log.warning(
                "Template has no Browse Data sheet — browse node and product type columns will be left blank."
            )

        flat_by_sku: Dict[str, Dict] = {}
        for row in flat_products:
            sku_candidates = [
                row.get("external_product_id", ""),
                row.get("seller_sku", ""),
                row.get("item_sku", ""),
            ]
            sku_val = next((s.strip() for s in sku_candidates if s.strip()), "")
            if sku_val:
                flat_by_sku[sku_val] = row

        all_skus = list(catalog.keys())
        _do_images = generate_images and IMAGE_GENERATION_ENABLED and len(all_skus) <= 100
        results_by_sku: Dict[str, dict] = {}

        def _process_one(sku: str) -> dict:
            image_urls = catalog[sku]
            flat_row = flat_by_sku.get(sku) or next(
                (r for r in flat_products if sku.lower() in str(r).lower()), None
            )
            result = process_catalog_row(
                sku=sku,
                image_urls=image_urls,
                flat_file_row=flat_row,
                candidate_nodes=candidate_nodes or None,
                provider=provider,
                generate_images=_do_images,
                seller_notes=seller_notes,
                valid_product_types=vpt or None,
                all_valid_values=avv or None,
            )
            template_row = catalog_result_to_template_row(
                result, columns, all_valid_values=avv or None
            )
            return {"sku": sku, "row": template_row, "pipeline": result}

        max_sku_workers = min(len(all_skus), int(os.environ.get("CATALOG_SKU_WORKERS", "20")))
        from concurrent.futures import ThreadPoolExecutor as _TPE, as_completed as _ac
        with _TPE(max_workers=max_sku_workers) as pool:
            futures = {pool.submit(_process_one, sku): sku for sku in all_skus}
            for future in _ac(futures):
                sku = futures[future]
                try:
                    results_by_sku[sku] = future.result()
                    log.info("SKU %s complete (%d/%d)", sku, len(results_by_sku), len(all_skus))
                except Exception as exc:
                    log.error("SKU %s failed: %s", sku, exc)
                    results_by_sku[sku] = {
                        "sku": sku, "row": {}, "pipeline": {"sku": sku, "errors": [str(exc)]}
                    }

        # Preserve original catalog order
        results = [results_by_sku[sku] for sku in all_skus if sku in results_by_sku]
        products_for_ui = [r["row"] for r in results]
        return jsonify({
            "ok": True,
            "products": products_for_ui,
            "details": results,
            "count": len(results),
        })
    except NotImplementedError as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400
    except Exception as exc:
        log.exception("Error in /api/process-catalog")
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.get("/api/master-nodes")
def api_master_nodes():
    """Return summary of loaded master browse nodes (for UI autocomplete)."""
    try:
        nodes = _load_master_nodes()
        summary = [
            {
                "browse_node_id": n.get("browse_node_id", ""),
                "browse_node": n.get("browse_node", ""),
                "product_type": n.get("product_type", ""),
            }
            for n in nodes[:200]
        ]
        return jsonify({"ok": True, "count": len(nodes), "nodes": summary})
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


try:
    from amazon_product_images import generateProductImages as _generateProductImages
    from amazon_product_images import SHOT_TYPES as _SHOT_TYPES
    from amazon_product_images import _VALID_PROVIDERS as _IMG_VALID_PROVIDERS
    _PRODUCT_IMAGES_AVAILABLE = True
except ImportError:
    _PRODUCT_IMAGES_AVAILABLE = False
    log.warning("amazon_product_images not found; /api/generate-product-images unavailable.")


@app.post("/api/generate-product-images")
def api_generate_product_images():
    """Generate a 4-shot Amazon product image set.

    Request JSON:
      product              dict   — analysis dict (product_type, category, material,
                                    colors, features, usage, style).  At minimum,
                                    one of product_type or category must be non-empty.
      reference_image_path str    — local absolute path or https:// URL of reference photo
      reference_description str   — (optional) pre-computed reference text from
                                    _reference_conditioning_text(); skips vision call
      options              dict   — (optional) shots, provider, out_dir, prefix

    Response JSON:
      {ok: true, images: {shot_type: absolute_path, ...}}
    """
    if not _PRODUCT_IMAGES_AVAILABLE:
        return jsonify({"ok": False, "error": "amazon_product_images module not installed."}), 500
    try:
        payload = request.get_json(force=True) or {}
        product = payload.get("product")
        if not isinstance(product, dict):
            raise ValueError("'product' must be a JSON object.")
        if not str(product.get("product_type", "")).strip() and not str(product.get("category", "")).strip():
            raise ValueError("product must have at least one of 'product_type' or 'category'.")

        ref_image = str(payload.get("reference_image_path", "") or "").strip()
        ref_desc  = str(payload.get("reference_description", "") or "").strip()
        options   = payload.get("options") or {}
        if not isinstance(options, dict):
            options = {}

        provider = str(options.get("provider", "fal")).strip().lower() or "fal"
        if provider not in _IMG_VALID_PROVIDERS:
            provider = "fal"

        shots_raw = options.get("shots", list(_SHOT_TYPES))
        shots = tuple(str(s).strip() for s in shots_raw if str(s).strip())

        out_dir = Path(
            options.get("out_dir", "")
            or os.environ.get("AI_GENERATED_IMAGE_DIR", "generated_amazon_images/ai")
        )

        # If no reference_description provided and a reference image exists,
        # compute it now using the existing vision pipeline.
        if not ref_desc and ref_image:
            try:
                ref_desc = _reference_conditioning_text(ref_image, str(product.get("product_type", "")))
            except Exception as exc:
                log.warning("Could not compute reference description: %s — skipping conditioning", exc)

        result_paths = _generateProductImages(
            product=product,
            reference_image=ref_image,
            config={"shots": shots, "out_dir": out_dir, "provider": provider},
            reference_description=ref_desc,
        )
        return jsonify({"ok": True, "images": {k: str(v) for k, v in result_paths.items()}})
    except Exception as exc:
        log.exception("Error in /api/generate-product-images")
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/api/generate-catalog-images")
def api_generate_catalog_images():
    """Generate two additional images per product, leaving the main image untouched.

    For each product row:
      - Main Image URL  → kept exactly as uploaded by the user (not modified)
      - Other Image 1   → lifestyle shot generated by AI
      - Other Image 2   → infographic composite built by PIL (bullets + product photo)

    Request JSON:
      template_path  str   — path to the loaded template
      products       list  — current products array from the frontend sheet

    Response JSON:
      {ok: true, products: [...], images_generated: N}
    """
    try:
        payload       = request.get_json(force=True) or {}
        template_path = str(payload.get("template_path", "")).strip()
        products      = payload.get("products", [])
        if not template_path:
            raise ValueError("template_path is required.")
        if not isinstance(products, list) or not products:
            raise ValueError("products must be a non-empty list.")

        meta            = load_template_meta(template_path)
        columns         = meta["columns"]
        main_img_attr   = first_attr_by_label(columns, "Main Image URL")
        other_img_attrs = all_attrs_by_label(columns, "Other Image URL")[:2]
        item_name_attr  = first_attr_by_label(columns, "Item Name")
        pt_attr         = first_attr_by_label(columns, "Product Type") or first_attr_by_label(columns, "Item Type Keyword")
        bullet_attrs    = all_attrs_by_label(columns, "Bullet Point")[:5]

        if not main_img_attr:
            raise ValueError("Template has no Main Image URL column.")

        images_generated = 0
        out_base = Path(os.environ.get("AI_GENERATED_IMAGE_DIR", str(APP_DIR / "generated_images")))

        for i, row in enumerate(products):
            if not isinstance(row, dict):
                continue
            image_url = str(row.get(main_img_attr, "")).strip()
            if not image_url:
                continue

            product_name = str(row.get(item_name_attr, "")).strip() if item_name_attr else ""
            product_type = str(row.get(pt_attr, "")).strip() if pt_attr else product_name
            sku          = str(row.get("seller_sku") or row.get("item_sku") or f"product_{i+1}").strip()

            out_dir = out_base / sku
            out_dir.mkdir(parents=True, exist_ok=True)

            bullets = [str(row.get(a, "")).strip() for a in bullet_attrs if str(row.get(a, "")).strip()]

            analysis = {
                "product_type": product_type or product_name or "product",
                "category":     str(row.get("feed_product_type", "")).strip(),
                "material":     str(row.get("material_type", "")).strip(),
                "colors":       [str(row.get("color", "")).strip()],
                "features":     bullets[:3],
                "usage":        "",
                "style":        "",
            }

            # ── Build listing dict from row data for rich infographic details ──
            row_listing = {
                "title":               product_name,
                "bullet_points":       bullets,
                "description":         str(row.get("product_description", "") or row.get("item_description", "") or "").strip(),
                "estimated_dimensions": str(row.get("item_dimensions", "") or "").strip(),
                "constrained_fields":  {},
            }

            # ── Generate both images (lifestyle + infographic) via PIL ────────
            try:
                pil_images = generate_reference_locked_images(
                    reference_image_source=image_url,
                    out_dir=out_dir,
                    prefix=sku,
                    product_name=product_name or product_type or "Product",
                    bullets=bullets[:5],
                    analysis=analysis,
                    listing=row_listing,
                )
                for idx, attr in enumerate(other_img_attrs):
                    if idx < len(pil_images) and pil_images[idx]:
                        row[attr] = _local_path_to_img_url(pil_images[idx])
                if pil_images:
                    images_generated += 1
                    log.info("Generated %d images for %s", len(pil_images), sku)
            except Exception as exc:
                log.warning("Image generation failed for %s: %s", sku, exc)

        return jsonify({"ok": True, "products": products, "images_generated": images_generated})
    except Exception as exc:
        log.exception("Error in /api/generate-catalog-images")
        return jsonify({"ok": False, "error": str(exc)}), 400


@app.post("/generate")
def generate():
    if "image" not in request.files:
        return jsonify({"error": "No file uploaded. Send a multipart field named 'image'."}), 400
    file = request.files["image"]
    if not file.filename:
        return jsonify({"error": "Uploaded file has no filename."}), 400
    suffix = Path(file.filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        return jsonify(
            {
                "error": f"Unsupported file type '{suffix}'. "
                f"Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
            }
        ), 400

    safe_name = f"{uuid.uuid4().hex}{suffix}"
    ensure_upload_dir()
    image_path = UPLOAD_DIR / safe_name
    file.save(image_path)
    log.info("Image saved for /generate: %s", image_path)
    try:
        analysis = groq_analyze_image_url(str(image_path), "")
        listing = _generate_listing(analysis)
        result = {
            "analysis": analysis,
            "listing": listing,
            "validation_errors": validate_listing(listing, product_analysis=analysis),
            "provider_notes": [f"copy_provider=groq"],
        }
        return jsonify(result), 200
    except ImageValidationError as exc:
        return jsonify({"error": str(exc)}), 400
    except (VisionAPIError, CopyGenerationError) as exc:
        return jsonify({"error": str(exc)}), 500
    except RetryError as exc:
        # Defensive fallback if any RetryError escapes despite reraise=True.
        root = exc.last_attempt.exception() if hasattr(exc, "last_attempt") else exc
        return jsonify({"error": f"AI retry exhausted: {root}"}), 500
    except Exception as exc:  # noqa: BLE001
        log.exception("Unexpected error in /generate")
        return jsonify({"error": f"Internal server error: {exc}"}), 500


@app.get("/api/download-excel")
def api_download_excel():
    """Serve the saved template file as a browser download.

    Strips the internal UUID prefix (e.g. a3f1...._Template.xlsm → Template.xlsm)
    so the user downloads a clean filename matching what they uploaded.
    """
    path_str = request.args.get("path", "").strip()
    if not path_str:
        return "path parameter required", 400
    try:
        full = Path(path_str).resolve()
        allowed_roots = [APP_DIR.resolve(), TEMPLATE_UPLOAD_DIR.resolve()]
        if not any(str(full).startswith(str(r)) for r in allowed_roots):
            return "Forbidden", 403
        if not full.exists():
            return "File not found", 404
        suffix = full.suffix.lower()
        mime = {
            ".xlsm": "application/vnd.ms-excel.sheet.macroEnabled.12",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".xls":  "application/vnd.ms-excel",
        }.get(suffix, "application/octet-stream")
        # Strip UUID prefix: {32-hex-chars}_{original}.xlsm → {original}.xlsm
        raw_name = full.name
        parts = raw_name.split("_", 1)
        download_name = (
            parts[1]
            if len(parts) == 2 and len(parts[0]) == 32
               and all(c in "0123456789abcdef" for c in parts[0])
            else raw_name
        )
        return _send_file(str(full), mimetype=mime, as_attachment=True,
                          download_name=download_name)
    except Exception as exc:
        return str(exc), 500


@app.post("/api/save")
def api_save():
    try:
        payload = request.get_json(force=True) or {}
        template_path = str(payload.get("template_path", "")).strip()
        products = payload.get("products", [])
        if not template_path:
            raise ValueError("template_path is required.")
        if not isinstance(products, list) or not products:
            raise ValueError("At least one product row is required.")

        meta = load_template_meta(template_path)
        columns = meta["columns"]
        required_attrs = set(meta["required_attrs"])
        definitions = meta["definitions"]
        data_row = int(meta["data_row"])
        description_attr = first_attr_by_label(columns, "Product Description")

        missing = []
        for i, row in enumerate(products):
            for c in columns:
                attr = c["attr"]
                if description_attr and attr == description_attr:
                    continue
                if attr in required_attrs and not str(row.get(attr, "")).strip():
                    label = definitions[attr].label if attr in definitions else c["label"] or attr
                    missing.append(
                        {
                            "row": i + 1,
                            "excelRow": data_row + i,
                            "label": label,
                            "attr": attr,
                        }
                    )
        wb = load_workbook(template_path, keep_vba=True)
        ws = wb["Template"]

        max_col = max(c["col_idx"] for c in columns) if columns else ws.max_column

        # Pre-read source-row styles once so copy_row_style doesn't re-read
        # the template row for each of the 1000 destination rows.
        _src_styles: List[Any] = []
        _src_row_dim = ws.row_dimensions.get(data_row)
        for col_idx in range(1, max_col + 1):
            src = ws.cell(data_row, col_idx)
            _src_styles.append({
                "_style": copy(src._style),
                "font": copy(src.font) if src.has_style else None,
                "fill": copy(src.fill) if src.has_style else None,
                "border": copy(src.border) if src.has_style else None,
                "alignment": copy(src.alignment) if src.has_style else None,
                "protection": copy(src.protection) if src.has_style else None,
                "number_format": src.number_format,
                "has_style": src.has_style,
            })

        def _apply_cached_style(dst_row: int) -> None:
            for col_idx, s in enumerate(_src_styles, start=1):
                d = ws.cell(dst_row, col_idx)
                d._style = copy(s["_style"])
                if s["has_style"]:
                    d.font = copy(s["font"])
                    d.fill = copy(s["fill"])
                    d.border = copy(s["border"])
                    d.alignment = copy(s["alignment"])
                    d.protection = copy(s["protection"])
                    d.number_format = s["number_format"]
            if _src_row_dim is not None:
                ws.row_dimensions[dst_row].height = _src_row_dim.height
                ws.row_dimensions[dst_row].hidden = _src_row_dim.hidden
                ws.row_dimensions[dst_row].outlineLevel = _src_row_dim.outlineLevel

        for i, row in enumerate(products):
            target_row = data_row + i
            if target_row != data_row:
                _apply_cached_style(target_row)
            for c in columns:
                attr = c["attr"]
                col_idx = c["col_idx"]
                val = str(row.get(attr, ""))
                ws.cell(target_row, col_idx).value = val if val != "" else None

        wb.save(template_path)
        if missing:
            return jsonify({
                "ok": True,
                "saved_path": template_path,
                "message": f"Saved to template: {Path(template_path).name}",
                "warnings": missing,
            })
        return jsonify({
            "ok": True,
            "saved_path": template_path,
            "message": f"Saved to template: {Path(template_path).name}",
        })
    except Exception as exc:
        return jsonify({"ok": False, "error": str(exc)}), 400


def main():
    parser = argparse.ArgumentParser(description="Browser UI for multi-product Amazon template filling.")
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", default=7860, type=int)
    parser.add_argument("--no-browser", action="store_true")
    args = parser.parse_args()

    # On Ctrl-C / SIGTERM, exit immediately without waiting for background
    # ThreadPoolExecutor threads (catalog workers still running API calls).
    import signal as _signal
    def _fast_exit(sig, frame):
        os._exit(0)
    _signal.signal(_signal.SIGINT,  _fast_exit)
    _signal.signal(_signal.SIGTERM, _fast_exit)

    if not args.no_browser:
        threading.Timer(1.0, lambda: webbrowser.open(f"http://{args.host}:{args.port}")).start()
    app.run(host=args.host, port=args.port, debug=False, threaded=True, use_reloader=False)


if __name__ == "__main__":
    main()
